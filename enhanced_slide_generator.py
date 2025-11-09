import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
from datetime import datetime, timedelta
import sys
import os

class EnhancedCatholicSlideGenerator:
    def __init__(self):
        self.base_url = "https://www.ibreviary.com/m2/"
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })

    def fetch_morning_prayer_detailed(self):
        """
        Fetch and parse morning prayer with detailed extraction
        """
        url = f"{self.base_url}breviario.php?s=lodi"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            full_text = soup.get_text()
            
            # Find the main invitatory antiphon (Antiphon 1)
            antiphon_1 = None
            invitatory_pattern = r'(Come, worship the Lord[^.]*alleluia\.)'
            invitatory_match = re.search(invitatory_pattern, full_text, re.IGNORECASE)
            if invitatory_match:
                antiphon_1 = invitatory_match.group(1).strip()
            
            # Find the first psalm verse that follows (priest's response)
            priest_response = None
            # Look for substantial psalm content after the invitatory
            psalm_patterns = [
                r'Come, let us sing to the Lord[^.]*\.',
                r'Let us approach him with praise[^.]*\.',
                r'The Lord is God, the mighty God[^.]*\.',
            ]
            
            for pattern in psalm_patterns:
                match = re.search(pattern, full_text, re.IGNORECASE)
                if match:
                    priest_response = match.group(0).strip()
                    break
            
            # If no specific match, get a meaningful psalm verse
            if not priest_response:
                psalm_verse_pattern = r'([A-Z][^.]*Lord[^.]*\.)'
                psalm_matches = list(re.finditer(psalm_verse_pattern, full_text))
                for match in psalm_matches:
                    verse = match.group(1).strip()
                    if (len(verse) > 30 and len(verse) < 200 and 
                        'Lord' in verse and 'Come' not in verse):
                        priest_response = verse
                        break
            
            return {
                'antiphon_1': antiphon_1,
                'priest_response': priest_response,
                'date': datetime.now().strftime('%B %d, %Y')
            }
            
        except Exception as e:
            print(f"Error fetching morning prayer: {e}")
            return None

    def fetch_daily_readings_detailed(self):
        """
        Fetch daily readings with better parsing
        """
        url = f"{self.base_url}letture.php?s=letture"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            full_text = soup.get_text()
            
            readings = {}
            
            # Extract first reading
            first_reading_match = re.search(r'First Reading(.{50,1500}?)(?=Responsorial Psalm|Second Reading|Gospel)', full_text, re.IGNORECASE | re.DOTALL)
            if first_reading_match:
                content = first_reading_match.group(1).strip()
                content = re.sub(r'\s+', ' ', content)
                # Clean up and extract the actual reading text
                reading_match = re.search(r'A reading from[^.]*\.(.+)', content, re.DOTALL)
                if reading_match:
                    readings['first_reading'] = reading_match.group(1).strip()[:1000]
                else:
                    readings['first_reading'] = content[:800]
            
            # Extract Gospel
            gospel_match = re.search(r'Gospel(.{50,1500})$', full_text, re.IGNORECASE | re.DOTALL)
            if gospel_match:
                content = gospel_match.group(1).strip()
                content = re.sub(r'\s+', ' ', content)
                # Clean up
                gospel_reading_match = re.search(r'A reading from[^.]*\.(.+)', content, re.DOTALL)
                if gospel_reading_match:
                    readings['gospel'] = gospel_reading_match.group(1).strip()[:1000]
                else:
                    readings['gospel'] = content[:800]
            
            return readings
            
        except Exception as e:
            print(f"Error fetching readings: {e}")
            return {}

    def create_enhanced_slides(self, prayer_data, readings_data, output_filename=None):
        """
        Create PowerPoint slides with enhanced formatting for elderly congregation
        """
        # Generate filename with date if not provided
        if output_filename is None:
            date_str = datetime.now().strftime("%Y-%m-%d")
            output_filename = f"enhanced_daily_mass_slides_{date_str}.pptx"
        prs = Presentation()
        
        # Set slide dimensions for widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 1. Title slide
        self._add_enhanced_title_slide(prs, prayer_data.get('date', ''))
        
        # 2. Antiphon 1 slide (Blue - audience and priest read together)
        if prayer_data and prayer_data.get('antiphon_1'):
            self._add_antiphon_slide(prs, prayer_data['antiphon_1'])
            
            # 3. Priest response slide (Red - priest only)
            if prayer_data.get('priest_response'):
                self._add_priest_response_slide(prs, prayer_data['priest_response'])
        
        # 4. First Reading slide
        if readings_data and readings_data.get('first_reading'):
            self._add_reading_slide(prs, "First Reading", readings_data['first_reading'])
        
        # 5. Gospel slide
        if readings_data and readings_data.get('gospel'):
            self._add_reading_slide(prs, "Gospel", readings_data['gospel'])
        
        # Ensure output directory exists
        output_dir = "output"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # Create full path for output file
        output_path = os.path.join(output_dir, output_filename)
        
        # Remove existing file if it exists (replacing old file with new for same date)
        if os.path.exists(output_path):
            os.remove(output_path)
            print(f"Replaced existing file: {output_path}")
        
        # Save the presentation
        prs.save(output_path)
        print(f"Enhanced slides saved as: {output_path}")

    def _add_enhanced_title_slide(self, prs, date_str):
        """Enhanced title slide with large, readable text"""
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "Daily Catholic Mass"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(60)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(0, 0, 100)  # Dark blue
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Morning Prayer and Readings"
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(40)
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
        date_frame = date_box.text_frame
        date_frame.text = date_str
        date_paragraph = date_frame.paragraphs[0]
        date_paragraph.font.size = Pt(32)
        date_paragraph.alignment = PP_ALIGN.CENTER
        date_paragraph.font.italic = True

    def _add_antiphon_slide(self, prs, antiphon_text):
        """Add Antiphon 1 slide (Blue text - congregation and priest together)"""
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Antiphon 1"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(0, 0, 150)
        
        # Instruction
        instruction_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(0.8))
        instruction_frame = instruction_box.text_frame
        instruction_frame.text = "(Audience and Priest together)"
        instruction_paragraph = instruction_frame.paragraphs[0]
        instruction_paragraph.font.size = Pt(24)
        instruction_paragraph.alignment = PP_ALIGN.CENTER
        instruction_paragraph.font.italic = True
        instruction_paragraph.font.color.rgb = RGBColor(100, 100, 100)
        
        # Antiphon text
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = antiphon_text
        content_paragraph = content_frame.paragraphs[0]
        content_paragraph.font.size = Pt(40)  # Large text for elderly
        content_paragraph.font.color.rgb = RGBColor(0, 100, 200)  # Blue
        content_paragraph.alignment = PP_ALIGN.CENTER
        content_paragraph.font.bold = True

    def _add_priest_response_slide(self, prs, response_text):
        """Add priest response slide (Red text - priest only)"""
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Priest Response"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(150, 0, 0)
        
        # Instruction
        instruction_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(0.8))
        instruction_frame = instruction_box.text_frame
        instruction_frame.text = "(Priest only)"
        instruction_paragraph = instruction_frame.paragraphs[0]
        instruction_paragraph.font.size = Pt(24)
        instruction_paragraph.alignment = PP_ALIGN.CENTER
        instruction_paragraph.font.italic = True
        instruction_paragraph.font.color.rgb = RGBColor(100, 100, 100)
        
        # Response text
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(4))
        content_frame = content_box.text_frame
        content_frame.text = response_text
        content_paragraph = content_frame.paragraphs[0]
        content_paragraph.font.size = Pt(36)  # Large text for elderly
        content_paragraph.font.color.rgb = RGBColor(200, 0, 0)  # Red
        content_paragraph.alignment = PP_ALIGN.CENTER
        content_paragraph.font.bold = True

    def _add_reading_slide(self, prs, reading_title, reading_text):
        """Add reading slides with large, readable text"""
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = reading_title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(0, 0, 100)
        
        # Reading text
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = reading_text
        content_frame.word_wrap = True
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(28)  # Large, readable text
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(12)  # Add space between paragraphs

def main():
    print("Enhanced Catholic Slide Generator")
    print("=" * 40)
    
    generator = EnhancedCatholicSlideGenerator()
    
    print("Fetching morning prayer content...")
    prayer_data = generator.fetch_morning_prayer_detailed()
    
    if prayer_data:
        print(f"✓ Found Antiphon 1: {prayer_data['antiphon_1'][:50]}...")
        print(f"✓ Found Priest Response: {prayer_data['priest_response'][:50]}...")
    
    print("\nFetching daily readings...")
    readings_data = generator.fetch_daily_readings_detailed()
    
    if readings_data:
        if 'first_reading' in readings_data:
            print(f"✓ Found First Reading: {readings_data['first_reading'][:50]}...")
        if 'gospel' in readings_data:
            print(f"✓ Found Gospel: {readings_data['gospel'][:50]}...")
    
    print("\nCreating enhanced slides...")
    generator.create_enhanced_slides(prayer_data, readings_data)
    print("✓ Enhanced slides created successfully!")

if __name__ == "__main__":
    main()