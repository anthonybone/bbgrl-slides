"""
BBGRL Slide Generator V2.py
Complete recreation of the 11_10 Morning Readings & Prayers structure
Based on the extracted PowerPoint analysis

This version codifies the exact structure to recreate the presentation 1:1
"""

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
from datetime import datetime, timedelta
import sys
import os

class BBGRLSlideGeneratorV2:
    def __init__(self):
        self.base_url = "https://www.ibreviary.com/m2/"
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })

    def get_complete_liturgical_structure(self):
        """
        Get the complete liturgical structure based on the extracted PowerPoint
        Total: 135 slides with exact content and sequence
        """
        
        structure = {
            "metadata": {
                "total_slides": 135,
                "date": datetime.now().strftime('%B %d, %Y'),
                "title": "Morning Readings & Prayers"
            },
            
            # Slides 1-2: Title/Opening slides (no text content)
            "opening_slides": [
                {"slide_num": 1, "content": ""},  # Title slide
                {"slide_num": 2, "content": ""}   # Blank/transition slide
            ],
            
            # Slides 3-36: PSALMODY Section
            "psalmody_section": {
                "slides": [
                    {"slide_num": 3, "title": "PSALMODY", "content": "(All) Ant. 1 Each morning, Lord, you fill us with your kindness.\nPsalm 90\nMay we live in the radiance of God"},
                    {"slide_num": 4, "speaker": "Priest", "content": "O Lord, you have been our refuge from one generation to the next. Before the mountains were born or the earth or the world brought forth, you are God, without beginning or end."},
                    {"slide_num": 5, "speaker": "People", "content": "You turn men back to dust and say: \"Go back, sons of men.\" To your eyes a thousand years are like yesterday, come and gone, no more than a watch in the night."},
                    {"slide_num": 6, "speaker": "Priest", "content": "You sweep men away like a dream, like the grass which springs up in the morning. In the morning it springs up and flowers: by evening it withers and fades."},
                    {"slide_num": 7, "speaker": "People", "content": "So we are destroyed in your anger, struck with terror in your fury. Our guilt lies open before you; our secrets in the light of your face."},
                    {"slide_num": 8, "speaker": "Priest", "content": "All our days pass away in your anger. Our life is over like a sigh. Our span is seventy years or eighty for those who are strong."},
                    {"slide_num": 9, "speaker": "People", "content": "And most of these are emptiness and pain. They pass swiftly and we are gone. Who understands the power of your anger and fears the strength of your fury?"},
                    {"slide_num": 10, "speaker": "Priest", "content": "Make us know the shortness of our life that we may gain wisdom of heart. Lord, relent! Is your anger for ever? Show pity to your servants."},
                    {"slide_num": 11, "speaker": "People", "content": "In the morning, fill us with your love; we shall exult and rejoice all our days. Give us joy to balance our affliction for the years when we knew misfortune."},
                    {"slide_num": 12, "speaker": "Priest", "content": "Show forth your work to your servants; let your glory shine on their children. Let the favor of the Lord be upon us: give success to the work of our hands, give success to the work of our hands."},
                    {"slide_num": 13, "speaker": "All", "content": "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be for ever. Amen."},
                    {"slide_num": 14, "speaker": "All", "content": "(All) Ant. Each morning, Lord, you fill us with your kindness."},
                    
                    # Antiphon 2 and Canticle (Slides 15-25)
                    {"slide_num": 15, "title": "Ant. 2", "content": "From the farthest bounds of earth, may God be praised!"},
                    {"slide_num": 16, "title": "Canticle: Isaiah 42:10-16", "content": "God, victor and savior"},
                    {"slide_num": 17, "speaker": "Priest", "content": "Sing to the Lord a new song, his praise from the end of the earth:"},
                    {"slide_num": 18, "speaker": "People", "content": "Let the sea and what fills it resound, the coastlands, and those who dwell in them. Let the steppe and its cities cry out, the villages where Kedar dwells;"},
                    {"slide_num": 19, "speaker": "Priest", "content": "let the inhabitants of Sela exult, and shout from the top of the mountains. Let them give glory to the Lord, and utter his praise in the coastlands."},
                    {"slide_num": 20, "speaker": "People", "content": "The Lord goes forth like a hero, like a warrior he stirs up his ardor; he shouts out his battle cry, against his enemies he shows his might:"},
                    {"slide_num": 21, "speaker": "Priest", "content": "I have looked away, and kept silence, I have said nothing, holding myself in; but now, I cry out as a woman in labor, gasping and panting."},
                    {"slide_num": 22, "speaker": "People", "content": "I will lay waste mountains and hills, all their herbage I will dry up; I will turn the rivers into marshes, and the marshes I will dry up."},
                    {"slide_num": 23, "speaker": "Priest", "content": "I will lead the blind on their journey; by paths unknown I will guide them. I will turn darkness into light before them, and make crooked ways straight."},
                    {"slide_num": 24, "speaker": "All", "content": "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be for ever. Amen."},
                    {"slide_num": 25, "speaker": "All", "content": "(All) Ant. From the farthest bounds of earth, may God be praised!"},
                    
                    # Antiphon 3 and Psalm 135 (Slides 26-36)
                    {"slide_num": 26, "speaker": "All", "content": "(All) Ant. 3 You who stand in his sanctuary, praise the name of the Lord"},
                    {"slide_num": 27, "title": "Psalm 135:1-12", "content": "Praise for the wonderful things God does for us"},
                    {"slide_num": 28, "speaker": "Priest", "content": "Praise the name of the Lord, praise him, servants of the Lord, who stand in the house of the Lord, in the courts of the house of our God."},
                    {"slide_num": 29, "speaker": "People", "content": "Praise the Lord for the Lord is good. Sing a psalm to his name for he is loving. For the Lord has chosen Jacob for himself and Israel for his own possession."},
                    {"slide_num": 30, "speaker": "Priest", "content": "For I know the Lord is great, that our Lord is high above all gods. The Lord does whatever he wills, in heaven, on earth, in the seas."},
                    {"slide_num": 31, "speaker": "People", "content": "He summons clouds from the ends of the earth; makes lightning produce the rain; from his treasuries he sends forth the wind."},
                    {"slide_num": 32, "speaker": "Priest", "content": "The first-born of the Egyptians he smote, of man and beast alike. Signs and wonders he worked in the midst of your land, O Egypt, against Pharaoh and all his servants."},
                    {"slide_num": 33, "speaker": "People", "content": "Nations in their greatness he struck and kings in their splendor he slew. Sihon, king of the Amorites,"},
                    {"slide_num": 34, "speaker": "People", "content": "Og, the king of Bashan, and all the kingdoms of Canaan. He let Israel inherit their land; on his people their land he bestowed."},
                    {"slide_num": 35, "speaker": "All", "content": "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be for ever. Amen."},
                    {"slide_num": 36, "speaker": "All", "content": "(All) Ant. You who stand in his sanctuary, praise the name of the Lord."}
                ]
            },
            
            # Slides 37-42: READING and RESPONSORY
            "reading_section": {
                "slides": [
                    {"slide_num": 37, "title": "READING", "subtitle": "[Doctors] Wisdom 7:13-14", "content": "Simply I learned about Wisdom, and ungrudgingly do I share— her riches I do not hide away;"},
                    {"slide_num": 38, "content": "For to men she is an unfailing treasure; those who gain this treasure win the friendship of God,"},
                    {"slide_num": 39, "content": "to whom the gifts they have from discipline commend them."},
                    {"slide_num": 40, "title": "RESPONSORY", "content": "Let the peoples proclaim the wisdom of the saints. — Let the peoples proclaim the wisdom of the saints."},
                    {"slide_num": 41, "content": "With joyful praise let the Church tell forth — the wisdom of the saints."},
                    {"slide_num": 42, "content": "Glory to the Father, and to the Son, and to the Holy Spirit. — Let the peoples proclaim the wisdom of the saints."}
                ]
            },
            
            # Slides 43-54: GOSPEL CANTICLE
            "gospel_canticle_section": {
                "slides": [
                    {"slide_num": 43, "title": "GOSPEL CANTICLE", "content": "Ant. Strengthened by Christ, blessed Peter has remained steadfast as a rock in his guidance of the Church."},
                    {"slide_num": 44, "title": "Canticle of Zechariah", "subtitle": "Luke 1:68-79"},
                    {"slide_num": 45, "content": "Blessed be the Lord, the God of Israel; he has come to his people and set them free."},
                    {"slide_num": 46, "content": "He has raised up for us a mighty savior, born of the house of his servant David."},
                    {"slide_num": 47, "content": "Through his holy prophets he promised of old that he would save us from our enemies, from the hands of all who hate us."},
                    {"slide_num": 48, "content": "He promised to show mercy to our fathers and to remember his holy covenant."},
                    {"slide_num": 49, "content": "This was the oath he swore to our father Abraham: to set us free from the hands of our enemies, free to worship him without fear, holy and righteous in his sight all the days of our life."},
                    {"slide_num": 50, "content": "You, my child, shall be called the prophet of the Most High; for you will go before the Lord to prepare"},
                    {"slide_num": 51, "content": "his way, to give his people knowledge of salvation by the forgiveness of their sins."},
                    {"slide_num": 52, "content": "In the tender compassion of our God the dawn from on high shall break upon us, to shine on those who dwell in darkness and the shadow of death, and to guide our feet into the way of peace."},
                    {"slide_num": 53, "content": "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be forever. Amen."},
                    {"slide_num": 54, "content": "Ant. Strengthened by Christ, blessed Peter has remained steadfast as a rock in his guidance of the Church."}
                ]
            },
            
            # Slides 55-63: INTERCESSIONS and CONCLUDING PRAYER
            "intercessions_section": {
                "slides": [
                    {"slide_num": 55, "title": "INTERCESSIONS", "content": "Christ is the Good Shepherd who laid down his life for his sheep. Let us praise and thank him as we pray:"},
                    {"slide_num": 56, "content": "(All): Nourish your people, Lord."},
                    {"slide_num": 57, "content": "Christ, you decided to show your merciful love through your holy shepherds, — let your mercy always reach us through them."},
                    {"slide_num": 58, "content": "Through your vicars you continue to perform the ministry of shepherd of souls, — direct us always through our leaders."},
                    {"slide_num": 59, "content": "Through your holy ones, the leaders of your people, you served as physician of our bodies and our spirits, — continue to fulfill your ministry of life and holiness in us."},
                    {"slide_num": 60, "content": "You taught your flock through the prudence and love of your saints, — grant us continual growth in holiness under the direction of our pastors."},
                    {"slide_num": 61, "title": "THE LORD'S PRAYER"},
                    {"slide_num": 62, "title": "CONCLUDING PRAYER", "content": "Priest: God our Father you will never allow the power of hell to prevail against your Church, founded on the rock of the apostle Peter. Let the prayers of Pope Leo the Great"},
                    {"slide_num": 63, "content": "keep us faithful to your truth and secure in your peace. We ask this through our Lord Jesus Christ, your Son, who lives and reigns with you and the Holy Spirit, God, for ever and ever. — Amen."}
                ]
            },
            
            # Slides 64-69: SACRED HEART HYMNS
            "sacred_heart_hymns": {
                "slides": [
                    {"slide_num": 64, "content": ""},  # Blank slide
                    {"slide_num": 65, "content": "Heart of Jesus meek and mild. Hear oh hear thy feeble child When the tempest's most severe. Heart of Jesus, hear. Sweetly, we'll rest on thy Sacred Heart. Never from Thee. O let us part!"},
                    {"slide_num": 66, "content": "Hear then, Thy loving children's pray'r O Heart of Jesus, Heart of Jesus hear."},
                    {"slide_num": 67, "content": ""},  # Blank slide
                    {"slide_num": 68, "content": "Oh Sacred Heart, Oh love divine. Do keep us near to Thee. And make our love so like to Thine That we may holy be. Heart of Jesus hear. Oh heart of love divine. Listen to our pray'r."},
                    {"slide_num": 69, "content": "Make us always Thine. Oh temple pure, Oh house of gold. Our heaven here below. What gifts unfurled, what wealth untold. From Thee do ever flow. Heart of Jesus hear. Oh Heart of love divine. Listen to our pray'r. Make us always Thine."}
                ]
            },
            
            # Slides 70-88: MASS READINGS
            "mass_readings": {
                "first_reading": {
                    "slides": [
                        {"slide_num": 70, "title": "FIRST READING", "subtitle": "Rom 9:1-5", "content": "A reading from the Letter of Saint Paul to the Romans\nBrothers and sisters: I speak the truth in Christ, I do not lie; my conscience joins with the Holy Spirit in bearing me witness"},
                        {"slide_num": 71, "content": "that I have great sorrow and constant anguish in my heart. For I could wish that I myself were accursed and cut off from Christ for the sake of my own people, my kindred according to the flesh."},
                        {"slide_num": 72, "content": "They are children of Israel; theirs the adoption, the glory, the covenants, the giving of the law, the worship, and the promises; theirs the patriarchs, and from them,"},
                        {"slide_num": 73, "content": "according to the flesh, is the Christ, who is over all, God blessed forever. Amen. The word of the word."}
                    ]
                },
                "responsorial_psalm": {
                    "slides": [
                        {"slide_num": 74, "title": "Responsorial Psalm", "subtitle": "Ps 147:12-13, 14-15, 19-20", "content": "℟. Praise the Lord, Jerusalem."},
                        {"slide_num": 75, "content": "Glorify the LORD, O Jerusalem; praise your God, O Zion. For he has strengthened the bars of your gates; he has blessed your children within you."},
                        {"slide_num": 76, "content": "℟. Praise the Lord, Jerusalem."},
                        {"slide_num": 77, "content": "He has granted peace in your borders; with the best of wheat he fills you. He sends forth his command to the earth; swiftly runs his word!"},
                        {"slide_num": 78, "content": "℟. Praise the Lord, Jerusalem."},
                        {"slide_num": 79, "content": "He has proclaimed his word to Jacob, his statutes and his ordinances to Israel. He has not done thus for any other nation;"},
                        {"slide_num": 80, "content": "his ordinances he has not made known to them. Alleluia."},
                        {"slide_num": 81, "content": "℟. Praise the Lord, Jerusalem."}
                    ]
                },
                "gospel_acclamation": {
                    "slides": [
                        {"slide_num": 82, "title": "ACCLAMATION BEFORE THE GOSPEL", "subtitle": "Jn 10:27", "content": "℟. Alleluia, alleluia."},
                        {"slide_num": 83, "content": "My sheep hear my voice, says the Lord; I know them, and they follow me.\n\n℟. Alleluia, alleluia."}
                    ]
                },
                "gospel": {
                    "slides": [
                        {"slide_num": 84, "title": "Gospel", "subtitle": "Lk 14:1-6", "content": "✠ A reading from the holy Gospel according to Luke\nOn a sabbath Jesus went to dine at the home of one of the leading Pharisees,"},
                        {"slide_num": 85, "content": "and the people there were observing him carefully. In front of him there was a man suffering from dropsy. Jesus spoke to the scholars of the law and Pharisees in reply, asking,"},
                        {"slide_num": 86, "content": "\"Is it lawful to cure on the sabbath or not?\" But they kept silent; so he took the man and, after he had healed him, dismissed him. Then he said to them \"Who among"},
                        {"slide_num": 87, "content": "you, if your son or ox falls into a cistern, would not immediately pull him out on the sabbath day?\" But they were unable to answer his question."},
                        {"slide_num": 88, "content": "At the end of the Gospel, the Priest, acclaims: The Gospel of the Lord. All reply: Praise to you, Lord Jesus Christ."}
                    ]
                }
            },
            
            # Slides 89-105: POST-COMMUNION PRAYERS
            "post_communion_prayers": {
                "slides": [
                    {"slide_num": 89, "content": ""},  # Blank slide
                    {"slide_num": 90, "content": "Soul of Christ, make me holy. Body of Christ, save me. Blood of Christ, inebriate me. Water from the side of Christ, wash me. Passion of Christ, make me strong. O good Jesus, hear me. Hide me within your wounds."},
                    {"slide_num": 91, "content": "Let me never be separated from You. Deliver me from the wicked enemy, Call me at the hour of my death. And tell me to come to you that with Your saints I may praise You forever. Amen."},
                    {"slide_num": 92, "title": "PRAYER OF THANKSGIVING:", "content": "Lord God, I thank you through the Sacred Heart of Jesus, who is pleased to offer You on our behalf continuous thanksgiving in the Eucharist."},
                    {"slide_num": 93, "content": "Thank You for all Your benefits, general as well as special: for creation, redemption, the sacraments; and especially for the Holy Eucharist, and for all graces You have given me."},
                    {"slide_num": 94, "content": "I wish to thank You, Lord Jesus, on behalf of myself and of all creatures, and to make amends to You, as far as I am able, for the ingratitude of many, of which You complain so vehemently."},
                    {"slide_num": 95, "content": "I wish that I was able to direct the hearts and minds of all to You and, together with them and for them, love You perfectly in return just as You rightly expect."},
                    {"slide_num": 96, "title": "NOVENA OF CONFIDENCE:", "content": "O Lord, Jesus Christ, to Your Most Sacred Heart I confide this intention..... (Mention your request)"},
                    {"slide_num": 97, "content": "Only look upon me, then do what Your heart inspires. Let Your Sacred Heart decide. I count on You. I trust in You. I throw myself on Your mercy."},
                    {"slide_num": 98, "content": "Lord Jesus! You will not fail me. Sacred Heart of Jesus, I believe in Your love for me. O Sacred Heart Of Jesus, Your kingdom come."},
                    {"slide_num": 99, "content": "O Sacred Heart of Jesus, I have asked for many favors, but I earnestly implore this one. Take it; place it in Your Sacred Heart. When the Eternal Father sees it covered with Your Precious Blood,"},
                    {"slide_num": 100, "content": "He will not refuse it. It will be no longer my prayer but Yours, O Jesus. O Sacred Heart of Jesus, I place my trust in you. Let me never be confounded. Amen."},
                    {"slide_num": 101, "title": "NOVENA PRAYER:", "content": "O most holy Heart of Jesus, fountain of every blessing, I adore You, I love You, and with a lively sorrow for my sins,"},
                    {"slide_num": 102, "content": "I offer You this poor heart of mine. Make me humble, patient, pure and wholly obedient to Your will. Grant, good Jesus that I may live in You and for You."},
                    {"slide_num": 103, "content": "Protect me in the midst of danger; comfort me in my afflictions; give me health of body, assistance in my temporal needs, Your blessing on all that I do, and the grace of a holy death....... Our Father…. Hail Mary.... Glory be to the Father....."},
                    {"slide_num": 104, "content": "Let us Pray, Heavenly Father, we rejoice in the gifts of love we have received from the Heart of Jesus, your Son."},
                    {"slide_num": 105, "content": "Open our hearts to share His life and continue to bless us with His love. We ask this in the name of Jesus the Lord. AMEN."}
                ]
            },
            
            # Slides 106-115: Blank transition slides
            "transition_slides": [
                {"slide_num": i, "content": ""} for i in range(106, 116)
            ],
            
            # Slides 116-122: JUBILEE PRAYER
            "jubilee_prayer": {
                "slides": [
                    {"slide_num": 116, "title": "THE JUBILEE PRAYER"},
                    {"slide_num": 117, "content": "Father in heaven, may the faith you have given us in your son, Jesus Christ, our brother, and the flame of charity"},
                    {"slide_num": 118, "content": "enkindled in our hearts by the Holy Spirit, reawaken in us the blessed hope for the coming of your Kingdom."},
                    {"slide_num": 119, "content": "May your grace transform us into tireless cultivators of the seeds of the Gospel. May those seeds transform from within both humanity and the whole"},
                    {"slide_num": 120, "content": "cosmos in the sure expectation of a new heaven and a new earth, when, with the powers of Evil vanquished, your glory will shine eternally."},
                    {"slide_num": 121, "content": "May the grace of the Jubilee reawaken in us, Pilgrims of Hope, a yearning for the treasures of heaven. May that same grace spread the"},
                    {"slide_num": 122, "content": "joy and peace of our Redeemer throughout the earth. To you our God, eternally blessed, be glory and praise for ever. Amen."}
                ]
            },
            
            # Slide 123: Blank
            "final_blank": {"slide_num": 123, "content": ""},
            
            # Slides 124-135: PRAYER TO ST. JOSEPH
            "st_joseph_prayer": {
                "slides": [
                    {"slide_num": 124, "content": "To you, O blessed Joseph, do we come in our tribulation, and having implored the help of your most holy Spouse, we confidently invoke your patronage also."},
                    {"slide_num": 125, "content": "Through that charity which bound you to the Immaculate Virgin Mother of God and through the"},
                    {"slide_num": 126, "content": "paternal love with which you embraced the Child Jesus, we humbly beg you graciously to regard the inheritance"},
                    {"slide_num": 127, "content": "which Jesus Christ has purchased by his Blood, and with your power and strength to aid us in our necessities"},
                    {"slide_num": 128, "content": "O most watchful guardian of the Holy Family, defend the chosen children of Jesus Christ;"},
                    {"slide_num": 129, "content": "O most loving father, ward off from us every contagion of error and corrupting influence;"},
                    {"slide_num": 130, "content": "O our most mighty protector, be kind to us and from heaven assist us in our struggle with the power of darkness"},
                    {"slide_num": 131, "content": "As once you rescued the Child Jesus from deadly peril, so now protect God's Holy Church from the snares of the"},
                    {"slide_num": 132, "content": "enemy and from all adversity; shield, too, each one of us by your constant protection,"},
                    {"slide_num": 133, "content": "so that, supported by your example and your aid, we may be able to live piously, to die in holiness, and to obtain eternal happiness in heaven. Amen."},
                    {"slide_num": 134, "content": "St. Joseph, husband of the Blessed Virgin Mary All praise, O God, for Joseph, The guardian of your Son, Who saved him from King Herod, When safety there was none."},
                    {"slide_num": 135, "content": "He taught the trade of builder, When they to Naz'reth came, And Joseph's love made \"Father\" To be, for Christ, God's name."}
                ]
            }
        }
        
        return structure

    def create_complete_presentation(self, output_filename=None):
        """
        Create the complete 135-slide presentation exactly as extracted
        """
        if output_filename is None:
            date_str = datetime.now().strftime("%m_%d")
            output_filename = f"{date_str} Morning Readings & Prayers .pptx"
        
        structure = self.get_complete_liturgical_structure()
        prs = Presentation()
        
        # Set slide dimensions for widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        print(f"Creating presentation: {structure['metadata']['title']}")
        print(f"Total slides to create: {structure['metadata']['total_slides']}")
        
        # Create all slides in sequence
        slide_count = 0
        
        # Opening slides (1-2)
        for slide_data in structure["opening_slides"]:
            slide_count += 1
            if slide_data["content"] == "":
                self._add_blank_slide(prs, slide_count)
            print(f"Created slide {slide_count}: Opening slide")
        
        # Psalmody section (3-36)
        for slide_data in structure["psalmody_section"]["slides"]:
            slide_count += 1
            self._add_liturgical_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Psalmody")
        
        # Reading section (37-42)
        for slide_data in structure["reading_section"]["slides"]:
            slide_count += 1
            self._add_liturgical_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Reading/Responsory")
        
        # Gospel Canticle section (43-54)
        for slide_data in structure["gospel_canticle_section"]["slides"]:
            slide_count += 1
            self._add_liturgical_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Gospel Canticle")
        
        # Intercessions section (55-63)
        for slide_data in structure["intercessions_section"]["slides"]:
            slide_count += 1
            self._add_liturgical_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Intercessions/Concluding Prayer")
        
        # Sacred Heart Hymns (64-69)
        for slide_data in structure["sacred_heart_hymns"]["slides"]:
            slide_count += 1
            if slide_data["content"] == "":
                self._add_blank_slide(prs, slide_count)
            else:
                self._add_hymn_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Sacred Heart Hymns")
        
        # Mass Readings (70-88)
        # First Reading
        for slide_data in structure["mass_readings"]["first_reading"]["slides"]:
            slide_count += 1
            self._add_reading_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: First Reading")
        
        # Responsorial Psalm
        for slide_data in structure["mass_readings"]["responsorial_psalm"]["slides"]:
            slide_count += 1
            self._add_reading_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Responsorial Psalm")
        
        # Gospel Acclamation
        for slide_data in structure["mass_readings"]["gospel_acclamation"]["slides"]:
            slide_count += 1
            self._add_reading_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Gospel Acclamation")
        
        # Gospel
        for slide_data in structure["mass_readings"]["gospel"]["slides"]:
            slide_count += 1
            self._add_reading_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Gospel")
        
        # Post-communion prayers (89-105)
        for slide_data in structure["post_communion_prayers"]["slides"]:
            slide_count += 1
            if slide_data["content"] == "":
                self._add_blank_slide(prs, slide_count)
            else:
                self._add_prayer_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Post-communion prayers")
        
        # Transition slides (106-115)
        for slide_data in structure["transition_slides"]:
            slide_count += 1
            self._add_blank_slide(prs, slide_count)
            print(f"Created slide {slide_count}: Transition slide")
        
        # Jubilee Prayer (116-122)
        for slide_data in structure["jubilee_prayer"]["slides"]:
            slide_count += 1
            self._add_prayer_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: Jubilee Prayer")
        
        # Final blank (123)
        slide_count += 1
        self._add_blank_slide(prs, slide_count)
        print(f"Created slide {slide_count}: Final blank")
        
        # St. Joseph Prayer (124-135)
        for slide_data in structure["st_joseph_prayer"]["slides"]:
            slide_count += 1
            self._add_prayer_slide(prs, slide_data, slide_count)
            print(f"Created slide {slide_count}: St. Joseph Prayer")
        
        # Save the presentation
        output_dir = "output"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        output_path = os.path.join(output_dir, output_filename)
        prs.save(output_path)
        
        print(f"\nPresentation created successfully!")
        print(f"File: {output_path}")
        print(f"Total slides created: {slide_count}")
        
        return output_path

    def _add_blank_slide(self, prs, slide_num):
        """Add a blank slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

    def _add_liturgical_slide(self, prs, slide_data, slide_num):
        """Add a liturgical slide with proper formatting"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Handle different slide types
        if "title" in slide_data:
            # Title slide
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(48)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)
            title_para.alignment = PP_ALIGN.CENTER
            
            if "subtitle" in slide_data:
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide_data["subtitle"]
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(32)
                subtitle_para.font.color.rgb = RGBColor(0, 51, 102)
                subtitle_para.alignment = PP_ALIGN.CENTER
            
            if slide_data.get("content"):
                content_y = 3.5 if "subtitle" in slide_data else 2.5
                content_box = slide.shapes.add_textbox(Inches(1), Inches(content_y), Inches(11.33), Inches(3))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = slide_data["content"]
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(28)
                content_para.alignment = PP_ALIGN.CENTER
        else:
            # Content slide
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = slide_data["content"]
            
            # Color coding based on speaker
            speaker = slide_data.get("speaker", "")
            if speaker == "Priest":
                color = RGBColor(200, 0, 0)  # Red for priest
                font_size = Pt(32)
            elif speaker == "People":
                color = RGBColor(0, 100, 200)  # Blue for people
                font_size = Pt(32)
            elif speaker == "All":
                color = RGBColor(100, 0, 100)  # Purple for all
                font_size = Pt(36)
            else:
                color = RGBColor(0, 51, 102)  # Dark blue default
                font_size = Pt(30)
            
            content_para = content_frame.paragraphs[0]
            content_para.font.size = font_size
            content_para.font.bold = True
            content_para.font.color.rgb = color
            content_para.alignment = PP_ALIGN.CENTER

    def _add_reading_slide(self, prs, slide_data, slide_num):
        """Add a reading slide (Mass readings)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        if "title" in slide_data:
            # Title reading slide
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(48)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(100, 0, 0)
            title_para.alignment = PP_ALIGN.CENTER
            
            if "subtitle" in slide_data:
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide_data["subtitle"]
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(32)
                subtitle_para.font.color.rgb = RGBColor(100, 0, 0)
                subtitle_para.alignment = PP_ALIGN.CENTER
                
                content_y = 3
            else:
                content_y = 2
            
            if slide_data.get("content"):
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(content_y), Inches(12.33), Inches(4))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = slide_data["content"]
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(28)
                content_para.alignment = PP_ALIGN.LEFT
        else:
            # Regular reading content
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = slide_data["content"]
            content_para = content_frame.paragraphs[0]
            content_para.font.size = Pt(30)
            content_para.alignment = PP_ALIGN.LEFT

    def _add_hymn_slide(self, prs, slide_data, slide_num):
        """Add a hymn slide with special formatting"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = slide_data["content"]
        
        content_para = content_frame.paragraphs[0]
        content_para.font.size = Pt(32)
        content_para.font.color.rgb = RGBColor(139, 0, 0)  # Sacred Heart red
        content_para.alignment = PP_ALIGN.CENTER
        content_para.font.italic = True

    def _add_prayer_slide(self, prs, slide_data, slide_num):
        """Add a prayer slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        if "title" in slide_data:
            # Prayer title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(139, 0, 0)
            title_para.alignment = PP_ALIGN.CENTER
            
            if slide_data.get("content"):
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(4))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.text = slide_data["content"]
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(28)
                content_para.alignment = PP_ALIGN.CENTER
        else:
            # Prayer content only
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.text = slide_data["content"]
            content_para = content_frame.paragraphs[0]
            content_para.font.size = Pt(30)
            content_para.alignment = PP_ALIGN.CENTER

def main():
    print("BBGRL Slide Generator V2")
    print("=" * 40)
    print("Recreating the exact 135-slide structure")
    
    generator = BBGRLSlideGeneratorV2()
    
    # Generate the current date version
    date_str = datetime.now().strftime("%m_%d")
    output_filename = f"{date_str} Morning Readings & Prayers .pptx"
    
    print(f"Creating: {output_filename}")
    
    # Create the complete presentation
    output_path = generator.create_complete_presentation(output_filename)
    
    print("✓ Complete presentation created successfully!")
    print(f"✓ File saved as: {output_path}")

if __name__ == "__main__":
    main()