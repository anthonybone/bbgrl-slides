"""Tests for responsory slide formatting (no speaker labels).

Ensures the generator creates three slides matching the expected pattern:
1. Title 'RESPONSORY' + four content lines (including em-dash repeat)
2. Two-line stanza with em-dash response
3. Glory stanza with em-dash repeated response.
"""

from pptx import Presentation

from bbgrl.generator.generator import bbgrlslidegeneratorv1


def test_responsory_format_no_speaker_labels():
    verses = [
        {
            "include_title": True,
            "text": (
                "The just are the friends of God,\n"
                "They live with him for ever.\n"
                "— The just are the friends of God,\n"
                "They live with him for ever."
            ),
        },
        {
            "text": (
                "God himself is their reward.\n"
                "— They live with him for ever."
            ),
        },
        {
            "text": (
                "Glory to the Father, and to the Son, and to the Holy Spirit.\n"
                "— The just are the friends of God,\n"
                "They live with him for ever."
            ),
        },
    ]

    liturgical_data = {"morning_prayer": {"reading": {"responsory": verses}}}
    gen = bbgrlslidegeneratorv1()
    prs = Presentation()
    start_count = 0
    end_count = gen._create_responsory_section(prs, liturgical_data, start_count)

    assert end_count - start_count == 3, "Should create exactly three responsory slides"
    assert len(prs.slides) == 3, "Presentation should have three slides created"

    # Slide 1: should have title 'RESPONSORY' and no speaker labels
    slide1_texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
    assert any(t.strip() == "RESPONSORY" for t in slide1_texts), "First slide must have RESPONSRY title"
    combined_slide1 = "\n".join(slide1_texts)
    assert "(All)" not in combined_slide1 and "(Priest)" not in combined_slide1, "No speaker labels expected"
    assert "The just are the friends of God," in combined_slide1
    assert "They live with him for ever." in combined_slide1
    assert "— The just are the friends of God," in combined_slide1

    # Slide 2 content
    slide2_texts = [shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text")]
    combined_slide2 = "\n".join(slide2_texts)
    assert "God himself is their reward." in combined_slide2
    assert "— They live with him for ever." in combined_slide2

    # Slide 3 content
    slide3_texts = [shape.text for shape in prs.slides[2].shapes if hasattr(shape, "text")]
    combined_slide3 = "\n".join(slide3_texts)
    assert "Glory to the Father, and to the Son" in combined_slide3
    assert "— The just are the friends of God," in combined_slide3
    assert "They live with him for ever." in combined_slide3
