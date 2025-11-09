import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
from datetime import datetime, timedelta
import sys
import os

class CatholicSlideGenerator:
    def __init__(self):
        self.base_url = "https://www.ibreviary.com/m2/"
        self.session = requests.Session()
        # Set headers to mimic a browser
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })

    def fetch_morning_prayer(self, date=None):
        """
        Fetch the morning prayer (Laudes) for a specific date
        If no date is provided, uses today's date
        """
        if date is None:
            date = datetime.now()
        
        # Format date for the URL (if needed - the site might use current date by default)
        url = f"{self.base_url}breviario.php?s=lodi"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Extract the morning prayer content
            return self._parse_morning_prayer(soup)
        except requests.exceptions.RequestException as e:
            print(f"Error fetching morning prayer: {e}")
            return None

    def fetch_daily_readings(self, date=None):
        """
        Fetch the daily readings for Mass
        """
        if date is None:
            date = datetime.now()
            
        url = f"{self.base_url}letture.php?s=letture"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            return self._parse_daily_readings(soup)
        except requests.exceptions.RequestException as e:
            print(f"Error fetching daily readings: {e}")
            return None

    def _parse_morning_prayer(self, soup):
        """
        Parse the morning prayer HTML to extract relevant content
        """
        prayer_data = {
            'antiphons': [],
            'psalms': [],
            'readings': [],
            'canticles': []
        }
        
        # Get the main content
        main_content = soup.get_text()
        
        # Look for specific antiphon patterns based on iBreviary structure
        # Pattern 1: "Ant. 1" followed by text
        antiphon_pattern1 = r'Ant\.\s*(\d+)\s+([^.]+(?:\.[^A-Z]*)*)'
        # Pattern 2: "Ant." followed by text (for main antiphons)
        antiphon_pattern2 = r'Ant\.\s+([^.]+\.(?:[^A-Z][^.]*\.)*)'
        
        # Find numbered antiphons first
        antiphon_matches1 = re.finditer(antiphon_pattern1, main_content, re.DOTALL)
        for match in antiphon_matches1:
            antiphon_number = match.group(1)
            antiphon_text = match.group(2).strip()
            if antiphon_text and len(antiphon_text) > 20 and len(antiphon_text) < 300:
                # Clean up the text
                antiphon_text = re.sub(r'\s+', ' ', antiphon_text)
                prayer_data['antiphons'].append({
                    'number': antiphon_number,
                    'text': antiphon_text
                })
        
        # If no numbered antiphons found, look for repeated antiphons
        if not prayer_data['antiphons']:
            # Look for the repeating invitatory antiphon
            invitatory_pattern = r'(Come, worship the Lord[^.]*alleluia\.)'
            matches = re.finditer(invitatory_pattern, main_content, re.IGNORECASE)
            antiphon_texts = []
            for match in matches:
                text = match.group(1).strip()
                if text not in antiphon_texts:  # Avoid duplicates
                    antiphon_texts.append(text)
            
            if antiphon_texts:
                prayer_data['antiphons'].append({
                    'number': '1',
                    'text': antiphon_texts[0]
                })
        
        # Extract psalm verses that follow the antiphons
        # Look for verses from psalms (usually start with capitals and have poetic structure)
        psalm_verse_pattern = r'([A-Z][^.!?]*[.!?])'
        potential_verses = re.findall(psalm_verse_pattern, main_content)
        
        # Filter for psalm-like content (usually has religious language)
        psalm_keywords = ['Lord', 'God', 'praise', 'glory', 'heaven', 'earth', 'blessed', 'holy']
        for verse in potential_verses:
            if (len(verse) > 30 and len(verse) < 200 and 
                any(keyword in verse for keyword in psalm_keywords) and
                len(prayer_data['psalms']) < 3):
                prayer_data['psalms'].append(verse.strip())
        
        return prayer_data

    def _parse_daily_readings(self, soup):
        """
        Parse the daily readings HTML
        """
        readings_data = {
            'first_reading': None,
            'responsorial_psalm': None,
            'second_reading': None,
            'gospel': None
        }
        
        # Get all text content
        full_text = soup.get_text()
        
        # Look for reading citations and content
        # First Reading pattern
        first_reading_match = re.search(r'First Reading([^G]+)(?=Second Reading|Responsorial Psalm|Gospel|$)', full_text, re.IGNORECASE | re.DOTALL)
        if first_reading_match:
            reading_text = first_reading_match.group(1).strip()
            # Clean up and extract meaningful content
            reading_text = re.sub(r'\s+', ' ', reading_text)
            # Look for the actual reading content (usually after the citation)
            content_match = re.search(r'A reading from[^.]*\.(.+)', reading_text, re.DOTALL)
            if content_match:
                readings_data['first_reading'] = content_match.group(1).strip()[:1000]
            else:
                readings_data['first_reading'] = reading_text[:800]

        # Gospel pattern
        gospel_match = re.search(r'Gospel([^$]+)', full_text, re.IGNORECASE | re.DOTALL)
        if gospel_match:
            gospel_text = gospel_match.group(1).strip()
            gospel_text = re.sub(r'\s+', ' ', gospel_text)
            # Look for the actual gospel content
            content_match = re.search(r'A reading from[^.]*\.(.+)', gospel_text, re.DOTALL)
            if content_match:
                readings_data['gospel'] = content_match.group(1).strip()[:1000]
            else:
                readings_data['gospel'] = gospel_text[:800]
        
        return readings_data

    def create_slides(self, prayer_data, readings_data, output_filename=None):
        """
        Create PowerPoint slides with the prayer and reading content
        """
        # Generate filename with date if not provided
        if output_filename is None:
            date_str = datetime.now().strftime("%Y-%m-%d")
            output_filename = f"daily_mass_slides_{date_str}.pptx"
        prs = Presentation()
        
        # Set slide dimensions for better readability
        prs.slide_width = Inches(13.33)  # Widescreen format
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs)
        
        # Morning prayer slides
        if prayer_data and prayer_data['antiphons']:
            self._add_morning_prayer_slides(prs, prayer_data)
        
        # Daily readings slides
        if readings_data:
            self._add_reading_slides(prs, readings_data)
        
        # Ensure output directory exists
        output_dir = "output"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # Create full path for output file
        output_path = os.path.join(output_dir, output_filename)
        
        # Remove existing file if it exists (replacing old file with new for same date)
        if os.path.exists(output_path):
            os.remove(output_path)
            print(f"Replaced existing file: {output_path}")
        
        # Save the presentation
        prs.save(output_path)
        print(f"Slides saved as: {output_path}")

    def _add_title_slide(self, prs):
        """Add title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Daily Catholic Mass"
        subtitle.text = f"Morning Prayer and Readings\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Format title for large text
        title.text_frame.paragraphs[0].font.size = Pt(54)
        subtitle.text_frame.paragraphs[0].font.size = Pt(36)

    def _add_morning_prayer_slides(self, prs, prayer_data):
        """Add morning prayer slides with color coding"""
        # Add antiphon slides
        for i, antiphon in enumerate(prayer_data['antiphons'][:3]):  # Limit to first 3
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = f"Antiphon {antiphon['number']}"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            
            # Content
            content = slide.placeholders[1]
            content.text = antiphon['text']
            
            # Format the text - Blue for antiphons (audience and priest read together)
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.color.rgb = RGBColor(0, 100, 200)  # Blue
                paragraph.alignment = PP_ALIGN.CENTER
            
            # Add instruction slide for priest's response (if this is the first antiphon)
            if i == 0 and len(prayer_data['psalms']) > 0:
                response_slide = prs.slides.add_slide(prs.slide_layouts[1])
                response_title = response_slide.shapes.title
                response_title.text = "Priest Response"
                response_title.text_frame.paragraphs[0].font.size = Pt(48)
                
                response_content = response_slide.placeholders[1]
                # Use first part of the psalm as the priest's response
                psalm_text = prayer_data['psalms'][0][:500] + "..." if len(prayer_data['psalms']) > 0 else "..."
                response_content.text = psalm_text
                
                # Format in red for priest-only text
                for paragraph in response_content.text_frame.paragraphs:
                    paragraph.font.size = Pt(32)
                    paragraph.font.color.rgb = RGBColor(200, 0, 0)  # Red
                    paragraph.alignment = PP_ALIGN.CENTER

    def _add_reading_slides(self, prs, readings_data):
        """Add daily reading slides"""
        if readings_data['first_reading']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "First Reading"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            
            content = slide.placeholders[1]
            # Truncate if too long for readability
            reading_text = readings_data['first_reading'][:800] + "..." if len(readings_data['first_reading']) > 800 else readings_data['first_reading']
            content.text = reading_text
            
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(28)
                paragraph.alignment = PP_ALIGN.LEFT

        if readings_data['gospel']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Gospel"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            
            content = slide.placeholders[1]
            gospel_text = readings_data['gospel'][:800] + "..." if len(readings_data['gospel']) > 800 else readings_data['gospel']
            content.text = gospel_text
            
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(28)
                paragraph.alignment = PP_ALIGN.LEFT

def main():
    generator = CatholicSlideGenerator()
    
    print("Fetching morning prayer...")
    prayer_data = generator.fetch_morning_prayer()
    
    print("Fetching daily readings...")
    readings_data = generator.fetch_daily_readings()
    
    if prayer_data or readings_data:
        print("Creating slides...")
        generator.create_slides(prayer_data, readings_data)
        print("Slides created successfully!")
    else:
        print("Failed to fetch content. Please check your internet connection.")

if __name__ == "__main__":
    main()