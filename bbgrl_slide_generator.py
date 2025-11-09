import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
from datetime import datetime, timedelta
import sys
import os

class EnhancedCatholicSlideGenerator:
    def __init__(self):
        self.base_url = "https://www.ibreviary.com/m2/"
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })

    def get_liturgy_of_hours_structure(self, is_first_hour=True):
        """
        Get the proper Liturgy of the Hours structure for Morning Prayer
        """
        liturgy_structure = {
            "date": datetime.now().strftime('%B %d, %Y'),
            "is_first_hour": is_first_hour,
            "invitatory": None,
            "morning_prayer": None
        }
        
        # Invitatory (only if it's the first hour of the day)
        if is_first_hour:
            liturgy_structure["invitatory"] = {
                "sequence": [
                    "Opening Verse: Lord, open my lips.",
                    "Response: And my mouth will proclaim your praise.",
                    "Invitatory Antiphon",
                    "Invitatory Psalm (95, 100, 67, or 24)",
                    "Repeat Invitatory Antiphon"
                ]
            }
        
        # Morning Prayer structure
        liturgy_structure["morning_prayer"] = {
            "sequence": [
                "Opening Verse: God, come to my assistance.",
                "Response: Lord, make haste to help me.",
                "Glory to the Father, and to the Son, and to the Holy Spirit.",
                "As it was in the beginning, is now, and will be forever. Amen.",
                "Optional: Alleluia (omit during Lent)",
                "Hymn",
                "Antiphon 1",
                "Psalm 1",
                "Glory to the Father",
                "Repeat Antiphon 1",
                "Antiphon 2", 
                "Psalm 2",
                "Glory to the Father",
                "Repeat Antiphon 2",
                "Antiphon 3",
                "Old Testament Canticle",
                "Glory to the Father",
                "Repeat Antiphon 3",
                "Short Reading (Scripture)",
                "Responsory",
                "Gospel Canticle Antiphon",
                "Benedictus (Luke 1:68–79)",
                "Glory to the Father",
                "Repeat Gospel Canticle Antiphon",
                "Intercessions",
                "The Lord's Prayer",
                "Concluding Prayer (Collect)",
                "Blessing or Dismissal",
                "Optional: Marian Antiphon"
            ]
        }
        
        return liturgy_structure

    def extract_liturgical_content(self, soup, full_text):
        """
        Extract specific liturgical content from iBreviary
        """
        content = {}
        
        # Clean up the text for better parsing
        clean_text = re.sub(r'\s+', ' ', full_text)
        clean_text = re.sub(r'[\r\n]+', ' ', clean_text)
        
        # Extract Invitatory Antiphon
        invitatory_patterns = [
            r'(?:Come, worship the Lord[^.]*alleluia\.)',
            r'(?:Come, let us worship[^.]*\.)',
            r'(?:The Lord has risen[^.]*alleluia\.)',
            r'(?:Come, let us adore[^.]*\.)'
        ]
        
        for pattern in invitatory_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                content['invitatory_antiphon'] = self._clean_liturgical_text(match.group(0))
                break
        
        # Extract Hymn - look for content after "Hymn" keyword
        hymn_patterns = [
            r'Hymn\s+([^.]+(?:\.[^.]*?){1,4}\.)',
            r'Hymn[:\s]+([A-Z][^.]+(?:\.[^.]*?){1,4}\.)'
        ]
        
        for pattern in hymn_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                hymn_text = self._clean_liturgical_text(match.group(1))
                if len(hymn_text) > 30:  # Make sure it's substantial content
                    content['hymn'] = hymn_text[:400] + "..." if len(hymn_text) > 400 else hymn_text
                    break
        
        # Extract Antiphons with better pattern matching
        antiphon_patterns = [
            (r'Ant\.\s*1[:\s]+([^.]+\.)', 'antiphon_1'),
            (r'Antiphon\s*1[:\s]+([^.]+\.)', 'antiphon_1'),
            (r'Ant\.\s*2[:\s]+([^.]+\.)', 'antiphon_2'),
            (r'Antiphon\s*2[:\s]+([^.]+\.)', 'antiphon_2'),
            (r'Ant\.\s*3[:\s]+([^.]+\.)', 'antiphon_3'),
            (r'Antiphon\s*3[:\s]+([^.]+\.)', 'antiphon_3'),
        ]
        
        for pattern, key in antiphon_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                antiphon_text = self._clean_liturgical_text(match.group(1))
                if len(antiphon_text) > 10:
                    content[key] = antiphon_text
                break
        
        # Extract Psalm content more carefully
        psalm_patterns = [
            r'Psalm\s*\d+[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
            r'(?:Come, let us sing to the Lord[^.]*\.)',
            r'(?:Let us approach him with praise[^.]*\.)',
        ]
        
        psalms = []
        for pattern in psalm_patterns:
            matches = re.finditer(pattern, clean_text, re.IGNORECASE)
            for match in matches:
                psalm_text = self._clean_liturgical_text(match.group(1) if match.groups() else match.group(0))
                if len(psalm_text) > 30 and psalm_text not in psalms:
                    psalms.append(psalm_text)
        
        if psalms:
            content['psalm_1'] = psalms[0][:300] + "..." if len(psalms[0]) > 300 else psalms[0]
            if len(psalms) > 1:
                content['psalm_2'] = psalms[1][:300] + "..." if len(psalms[1]) > 300 else psalms[1]
        
        # Extract Old Testament Canticle
        canticle_patterns = [
            r'Canticle[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,3}\.)',
            r'Old Testament Canticle[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,3}\.)',
        ]
        
        for pattern in canticle_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                canticle_text = self._clean_liturgical_text(match.group(1))
                if len(canticle_text) > 20:
                    content['old_testament_canticle'] = canticle_text[:300] + "..." if len(canticle_text) > 300 else canticle_text
                    break
        
        # Extract Short Reading
        reading_patterns = [
            r'Short Reading[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
            r'Reading[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
        ]
        
        for pattern in reading_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                reading_text = self._clean_liturgical_text(match.group(1))
                if len(reading_text) > 20:
                    content['short_reading'] = reading_text[:250] + "..." if len(reading_text) > 250 else reading_text
                    break
        
        # Extract Responsory
        responsory_patterns = [
            r'Responsory[:\s]*([A-Z][^.]+\.)',
            r'℟[:\s]*([A-Z][^.]+\.)',
            r'Response[:\s]*([A-Z][^.]+\.)',
        ]
        
        for pattern in responsory_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                responsory_text = self._clean_liturgical_text(match.group(1))
                if len(responsory_text) > 10:
                    content['responsory'] = responsory_text
                    break
        
        # Extract Gospel Canticle Antiphon
        gospel_antiphon_patterns = [
            r'Gospel Canticle Antiphon[:\s]*([A-Z][^.]+\.)',
            r'Benedictus Antiphon[:\s]*([A-Z][^.]+\.)',
            r'Canticle of Zechariah[:\s]*([A-Z][^.]+\.)',
        ]
        
        for pattern in gospel_antiphon_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                antiphon_text = self._clean_liturgical_text(match.group(1))
                if len(antiphon_text) > 10:
                    content['gospel_canticle_antiphon'] = antiphon_text
                    break
        
        # Extract Intercessions
        intercessions_pattern = r'Intercessions[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)'
        intercessions_match = re.search(intercessions_pattern, clean_text, re.IGNORECASE)
        if intercessions_match:
            intercessions_text = self._clean_liturgical_text(intercessions_match.group(1))
            if len(intercessions_text) > 20:
                content['intercessions'] = intercessions_text[:300] + "..." if len(intercessions_text) > 300 else intercessions_text
        
        # Extract Concluding Prayer
        prayer_patterns = [
            r'Concluding Prayer[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
            r'Prayer[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
            r'Collect[:\s]*([A-Z][^.]+(?:\.[^.]*?){1,2}\.)',
        ]
        
        for pattern in prayer_patterns:
            match = re.search(pattern, clean_text, re.IGNORECASE)
            if match:
                prayer_text = self._clean_liturgical_text(match.group(1))
                if len(prayer_text) > 20:
                    content['concluding_prayer'] = prayer_text[:250] + "..." if len(prayer_text) > 250 else prayer_text
                    break
        
        return content

    def _clean_liturgical_text(self, text):
        """
        Clean up extracted liturgical text
        """
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text.strip())
        
        # Remove common web artifacts
        text = re.sub(r'(?:More|Breviary|Missal|Prayers)\s*', '', text, flags=re.IGNORECASE)
        text = re.sub(r'\s*\|\s*', ' ', text)
        text = re.sub(r'^\s*[:\-]\s*', '', text)
        
        # Fix punctuation
        text = re.sub(r'\s+([,.!?;:])', r'\1', text)
        text = re.sub(r'([.!?])\s*([A-Z])', r'\1 \2', text)
        
        return text.strip()

    def fetch_morning_prayer_detailed(self):
        """
        Fetch and parse morning prayer with detailed extraction from iBreviary
        """
        url = f"{self.base_url}breviario.php?s=lodi"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            full_text = soup.get_text()
            
            # Get the liturgy structure
            liturgy = self.get_liturgy_of_hours_structure(is_first_hour=True)
            
            # Extract actual liturgical content
            liturgy['extracted_content'] = self.extract_liturgical_content(soup, full_text)
            
            return liturgy
            
        except Exception as e:
            print(f"Error fetching morning prayer: {e}")
            # Return basic structure even if fetching fails
            return self.get_liturgy_of_hours_structure(is_first_hour=True)

    def fetch_daily_readings_detailed(self):
        """
        Fetch daily readings with better parsing
        """
        url = f"{self.base_url}letture.php?s=letture"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            
            full_text = soup.get_text()
            
            readings = {}
            
            # Extract first reading
            first_reading_match = re.search(r'First Reading(.{50,1500}?)(?=Responsorial Psalm|Second Reading|Gospel)', full_text, re.IGNORECASE | re.DOTALL)
            if first_reading_match:
                content = first_reading_match.group(1).strip()
                content = re.sub(r'\s+', ' ', content)
                # Clean up and extract the actual reading text
                reading_match = re.search(r'A reading from[^.]*\.(.+)', content, re.DOTALL)
                if reading_match:
                    readings['first_reading'] = reading_match.group(1).strip()[:1000]
                else:
                    readings['first_reading'] = content[:800]
            
            # Extract Gospel
            gospel_match = re.search(r'Gospel(.{50,1500})$', full_text, re.IGNORECASE | re.DOTALL)
            if gospel_match:
                content = gospel_match.group(1).strip()
                content = re.sub(r'\s+', ' ', content)
                # Clean up
                gospel_reading_match = re.search(r'A reading from[^.]*\.(.+)', content, re.DOTALL)
                if gospel_reading_match:
                    readings['gospel'] = gospel_reading_match.group(1).strip()[:1000]
                else:
                    readings['gospel'] = content[:800]
            
            return readings
            
        except Exception as e:
            print(f"Error fetching readings: {e}")
            return {}

    def create_enhanced_slides(self, prayer_data, readings_data, output_filename=None):
        """
        Create PowerPoint slides with enhanced formatting for elderly congregation
        """
        # Generate filename with date if not provided
        if output_filename is None:
            date_str = datetime.now().strftime("%Y-%m-%d")
            output_filename = f"enhanced_daily_mass_slides_{date_str}.pptx"
        prs = Presentation()
        
        # Set slide dimensions for widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 1. Title slide
        self._add_enhanced_title_slide(prs, prayer_data.get('date', ''))
        
        # 2. Liturgy of the Hours - Invitatory (if first hour)
        extracted_content = prayer_data.get('extracted_content', {}) if prayer_data else {}
        
        if prayer_data and prayer_data.get('is_first_hour') and prayer_data.get('invitatory'):
            self._add_liturgy_sequence_slides(prs, "Invitatory", prayer_data['invitatory']['sequence'], extracted_content)
        
        # 3. Liturgy of the Hours - Morning Prayer
        if prayer_data and prayer_data.get('morning_prayer'):
            self._add_liturgy_sequence_slides(prs, "Morning Prayer", prayer_data['morning_prayer']['sequence'], extracted_content)
        
        # 4. First Reading slide
        if readings_data and readings_data.get('first_reading'):
            self._add_reading_slide(prs, "First Reading", readings_data['first_reading'])
        
        # 5. Gospel slide
        if readings_data and readings_data.get('gospel'):
            self._add_reading_slide(prs, "Gospel", readings_data['gospel'])
        
        # Ensure output directory exists
        output_dir = "output"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        # Create full path for output file
        output_path = os.path.join(output_dir, output_filename)
        
        # Remove existing file if it exists (replacing old file with new for same date)
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
                print(f"Replaced existing file: {output_path}")
            except PermissionError:
                print(f"Warning: Could not replace existing file (may be open in PowerPoint): {output_path}")
                # Create a new filename with timestamp
                timestamp = datetime.now().strftime("%H%M%S")
                base_name, ext = os.path.splitext(output_filename)
                output_filename = f"{base_name}_{timestamp}{ext}"
                output_path = os.path.join(output_dir, output_filename)
                print(f"Creating new file instead: {output_path}")
        
        # Save the presentation
        prs.save(output_path)
        print(f"Enhanced slides saved as: {output_path}")

    def _add_liturgy_sequence_slides(self, prs, section_title, sequence, extracted_content=None):
        """
        Add slides for Liturgy of the Hours sequences (Invitatory and Morning Prayer)
        """
        if not extracted_content:
            extracted_content = {}
            
        for i, item in enumerate(sequence):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Section title at top
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.text = f"{section_title} - Step {i + 1}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            title_para.alignment = PP_ALIGN.CENTER
            
            # Get the actual content for this liturgical element
            content_text = self._get_liturgical_content_text(item, extracted_content)
            
            # Main content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_bottom = Inches(0.2)
            content_frame.margin_top = Inches(0.2)
            content_frame.margin_left = Inches(0.2)
            content_frame.margin_right = Inches(0.2)
            
            # Determine text color and style based on content type
            text_color = self._get_liturgy_text_color(item)
            
            content_frame.text = content_text
            content_para = content_frame.paragraphs[0]
            
            # Adjust font size based on content length
            if len(content_text) <= 150:
                content_para.font.size = Pt(44)
            elif len(content_text) <= 300:
                content_para.font.size = Pt(36)
            elif len(content_text) <= 600:
                content_para.font.size = Pt(30)
            else:
                content_para.font.size = Pt(24)
                
            content_para.font.bold = True
            content_para.font.color.rgb = text_color
            content_para.alignment = PP_ALIGN.CENTER
            content_para.line_spacing = 1.2

    def _get_liturgical_content_text(self, item, extracted_content):
        """
        Get the actual liturgical content text for each item
        """
        item_lower = item.lower()
        
        # Static liturgical texts
        if "opening verse: lord, open my lips" in item_lower:
            return "Lord, open my lips."
        elif "response: and my mouth will proclaim" in item_lower:
            return "And my mouth will proclaim your praise."
        elif "opening verse: god, come to my assistance" in item_lower:
            return "God, come to my assistance."
        elif "response: lord, make haste to help me" in item_lower:
            return "Lord, make haste to help me."
        elif "glory to the father" in item_lower:
            return "Glory to the Father, and to the Son, and to the Holy Spirit. As it was in the beginning, is now, and will be forever. Amen."
        elif "alleluia" in item_lower and "omit during lent" in item_lower:
            return "Alleluia! (Omit during Lent)"
        elif "the lord's prayer" in item_lower:
            return "Our Father, who art in heaven, hallowed be thy name; thy kingdom come, thy will be done on earth as it is in heaven. Give us this day our daily bread, and forgive us our trespasses, as we forgive those who trespass against us; and lead us not into temptation, but deliver us from evil. Amen."
        elif "benedictus" in item_lower:
            return "Blessed be the Lord, the God of Israel; he has come to his people and set them free. He has raised up for us a mighty savior, born of the house of his servant David..."
        
        # Dynamic content from iBreviary
        elif "invitatory antiphon" in item_lower:
            return extracted_content.get('invitatory_antiphon', "Come, let us worship the Lord.")
        elif "invitatory psalm" in item_lower:
            return extracted_content.get('psalm_1', "Come, let us sing to the Lord and shout with joy to the Rock who saves us...")
        elif "hymn" in item_lower:
            return extracted_content.get('hymn', "[Hymn for today - to be sung or recited]")
        elif "antiphon 1" in item_lower and "repeat" not in item_lower:
            return extracted_content.get('antiphon_1', "[Antiphon 1 for today]")
        elif "repeat antiphon 1" in item_lower:
            return extracted_content.get('antiphon_1', "[Repeat Antiphon 1]")
        elif "psalm 1" in item_lower:
            return extracted_content.get('psalm_1', "[Psalm 1 for today]")
        elif "antiphon 2" in item_lower and "repeat" not in item_lower:
            return extracted_content.get('antiphon_2', "[Antiphon 2 for today]")
        elif "repeat antiphon 2" in item_lower:
            return extracted_content.get('antiphon_2', "[Repeat Antiphon 2]")
        elif "psalm 2" in item_lower:
            return extracted_content.get('psalm_2', "[Psalm 2 for today]")
        elif "antiphon 3" in item_lower and "repeat" not in item_lower:
            return extracted_content.get('antiphon_3', "[Antiphon 3 for today]")
        elif "repeat antiphon 3" in item_lower:
            return extracted_content.get('antiphon_3', "[Repeat Antiphon 3]")
        elif "old testament canticle" in item_lower:
            return extracted_content.get('old_testament_canticle', "[Old Testament Canticle for today]")
        elif "short reading" in item_lower:
            return extracted_content.get('short_reading', "[Short Scripture reading for today]")
        elif "responsory" in item_lower:
            return extracted_content.get('responsory', "[Responsory for today]")
        elif "gospel canticle antiphon" in item_lower and "repeat" not in item_lower:
            return extracted_content.get('gospel_canticle_antiphon', "[Gospel Canticle Antiphon for today]")
        elif "repeat gospel canticle antiphon" in item_lower:
            return extracted_content.get('gospel_canticle_antiphon', "[Repeat Gospel Canticle Antiphon]")
        elif "intercessions" in item_lower:
            return extracted_content.get('intercessions', "[Intercessions for today]")
        elif "concluding prayer" in item_lower:
            return extracted_content.get('concluding_prayer', "[Concluding Prayer/Collect for today]")
        elif "blessing" in item_lower or "dismissal" in item_lower:
            return "May almighty God bless you, the Father, and the Son, and the Holy Spirit. Amen. Go in peace."
        elif "marian antiphon" in item_lower:
            return "Hail Holy Queen, Mother of mercy, our life, our sweetness, and our hope..."
        
        # Default fallback
        else:
            return item
    
    def _get_liturgy_text_color(self, item):
        """
        Determine text color based on liturgy item type
        """
        item_lower = item.lower()
        
        # Responses and congregation parts - Blue
        if any(keyword in item_lower for keyword in ['response:', 'glory to the father', 'amen', 'alleluia', 'repeat']):
            return RGBColor(0, 100, 200)  # Blue for congregation responses
        
        # Priest-only parts - Red
        elif any(keyword in item_lower for keyword in ['opening verse:', 'blessing', 'dismissal', 'concluding prayer']):
            return RGBColor(200, 0, 0)  # Red for priest only
        
        # Psalms and readings - Purple
        elif any(keyword in item_lower for keyword in ['psalm', 'canticle', 'reading', 'benedictus', 'magnificat']):
            return RGBColor(128, 0, 128)  # Purple for psalms/readings
        
        # Default - Dark blue
        else:
            return RGBColor(0, 51, 102)  # Dark blue for general liturgical text

    def _add_enhanced_title_slide(self, prs, date_str):
        """Enhanced title slide with large, readable text"""
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Main title with proper text boundaries - positioned higher for better balance, full width
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(2.5))
        title_frame = title_box.text_frame
        title_frame.text = "Daily Catholic Mass"
        title_frame.word_wrap = True
        title_frame.margin_left = Inches(0.3)
        title_frame.margin_right = Inches(0.3)
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(60)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.color.rgb = RGBColor(0, 0, 100)  # Dark blue
        
        # Subtitle with proper boundaries - more space, full width
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.5), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Liturgy of the Hours - Morning Prayer and Readings"
        subtitle_frame.word_wrap = True
        subtitle_frame.margin_left = Inches(0.3)
        subtitle_frame.margin_right = Inches(0.3)
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(40)
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Date with proper boundaries - lower position with more space, full width
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.5), Inches(1.5))
        date_frame = date_box.text_frame
        date_frame.text = date_str
        date_frame.word_wrap = True
        date_frame.margin_left = Inches(0.3)
        date_frame.margin_right = Inches(0.3)
        date_paragraph = date_frame.paragraphs[0]
        date_paragraph.font.size = Pt(32)
        date_paragraph.alignment = PP_ALIGN.CENTER
        date_paragraph.font.italic = True

    def _add_reading_slide(self, prs, reading_title, reading_text):
        """Add reading slides with large, readable text"""
        # Split long text across multiple slides if necessary
        max_chars_per_slide = 1000  # Increased from 800 since we have more space
        text_chunks = self._split_text_for_slides(reading_text, max_chars_per_slide)
        
        for i, chunk in enumerate(text_chunks):
            slide_layout = prs.slide_layouts[6]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Title (add part indicator if multiple slides) - full width
            title_text = reading_title if len(text_chunks) == 1 else f"{reading_title} (Part {i+1}/{len(text_chunks)})"
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.5), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.word_wrap = True
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(48)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.CENTER
            title_paragraph.font.color.rgb = RGBColor(0, 0, 100)
            
            # Reading text with proper boundaries - extended to use full slide width
            content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.8), Inches(12.8), Inches(6.2))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0.4)
            content_frame.margin_right = Inches(0.4)
            content_frame.margin_top = Inches(0.2)
            content_frame.margin_bottom = Inches(0.2)
            
            content_frame.text = chunk
            
            # Set font size based on text length - larger sizes for extended space
            char_count = len(chunk)
            if char_count <= 400:
                font_size = 32
            elif char_count <= 600:
                font_size = 30
            elif char_count <= 800:
                font_size = 28
            elif char_count <= 1000:
                font_size = 26
            else:
                font_size = 24
            
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_after = Pt(12)  # Add space between paragraphs
    
    def _split_text_for_slides(self, text, max_chars_per_slide):
        """Split text into chunks suitable for individual slides"""
        if len(text) <= max_chars_per_slide:
            return [text]
        
        chunks = []
        current_chunk = ""
        sentences = text.split('. ')
        
        for sentence in sentences:
            # If adding this sentence would exceed the limit, start a new chunk
            if len(current_chunk + sentence) > max_chars_per_slide and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = sentence + '. '
            else:
                current_chunk += sentence + '. '
        
        # Add the last chunk if it has content
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        # If we still have chunks that are too long, split them more aggressively
        final_chunks = []
        for chunk in chunks:
            if len(chunk) <= max_chars_per_slide:
                final_chunks.append(chunk)
            else:
                # Split by words if sentences are still too long
                words = chunk.split(' ')
                current_word_chunk = ""
                for word in words:
                    if len(current_word_chunk + word) > max_chars_per_slide and current_word_chunk:
                        final_chunks.append(current_word_chunk.strip())
                        current_word_chunk = word + ' '
                    else:
                        current_word_chunk += word + ' '
                if current_word_chunk.strip():
                    final_chunks.append(current_word_chunk.strip())
        
        return final_chunks if final_chunks else [text[:max_chars_per_slide]]

def main(include_invitatory=True, preview_mode=False):
    print("Enhanced Catholic Slide Generator")
    print("=" * 40)
    
    generator = EnhancedCatholicSlideGenerator()
    
    print("Fetching morning prayer content...")
    prayer_data = generator.fetch_morning_prayer_detailed()
    
    # Override the is_first_hour setting if specified
    if prayer_data and not include_invitatory:
        prayer_data['is_first_hour'] = False
    
    if prayer_data:
        print(f"✓ Found Liturgy structure for: {prayer_data['date']}")
        if prayer_data.get('is_first_hour'):
            print("✓ Including Invitatory (first hour of day)")
        print("✓ Including Morning Prayer sequence")
        
        # Show extracted content if available
        if prayer_data.get('extracted_content'):
            extracted = prayer_data['extracted_content']
            print(f"✓ Extracted {len(extracted)} liturgical elements:")
            
            content_items = [
                ('invitatory_antiphon', 'Invitatory Antiphon'),
                ('hymn', 'Hymn'),
                ('antiphon_1', 'Antiphon 1'),
                ('antiphon_2', 'Antiphon 2'), 
                ('antiphon_3', 'Antiphon 3'),
                ('psalm_1', 'Psalm 1'),
                ('psalm_2', 'Psalm 2'),
                ('old_testament_canticle', 'Old Testament Canticle'),
                ('short_reading', 'Short Reading'),
                ('responsory', 'Responsory'),
                ('gospel_canticle_antiphon', 'Gospel Canticle Antiphon'),
                ('intercessions', 'Intercessions'),
                ('concluding_prayer', 'Concluding Prayer')
            ]
            
            for key, name in content_items:
                if extracted.get(key):
                    preview = extracted[key][:50] + "..." if len(extracted[key]) > 50 else extracted[key]
                    print(f"  • {name}: {preview}")
        
        # Preview mode - show what will be on slides
        if preview_mode:
            print("\n" + "=" * 50)
            print("SLIDE CONTENT PREVIEW")
            print("=" * 50)
            
            if prayer_data.get('is_first_hour') and prayer_data.get('invitatory'):
                print("\n--- INVITATORY SLIDES ---")
                for i, item in enumerate(prayer_data['invitatory']['sequence']):
                    content_text = generator._get_liturgical_content_text(item, extracted)
                    print(f"\nSlide {i+1}: {item}")
                    print(f"Content: {content_text[:100]}...")
            
            if prayer_data.get('morning_prayer'):
                print("\n--- MORNING PRAYER SLIDES ---")
                for i, item in enumerate(prayer_data['morning_prayer']['sequence']):
                    content_text = generator._get_liturgical_content_text(item, extracted)
                    print(f"\nSlide {i+1}: {item}")
                    print(f"Content: {content_text[:100]}...")
            return
    
    print("\nFetching daily readings...")
    readings_data = generator.fetch_daily_readings_detailed()
    
    if readings_data:
        if 'first_reading' in readings_data:
            print(f"✓ Found First Reading: {readings_data['first_reading'][:50]}...")
        if 'gospel' in readings_data:
            print(f"✓ Found Gospel: {readings_data['gospel'][:50]}...")
    
    print("\nCreating enhanced slides...")
    generator.create_enhanced_slides(prayer_data, readings_data)
    print("✓ Enhanced slides created successfully!")

if __name__ == "__main__":
    # Check for command line arguments
    include_invitatory = True
    preview_mode = False
    
    for arg in sys.argv[1:]:
        if arg.lower() in ['--no-invitatory', '-n']:
            include_invitatory = False
            print("Note: Running without Invitatory (not first hour)")
        elif arg.lower() in ['--preview', '-p']:
            preview_mode = True
            print("Note: Running in preview mode (showing slide content without generating PowerPoint)")
    
    main(include_invitatory, preview_mode)