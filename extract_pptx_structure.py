"""
Script to extract structure and content from PowerPoint file for analysis
"""
from pptx import Presentation
import os

def extract_pptx_structure(pptx_path):
    """Extract all slide content and structure from PowerPoint file"""
    
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
    
    try:
        prs = Presentation(pptx_path)
        
        print("=" * 60)
        print(f"POWERPOINT STRUCTURE ANALYSIS")
        print(f"File: {pptx_path}")
        print(f"Total Slides: {len(prs.slides)}")
        print("=" * 60)
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n--- SLIDE {i} ---")
            
            # Extract all text from the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_text.append(text)
                elif hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            # Print slide content
            if slide_text:
                for j, text in enumerate(slide_text):
                    print(f"Text {j+1}: {text}")
            else:
                print("(No text content)")
            
            print("-" * 40)
        
        return True
        
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        return False

if __name__ == "__main__":
    # Path to the PowerPoint file
    pptx_file = "output/11_10 Morning Readings & Prayers .pptx"
    extract_pptx_structure(pptx_file)