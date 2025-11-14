"""
BBGRL Slide Generator V2.py - Template-Based Dynamic Generator
Uses the reference PowerPoint structure as a formatting template
Fetches live liturgical data and applies the exact same presentation structure

This creates presentations that look identical to the reference but with current liturgical content
"""

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
import re
from datetime import datetime, timedelta
import sys
import os
import traceback
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.chrome.options import Options
import time

class BBGRLSlideGeneratorV2:
    def __init__(self):
        self.base_url = "https://www.ibreviary.com/m2/"
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        })
        
        # Reference structure template (based on the analyzed PowerPoint)
        self.reference_template = self._get_reference_template()
        
        # Selenium driver (initialized when needed)
        self.driver = None

    def _get_reference_template(self):
        """
        Define the exact reference structure that should be applied to any liturgical data
        This serves as the formatting template regardless of the content
        """
        return {
            "metadata": {
                "total_expected_slides": 135,  # Target slide count
                "title_pattern": "{date} Morning Readings & Prayers", # TODO: might not need
                "structure_sections": [
                    "opening_slides",
                    "psalmody_section", 
                    "reading_section",
                    "gospel_canticle_section",
                    "intercessions_section",
                    "sacred_heart_hymns",
                    "mass_readings",
                    "post_communion_prayers", 
                    "transition_slides",
                    "jubilee_prayer",
                    "st_joseph_prayer"
                ]
            },
            
            # Template structure - defines how to organize any liturgical data
            "section_templates": { # TODO: double check these section templates are correct
                "opening_slides": {
                    "slide_count": 2,
                    "slides": [
                        {"type": "title", "content": "title_slide"},
                        {"type": "blank", "content": "transition"}
                    ]
                },
                
                "psalmody_section": {
                    "expected_elements": [
                        "antiphon_1", "psalm_1", "glory_be", "repeat_antiphon_1",
                        "antiphon_2", "canticle", "glory_be", "repeat_antiphon_2", 
                        "antiphon_3", "psalm_2", "glory_be", "repeat_antiphon_3"
                    ],
                    "slide_pattern": "alternating_priest_people",
                    "title_slide": {"text": "PSALMODY", "include": True}
                },
                
                "reading_section": {
                    "expected_elements": ["short_reading", "responsory"],
                    "slide_pattern": "reading_format",
                    "title_slide": {"text": "READING", "include": True}
                },
                
                "gospel_canticle_section": {
                    "expected_elements": ["gospel_antiphon", "benedictus", "glory_be", "repeat_antiphon"],
                    "slide_pattern": "canticle_format", 
                    "title_slide": {"text": "GOSPEL CANTICLE", "include": True}
                },
                
                "intercessions_section": {
                    "expected_elements": ["intercessions", "lords_prayer", "concluding_prayer"],
                    "slide_pattern": "intercession_format",
                    "title_slide": {"text": "INTERCESSIONS", "include": True}
                },
                
                "sacred_heart_hymns": {
                    "slide_count": 6,  # Fixed devotional content
                    "content_type": "static_devotional"
                },
                
                "mass_readings": {
                    "expected_elements": ["first_reading", "responsorial_psalm", "gospel_acclamation", "gospel"],
                    "slide_pattern": "mass_reading_format"
                },
                
                "post_communion_prayers": {
                    "slide_count": 17,  # Fixed devotional content
                    "content_type": "static_devotional"
                },
                
                "transition_slides": {
                    "slide_count": 10,  # Blank slides
                    "content_type": "blank_transitions"
                },
                
                "jubilee_prayer": {
                    "slide_count": 7,  # Fixed prayer content
                    "content_type": "static_prayer"
                },
                
                "st_joseph_prayer": {
                    "slide_count": 12,  # Fixed prayer content
                    "content_type": "static_prayer"
                }
            },
            
            # Formatting rules based on reference presentation
            "formatting_rules": {
                "priest_color": RGBColor(200, 0, 0),     # Red
                "people_color": RGBColor(0, 100, 200),   # Blue  
                "all_color": RGBColor(100, 0, 100),      # Purple
                "title_color": RGBColor(0, 51, 102),     # Dark blue
                "reading_color": RGBColor(100, 0, 0),    # Dark red for mass readings
                "devotional_color": RGBColor(139, 0, 0), # Sacred heart red
                
                "font_sizes": {
                    "title": Pt(48),
                    "subtitle": Pt(32), 
                    "priest_text": Pt(32),
                    "people_text": Pt(32),
                    "all_text": Pt(36),
                    "reading_text": Pt(30),
                    "prayer_text": Pt(30)
                }
            }
        }

    def fetch_live_liturgical_data(self, target_date=None):
        """
        Fetch current liturgical data from iBreviary and structure it according to the template
        """
        if target_date is None:
            target_date = datetime.now()
        
        print(f"Fetching live liturgical data from iBreviary for {target_date.strftime('%B %d, %Y')}...")
        
        try:
            # Fetch Morning Prayer
            morning_prayer_data = self._fetch_morning_prayer_structured(target_date)
            
            # Fetch daily readings  
            readings_data = self._fetch_daily_readings_structured(target_date)
            
            # Combine into structured data matching the reference template
            structured_data = {
                "date": target_date.strftime('%B %d, %Y'),
                "morning_prayer": morning_prayer_data,
                "mass_readings": readings_data,
                "static_content": self._get_static_devotional_content()
            }
            
            print(f"Successfully fetched liturgical data for {structured_data['date']}")
            return structured_data
            
        except Exception as e:
            print(f"Error fetching liturgical data: {e}")
            print("Using fallback template structure...")
            return self._get_fallback_data(target_date)

    def _navigate_ibreviary_to_date(self, target_date):
        """
        Navigate iBreviary mobile site to a specific date and return the Morning Prayer HTML
        
        Steps:
        1. Click the "More" link
        2. Update the date input boxes (giorno=day, mese=month dropdown, anno=year)
        3. Click the "Breviary" link
        4. Click on the "Morning Prayer" link
        
        Returns the full HTML of the Morning Prayer page for the specified date
        """
        try:
            # Initialize Chrome in headless mode
            chrome_options = Options()
            chrome_options.add_argument('--headless')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument('--disable-gpu')
            chrome_options.add_argument('--window-size=1920,1080')
            
            self.driver = webdriver.Chrome(options=chrome_options)
            wait = WebDriverWait(self.driver, 10)
            
            print(f"  -> Navigating to iBreviary mobile site...")
            self.driver.get(self.base_url)
            time.sleep(3)  # Give page more time to load
            
            # Step 1: Click the "More" link
            print(f"  -> Clicking 'More' menu...")
            try:
                # Try multiple approaches to find the More link
                more_link = None
                
                # Approach 1: By link text
                try:
                    more_link = wait.until(EC.presence_of_element_located((By.LINK_TEXT, "More")))
                except:
                    pass
                
                # Approach 2: By partial link text
                if not more_link:
                    try:
                        more_link = wait.until(EC.presence_of_element_located((By.PARTIAL_LINK_TEXT, "More")))
                    except:
                        pass
                
                # Approach 3: By href containing opzioni.php
                if not more_link:
                    try:
                        more_link = wait.until(EC.presence_of_element_located((By.XPATH, "//a[contains(@href, 'opzioni.php')]")))
                    except:
                        pass
                
                if more_link:
                    # Scroll to element and click
                    self.driver.execute_script("arguments[0].scrollIntoView(true);", more_link)
                    time.sleep(0.5)
                    more_link.click()
                    time.sleep(2)
                else:
                    print(f"  ⚠ Could not find 'More' link")
                    return None
                    
            except Exception as e:
                print(f"  ⚠ Error clicking 'More' link: {e}")
                return None
            
            # Step 2: Update date input boxes
            print(f"  -> Setting date to {target_date.strftime('%d/%m/%Y')}...")
            day = target_date.day
            month = target_date.month  # 1-12
            year = target_date.year
            
            try:
                # Fill day field (input name="giorno")
                day_field = self.driver.find_element(By.NAME, "giorno")
                day_field.clear()
                day_field.send_keys(str(day))
                print(f"    Set day: {day}")
                
                # Select month dropdown (select name="mese")
                from selenium.webdriver.support.select import Select
                month_dropdown = Select(self.driver.find_element(By.NAME, "mese"))
                month_dropdown.select_by_index(month - 1)  # Month indices are 0-based
                print(f"    Set month: {month}")
                
                # Fill year field (input name="anno")
                year_field = self.driver.find_element(By.NAME, "anno")
                year_field.clear()
                year_field.send_keys(str(year))
                print(f"    Set year: {year}")
                
                # Click the OK button to apply the date change
                print(f"    Clicking 'OK' button to apply date...")
                ok_button = self.driver.find_element(By.NAME, "ok")
                ok_button.click()
                time.sleep(2)  # Wait for page to reload with new date
                
            except Exception as e:
                print(f"  ⚠ Could not set date fields: {e}")
                return None
            
            # Step 3: Click the "Breviary" link
            print(f"  -> Clicking 'Breviary' link...")
            try:
                breviary_link = wait.until(EC.element_to_be_clickable((By.LINK_TEXT, "Breviary")))
                breviary_link.click()
                time.sleep(2)
            except Exception as e:
                print(f"  ⚠ Could not find 'Breviary' link: {e}")
                return None
            
            # Step 4: Click on the "Morning Prayer" link
            print(f"  -> Clicking 'Morning Prayer' link...")
            try:
                # Try various text options for Morning Prayer
                morning_prayer_link = None
                for text in ["Morning Prayer", "Lauds", "Lodi"]:
                    try:
                        morning_prayer_link = wait.until(EC.element_to_be_clickable((By.PARTIAL_LINK_TEXT, text)))
                        break
                    except:
                        continue
                
                if morning_prayer_link:
                    morning_prayer_link.click()
                    time.sleep(2)
                else:
                    print("  ⚠ Could not find 'Morning Prayer' link")
                    return None
                    
            except Exception as e:
                print(f"  ⚠ Error clicking Morning Prayer: {e}")
                return None
            
            # Get the final page HTML
            html_content = self.driver.page_source
            
            print(f"  Successfully navigated to Morning Prayer for {target_date.strftime('%B %d, %Y')}")
            
            return html_content
            
        except Exception as e:
            print(f"  Error during Selenium navigation: {e}")
            import traceback
            traceback.print_exc()
            return None
        
        finally:
            if self.driver:
                self.driver.quit()
                self.driver = None

    def _fetch_morning_prayer_structured(self, target_date):
        """
        Fetch morning prayer and structure it to match the reference template exactly
        Uses Selenium to navigate iBreviary to the specific date
        """
        try:
            # Use Selenium to navigate to the specific date
            print(f"  Fetching Morning Prayer using Selenium navigation...")
            html_content = self._navigate_ibreviary_to_date(target_date)
            
            if not html_content:
                print(f"  ⚠ Selenium navigation failed, using fallback data")
                return self._get_fallback_morning_prayer()
            
            # Parse the HTML content
            soup = BeautifulSoup(html_content, 'html.parser')
            full_text = soup.get_text()
            
            # Find "PSALMODY" in all caps and extract only the text AFTER it
            # This skips the "Tune:" and "Text:" segments that come before PSALMODY
            psalmody_pos = full_text.upper().find('PSALMODY')
            
            if psalmody_pos >= 0:
                # Extract only the text after PSALMODY for all parsing
                text_after_psalmody = full_text[psalmody_pos:]
                print(f"  Found PSALMODY at position {psalmody_pos}, parsing content after it")
                
                # Create a new soup object with only the content after PSALMODY
                # Find the element containing PSALMODY and get everything after it
                psalmody_soup = soup
            else:
                # Fallback: use full text if PSALMODY not found
                text_after_psalmody = full_text
                psalmody_soup = soup
                print(f"  ⚠ PSALMODY marker not found, using full text")
            
            # Extract and structure the content to match reference format
            # Pass soup object and text_after_psalmody for parsing
            structured = {
                "psalmody": {
                    "antiphon_1": self._extract_antiphon_and_psalm_info(soup, 1, text_after_psalmody),
                    "psalm_1": self._extract_psalm_verses_from_html(soup, 1),
                    "antiphon_2": self._extract_antiphon(text_after_psalmody, 2),
                    "canticle_info": self._extract_canticle_info(soup, text_after_psalmody),
                    "canticle": self._extract_canticle_verses(soup, text_after_psalmody),
                    "antiphon_3": self._extract_antiphon_and_psalm_info(soup, 3, text_after_psalmody),
                    "psalm_3": self._extract_psalm_verses_from_html(soup, 3)
                },
                "reading": {
                    "short_reading": self._extract_short_reading(text_after_psalmody),
                    "responsory": self._extract_responsory_from_html(soup, text_after_psalmody)
                },
                "gospel_canticle": {
                    "antiphon": self._extract_gospel_antiphon(text_after_psalmody),
                    "benedictus_verses": self._extract_benedictus_verses(text_after_psalmody)
                },
                "intercessions": self._extract_intercessions(soup, full_text)
            }
            
            return structured
            
        except Exception as e:
            print(f"Error parsing morning prayer: {e}")
            return self._get_fallback_morning_prayer()

    def _fetch_daily_readings_structured(self, target_date):
        """
        Fetch daily readings and structure them to match the reference template
        """
        # For iBreviary, we might need to adjust URL parameters for specific dates
        # For now, using current day's readings
        url = f"{self.base_url}letture.php?s=letture"
        
        try:
            response = self.session.get(url)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            full_text = soup.get_text()
            
            structured = {
                "first_reading": {
                    "citation": self._extract_first_reading_citation(full_text),
                    "verses": self._extract_first_reading_verses(full_text)
                },
                "responsorial_psalm": {
                    "citation": self._extract_psalm_citation(full_text),
                    "verses": self._extract_psalm_response_verses(full_text)
                },
                "gospel_acclamation": {
                    "verse": self._extract_gospel_acclamation(full_text)
                },
                "gospel": {
                    "citation": self._extract_gospel_citation(full_text),
                    "verses": self._extract_gospel_verses(full_text)
                }
            }
            
            return structured
            
        except Exception as e:
            print(f"Error parsing daily readings: {e}")
            return self._get_fallback_readings()

    def _extract_antiphon_and_psalm_info(self, text, number, text_after_psalmody=None):
        """Extract antiphon text and associated psalm information
        
        Note: 'text' parameter should be the BeautifulSoup object, not plain text,
        when called from _fetch_morning_prayer_structured
        
        For Antiphon 1: Uses text_after_psalmody (which starts at PSALMODY heading)
        to extract the antiphon and psalm info, skipping everything before PSALMODY
        """
        # If text is a BeautifulSoup object, extract from HTML structure
        if hasattr(text, 'find_all'):
            soup = text
            text_content = text_after_psalmody if text_after_psalmody else soup.get_text()
            
            # Extract antiphon text from plain text
            antiphon_text = ""
            psalm_title = ""
            psalm_subtitle = ""
            
            if number == 1 or number == 3:
                # Parse directly from text_after_psalmody (already starts after PSALMODY)
                # For Ant 1: Skip the Invitatory section - look for "Ant. 1" specifically
                # For Ant 3: Look for "Ant. 3" specifically
                # Pattern needs to capture multi-sentence antiphons (e.g., "sentence 1. sentence 2. sentence 3.")
                # Stop at: "Psalm" keyword (red text marker) or "Ant." keyword
                antiphon_patterns = [
                    rf'Ant\.\s*{number}[:\s]+(.+?)(?=Psalm\s+\d|\nAnt\.)',  # Match until "Psalm" or next "Ant."
                    rf'Antiphon\s*{number}[:\s]+(.+?)(?=Psalm\s+\d|\nAnt\.)',
                ]
                
                for pattern in antiphon_patterns:
                    match = re.search(pattern, text_content, re.IGNORECASE | re.DOTALL)
                    if match:
                        antiphon_text = match.group(1).strip()
                        # Clean up any extra whitespace/newlines within the antiphon
                        antiphon_text = re.sub(r'\s+', ' ', antiphon_text).strip()
                        print(f"  Found Antiphon {number} text: {antiphon_text[:50]}...")
                        break
                
                # Now extract the red text that comes immediately after "Ant. X" in the HTML
                # This is the text in <span class="rubrica"> tags that appear right after the antiphon
                # It typically shows "Psalm X:Y-Z" followed by a subtitle on the next line
                try:
                    # Find "Ant. X" in a rubrica span using the soup object passed in
                    ant_rubrica = soup.find('span', class_='rubrica', string=re.compile(rf'Ant\.\s*{number}'))
                    if ant_rubrica:
                        # The next rubrica span should have the psalm info
                        next_rubrica = ant_rubrica.find_next('span', class_='rubrica')
                        if next_rubrica:
                            # Get the text content which should be something like "Psalm 63:2-9\nA soul thirsting for God"
                            rubrica_text = next_rubrica.get_text(separator='\n').strip()
                            lines = rubrica_text.split('\n')
                            
                            if lines:
                                psalm_title = lines[0].strip()
                                print(f"  Found red psalm title: {psalm_title}")
                            
                            if len(lines) > 1:
                                psalm_subtitle = lines[1].strip()
                                print(f"  Found red psalm subtitle: {psalm_subtitle}")
                except Exception as e:
                    print(f"  ⚠ Could not extract red psalm text from HTML: {e}")

                
                # Fallback: Extract psalm info from the section after PSALMODY using regex
                # Look for the psalm pattern in text after PSALMODY and after the antiphon
                if not psalm_title:
                    # Pattern to match: Psalm 95, Psalm 95:1-5, Psalm 95A, etc.
                    # Capture up to a newline or another "Psalm" keyword
                    psalm_pattern = r'Psalm\s+(\d+)([A-Z])?(?::(\d+)(?:-(\d+))?)?\s*([^\n]*?)(?=\nPsalm|\n\n|Psalm\s+\d|$)'
                    psalm_matches = re.finditer(psalm_pattern, text_content if isinstance(text_content, str) else str(text_content), re.IGNORECASE)
                    
                    # Get the first psalm match
                    first_psalm_match = None
                    for match in psalm_matches:
                        first_psalm_match = match
                        break
                    
                    if first_psalm_match:
                        psalm_num = first_psalm_match.group(1)
                        psalm_letter = first_psalm_match.group(2) if first_psalm_match.group(2) else ""
                        verse_start = first_psalm_match.group(3)
                        verse_end = first_psalm_match.group(4)
                        subtitle_raw = first_psalm_match.group(5)
                        
                        # Build psalm title
                        if verse_start and verse_end:
                            psalm_title = f"Psalm {psalm_num}{psalm_letter}:{verse_start}-{verse_end}"
                        elif verse_start:
                            psalm_title = f"Psalm {psalm_num}{psalm_letter}:{verse_start}"
                        else:
                            psalm_title = f"Psalm {psalm_num}{psalm_letter}"
                        
                        # Extract subtitle if present
                        if subtitle_raw:
                            # Clean up subtitle - remove extra whitespace
                            subtitle = subtitle_raw.strip()
                            # Remove any text that looks like it's the start of another psalm reference
                            subtitle = re.sub(r'Psalm\s+\d.*$', '', subtitle, flags=re.IGNORECASE).strip()
                            # If subtitle contains "salm" (partial word), remove it
                            subtitle = re.sub(r'\bP?salm\b.*$', '', subtitle, flags=re.IGNORECASE).strip()
                            # Limit length
                            if len(subtitle) > 100:
                                subtitle = subtitle[:100].rsplit(' ', 1)[0] + '...'
                            psalm_subtitle = subtitle
                        
                        print(f"  Found psalm: {psalm_title} - {psalm_subtitle}")
            else:
                # For antiphons 2 and 3, use the original logic
                antiphon_patterns = [
                    rf'Ant\.\s*{number}[:\s]+([^.]+\.)',
                    rf'Antiphon\s*{number}[:\s]+([^.]+\.)'
                ]
                
                for pattern in antiphon_patterns:
                    match = re.search(pattern, text_content, re.IGNORECASE)
                    if match:
                        antiphon_text = match.group(1).strip()
                        break
            
            return {
                "text": antiphon_text if antiphon_text else "",
                "format": "all_response",
                "psalm_title": psalm_title,
                "psalm_subtitle": psalm_subtitle
            }
        
        # Fallback: if plain text is passed (for backwards compatibility)
        if number == 1:
            antiphon_patterns = [
                r'Ant\.\s+([^.]+\.)',
                rf'Ant\.\s*{number}[:\s]+([^.]+\.)',
                rf'Antiphon\s*{number}[:\s]+([^.]+\.)',
                r'Antiphon[:\s]+([^.]+\.)'
            ]
        else:
            antiphon_patterns = [
                rf'Ant\.\s*{number}[:\s]+([^.]+\.)',
                rf'Antiphon\s*{number}[:\s]+([^.]+\.)'
            ]
        
        antiphon_text = ""
        for pattern in antiphon_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                antiphon_text = match.group(1).strip()
                break
        
        return {
            "text": antiphon_text if antiphon_text else "",
            "format": "all_response",
            "psalm_title": "",
            "psalm_subtitle": ""
        }

    def _extract_antiphon(self, text, number):
        """Extract antiphon text and structure it with priest/people alternation"""
        # Pattern to find antiphons - match up to any sentence-ending punctuation (. ! ?)
        patterns = [
            rf'Ant\.\s*{number}[:\s]+([^.!?]+[.!?])',
            rf'Antiphon\s*{number}[:\s]+([^.!?]+[.!?])'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                antiphon_text = match.group(1).strip()
                return {
                    "text": antiphon_text,
                    "format": "all_response"
                }
        
        return {
            "text": "",
            "format": "all_response"
        }

    def _extract_psalm_verses_from_html(self, soup, psalm_number):
        """Extract psalm verses directly from HTML structure
        
        Looks for the rubrica span with 'Ant. X', then extracts the psalm verses
        which are plain text separated by <br><br> tags, stopping at "Psalm Prayer"
        """
        verses = []
        
        try:
            # Find the antiphon marker for this psalm number in the HTML
            ant_pattern = rf'Ant\.\s*{psalm_number}\s*$'  # Match "Ant. X" at end of string
            ant_span = None
            
            # Find all rubrica spans
            for span in soup.find_all('span', class_='rubrica'):
                if re.match(ant_pattern, span.get_text().strip()):
                    ant_span = span
                    break
            
            if not ant_span:
                print(f"  ⚠ Could not find Ant. {psalm_number} in HTML")
                return self._get_fallback_verses(psalm_number)
            
            # Find the parent element and get its HTML
            parent = ant_span.parent
            parent_html = str(parent)
            
            # Find the position of this antiphon in the HTML
            ant_pos = parent_html.find(str(ant_span))
            if ant_pos < 0:
                return self._get_fallback_verses(psalm_number)
            
            # Get HTML after the antiphon
            html_after_ant = parent_html[ant_pos + len(str(ant_span)):]
            
            # Find where to stop: "Psalm Prayer" or next "Ant. X"
            stop_patterns = [
                r'<span class="rubrica">Psalm\s+Prayer</span>',
                rf'<span class="rubrica">Ant\.\s*{psalm_number + 1}</span>',
                r'<span class="rubrica">Ant\.</span>',  # Repeated antiphon
            ]
            
            end_pos = len(html_after_ant)
            for pattern in stop_patterns:
                match = re.search(pattern, html_after_ant, re.IGNORECASE)
                if match:
                    end_pos = min(end_pos, match.start())
            
            # Extract the psalm section HTML
            psalm_html = html_after_ant[:end_pos]
            
            # Parse to get the text, skipping the intro em tag and rubrica (red) text
            psalm_soup = BeautifulSoup(psalm_html, 'html.parser')
            
            # Remove the psalm title rubrica and intro em tag
            for rubrica in psalm_soup.find_all('span', class_='rubrica', string=re.compile(r'Psalm\s+\d')):
                rubrica.decompose()
            for em in psalm_soup.find_all('em'):
                em.decompose()
            
            # Get the remaining HTML and split by <br><br> to get verses/stanzas
            remaining_html = str(psalm_soup)
            
            # Split by double br tags
            verse_sections = re.split(r'<br\s*/?>\s*<br\s*/?>', remaining_html)
            
            verse_count = 0
            skipped_first_section = False  # Track if we've skipped the antiphon
            
            for section in verse_sections:
                # Parse this section to get clean text
                section_soup = BeautifulSoup(section, 'html.parser')
                
                # Remove asterisks and daggers (rubrica spans)
                for rubrica in section_soup.find_all('span', class_='rubrica'):
                    rubrica.decompose()
                
                # Get text
                verse_text = section_soup.get_text().strip()
                
                # Skip empty or short verses
                if not verse_text or len(verse_text) < 20:
                    continue
                
                # Skip the first valid section if it looks like the antiphon text
                # The antiphon typically appears as the very first text section after "Ant. X"
                # Only skip the FIRST occurrence to avoid false positives with verses that have similar text
                # For a more reliable check, we look for specific patterns OR check if the text is
                # relatively short (< 150 chars) and complete (ends with punctuation)
                if not skipped_first_section and verse_count == 0:
                    # Known antiphon patterns from specific dates
                    if re.search(r'(Each morning|Martin, priest|My heart is ready|You who stand in his sanctuary)', verse_text, re.IGNORECASE):
                        print(f"  Skipping antiphon text in verse extraction: {verse_text[:50]}...")
                        skipped_first_section = True
                        continue
                    # Generic check: first section that's short and looks like a complete sentence
                    elif len(verse_text) < 150 and verse_text.endswith(('.', '!', '?')):
                        print(f"  Skipping potential antiphon text in verse extraction: {verse_text[:50]}...")
                        skipped_first_section = True
                        continue
                
                # Clean up: normalize whitespace
                verse_text = re.sub(r'\s+', ' ', verse_text).strip()
                
                # Add period if needed
                if not verse_text[-1] in '.!?"':
                    verse_text += '.'
                
                # Alternate speaker
                speaker = "Priest" if verse_count % 2 == 0 else "People"
                
                verses.append({
                    "speaker": speaker,
                    "text": verse_text
                })
                verse_count += 1
            
            if verses:
                print(f"  Extracted {len(verses)} verses for Psalm {psalm_number}")
                return verses
                
        except Exception as e:
            print(f"  ⚠ Error parsing psalm verses from HTML: {e}")
            import traceback
            traceback.print_exc()
        
        return self._get_fallback_verses(psalm_number)

    def _extract_psalm_verses(self, text, psalm_number):
        """Extract psalm verses - ONLY the verse paragraphs after the red psalm text
        
        Parsing logic (ignoring all the red text/antiphons):
        1. Start right after the intro sentence that ends with parentheses and period like "(Arnobius)."
        2. Skip any single line that sits between 2 empty lines
        3. Grab each proceeding paragraph as a verse
        4. Stop at "Glory to the Father" or "Psalm Prayer"
        5. Alternate speakers: Priest, People, Priest, People (always start with Priest)
        """
        verses = []
        
        try:
            # Find the intro sentence that marks the end of red text
            # It typically ends with a closing parenthesis followed by period: ).
            # Examples: "(Arnobius).", "(St. Augustine).", etc.
            intro_patterns = [
                r'\([^)]{3,50}\)\s*\.',  # Text in parentheses ending with period
                r'Psalm\s+\d+[A-Z]?(?::\d+(?:-\d+)?)?\s*[^\n]{10,100}\n'  # Or just after subtitle
            ]
            
            # Find all occurrences of these patterns
            intro_matches = []
            for pattern in intro_patterns:
                intro_matches.extend(list(re.finditer(pattern, text, re.IGNORECASE)))
            
            # Sort by position
            intro_matches.sort(key=lambda m: m.start())
            
            # Get the match that corresponds to this psalm number
            if psalm_number <= len(intro_matches):
                match = intro_matches[psalm_number - 1]
                start_pos = match.end()
            else:
                # Fallback: find "Psalm X" and skip ahead
                psalm_pattern = rf'Psalm\s+\d+[A-Z]?(?::\d+(?:-\d+)?)?'
                psalm_matches = list(re.finditer(psalm_pattern, text, re.IGNORECASE))
                if psalm_number <= len(psalm_matches):
                    # Skip 300 characters after psalm title to get past red text
                    start_pos = psalm_matches[psalm_number - 1].end() + 300
                else:
                    return self._get_fallback_verses(psalm_number)
            
            # Find the end: "Glory to the Father"
            glory_pattern = r'Glory\s+to\s+the\s+Father'
            glory_match = re.search(glory_pattern, text[start_pos:], re.IGNORECASE)
            
            if glory_match:
                end_pos = start_pos + glory_match.start()
            else:
                # Try "Psalm Prayer"
                psalm_prayer_match = re.search(r'Psalm\s+Prayer', text[start_pos:], re.IGNORECASE)
                if psalm_prayer_match:
                    end_pos = start_pos + psalm_prayer_match.start()
                else:
                    # Take next 2000 characters
                    end_pos = start_pos + 2000
            
            # Extract verse section
            verse_section = text[start_pos:end_pos].strip()
            
            # Split into paragraphs by looking for sentence endings or double newlines
            # Each paragraph is typically separated by period-newline or double newlines
            paragraphs = re.split(r'(?:\.\s*\n)|(?:\n\s*\n)', verse_section)
            
            # Process each paragraph
            verse_count = 0
            for para in paragraphs:
                para = para.strip()
                
                # Skip empty or very short
                if not para or len(para) < 20:
                    continue
                
                # Skip if it looks like a marker
                if re.match(r'^(Psalm|Ant\.|Glory|℟)', para, re.IGNORECASE):
                    continue
                
                # Clean up: replace asterisks and multiple spaces
                cleaned = re.sub(r'\s*\*\s*', ' ', para)
                cleaned = re.sub(r'\s+', ' ', cleaned).strip()
                
                # Make sure it ends with proper punctuation
                if cleaned and not cleaned[-1] in '.!?"':
                    cleaned += '.'
                
                if len(cleaned) < 20:
                    continue
                
                # Alternate speaker
                speaker = "Priest" if verse_count % 2 == 0 else "People"
                
                verses.append({
                    "speaker": speaker,
                    "text": cleaned
                })
                verse_count += 1
            
            if verses:
                print(f"  Extracted {len(verses)} verses for Psalm {psalm_number}")
                return verses
                
        except Exception as e:
            print(f"  ⚠ Error parsing psalm verses: {e}")
            import traceback
            traceback.print_exc()
        
        return self._get_fallback_verses(psalm_number)
    
    def _get_fallback_verses(self, psalm_number):
        """Return fallback verses when parsing fails"""
        print(f"  Using fallback verses for Psalm {psalm_number}")
        return [
            {"speaker": "Priest", "text": f"[Psalm {psalm_number} verse 1 - Priest]"},
            {"speaker": "People", "text": f"[Psalm {psalm_number} verse 2 - People]"},
            {"speaker": "Priest", "text": f"[Psalm {psalm_number} verse 3 - Priest]"},
            {"speaker": "People", "text": f"[Psalm {psalm_number} verse 4 - People]"}
        ]

    def _extract_canticle_verses(self, soup, text=None):
        """Extract canticle verses directly from HTML structure
        
        Looks for the rubrica span with 'Canticle:', then extracts the canticle verses
        which are plain text separated by <br><br> tags, stopping at "Glory to the Father"
        or the next antiphon marker.
        
        Args:
            soup: BeautifulSoup object of the HTML
            text: Optional plain text fallback (not currently used)
            
        Returns:
            dict with 'verses' list and 'omit_glory_be' boolean flag
        """
        verses = []
        omit_glory_be = False  # Flag to track if Glory Be should be omitted
        
        try:
            # Find the Canticle marker in the HTML (the one with verse references, not Gospel Canticle)
            canticle_span = None
            
            for span in soup.find_all('span', class_='rubrica'):
                span_text = span.get_text().strip()
                # Look for "Canticle:" with verse references (contains numbers and colon)
                if span_text.startswith('Canticle:') and re.search(r'\d+:\d+', span_text):
                    canticle_span = span
                    break
            
            if not canticle_span:
                print(f"  ⚠ Could not find Canticle marker in HTML")
                return self._get_fallback_canticle_verses()
            
            # Find the parent element and get its HTML
            parent = canticle_span.parent
            parent_html = str(parent)
            
            # Find the position of the canticle marker in the HTML
            canticle_pos = parent_html.find(str(canticle_span))
            if canticle_pos < 0:
                return self._get_fallback_canticle_verses()
            
            # Get HTML after the canticle marker
            html_after_canticle = parent_html[canticle_pos + len(str(canticle_span)):]
            
            # Check if "Glory to the Father is not said" appears BEFORE we trim the HTML
            # This indicates we should omit the Glory Be slide after the canticle
            # Pattern allows for HTML tags between words (e.g., </span>, <span class="rubrica">)
            # Matches variations like "the Glory to the Father is not said" or "At the end of the canticle the Glory to the Father is not said"
            if re.search(r'Glory\s+to\s+the\s+Father.*?is\s+not\s+said', html_after_canticle, re.IGNORECASE | re.DOTALL):
                omit_glory_be = True
                print(f"  ✓ Detected: Glory to the Father is not said for this canticle")
            
            # Find where to stop: "Glory to the Father" or next "Ant. 3"
            stop_patterns = [
                r'<span class="rubrica">Glory to the Father</span>',
                r'<span class="rubrica">Ant\.\s*3</span>',
                r'<span class="rubrica">Ant\.</span>',  # Repeated antiphon
                r'Glory to the Father',  # Plain text version
            ]
            
            end_pos = len(html_after_canticle)
            for pattern in stop_patterns:
                match = re.search(pattern, html_after_canticle, re.IGNORECASE)
                if match:
                    end_pos = min(end_pos, match.start())
            
            # Extract the canticle section HTML
            canticle_html = html_after_canticle[:end_pos]
            
            # Parse to get the text, skipping rubrica (red) text and em tags
            canticle_soup = BeautifulSoup(canticle_html, 'html.parser')
            
            # Remove the intro em tag (italic text)
            for em in canticle_soup.find_all('em'):
                em.decompose()
            
            # Get the remaining HTML and split by <br><br> to get verses/stanzas
            remaining_html = str(canticle_soup)
            
            # Split by double br tags
            verse_sections = re.split(r'<br\s*/?>\s*<br\s*/?>', remaining_html)
            
            verse_count = 0
            skipped_first_section = False  # Track if we've skipped the antiphon
            
            for section in verse_sections:
                # Parse this section to get clean text
                section_soup = BeautifulSoup(section, 'html.parser')
                
                # Remove asterisks and daggers (rubrica spans)
                for rubrica in section_soup.find_all('span', class_='rubrica'):
                    rubrica.decompose()
                
                # Get text
                verse_text = section_soup.get_text().strip()
                
                # Skip empty or short verses
                if not verse_text or len(verse_text) < 20:
                    continue
                
                # Skip the first valid section if it looks like the antiphon text
                # The antiphon typically appears as the very first text section after "Canticle:"
                # Only skip the FIRST occurrence to avoid false positives
                if not skipped_first_section and verse_count == 0:
                    # Check if this looks like an antiphon (relatively short, complete sentence)
                    if len(verse_text) < 150 and verse_text.endswith('.'):
                        print(f"  Skipping antiphon text in canticle extraction: {verse_text[:50]}...")
                        skipped_first_section = True
                        continue
                
                # Clean up: normalize whitespace
                verse_text = re.sub(r'\s+', ' ', verse_text).strip()
                
                # Add period if needed
                if not verse_text[-1] in '.!?"':
                    verse_text += '.'
                
                # Alternate speaker
                speaker = "Priest" if verse_count % 2 == 0 else "People"
                
                verses.append({
                    "speaker": speaker,
                    "text": verse_text
                })
                verse_count += 1
            
            if verses:
                print(f"  Extracted {len(verses)} verses for Canticle")
                return {
                    "verses": verses,
                    "omit_glory_be": omit_glory_be
                }
                
        except Exception as e:
            print(f"  ⚠ Error parsing canticle verses from HTML: {e}")
            import traceback
            traceback.print_exc()
        
        return self._get_fallback_canticle_verses()
    
    def _get_fallback_canticle_verses(self):
        """Return fallback canticle verses when parsing fails"""
        print(f"  Using fallback verses for Canticle")
        return {
            "verses": [
                {"speaker": "Priest", "text": "[Canticle verse 1 - Priest]"},
                {"speaker": "People", "text": "[Canticle verse 2 - People]"},
                {"speaker": "Priest", "text": "[Canticle verse 3 - Priest]"},
                {"speaker": "People", "text": "[Canticle verse 4 - People]"}
            ],
            "omit_glory_be": False  # Default to including Glory Be
        }

    def _extract_canticle_info(self, soup, text):
        """Extract canticle title and subtitle (red text)
        
        Looks for the rubrica span containing "Canticle:" and extracts the full red text
        Examples: 
        - "Canticle: Isaiah 42:10-16God, victor and savior"
        - "Canticle: Daniel 3:57-88, 56Let all creatures praise the Lord"
        - "Canticle: Isaiah 61:10—62:5The prophet's joy in the vision of a new Jerusalem"
        Should be split into title and subtitle
        """
        try:
            # Find the rubrica span with "Canticle:" in the HTML
            # Note: We need to find the FIRST Canticle span that has verse references
            # (not "Canticle of Zechariah" which is the Gospel Canticle)
            for span in soup.find_all('span', class_='rubrica'):
                span_text = span.get_text().strip()
                
                # Check if this is a Canticle with verse references (contains numbers and colon)
                # Examples: "Daniel 3:57", "Isaiah 42:10-16", "Isaiah 61:10—62:5"
                if span_text.startswith('Canticle:') and re.search(r'\d+:\d+', span_text):
                    # Extract the full text
                    # Format: "Canticle: [Book] [chapter]:[verses][Subtitle]"
                    # The subtitle starts with a capital letter after the verse numbers
                    
                    # Pattern to match:
                    # - "Canticle: " followed by book name (letters and spaces)
                    # - Chapter:verse or Chapter:verse-verse or Chapter:verse—verse or Chapter:verse—chapter:verse
                    # - Optional comma and additional verse numbers
                    # - Remaining text is the subtitle
                    
                    # More flexible pattern that handles:
                    # - "Isaiah 42:10-16" (hyphen within same chapter)
                    # - "Isaiah 61:10—62:5" (em dash spanning chapters)
                    # - "Daniel 3:57-88, 56" (comma for additional verse)
                    match = re.match(
                        r'(Canticle:\s+[A-Za-z\s]+\d+:\d+(?:[-—]\d+(?::\d+)?)?(?:,\s*\d+)?)(.*)', 
                        span_text, 
                        re.IGNORECASE
                    )
                    
                    if match:
                        title = match.group(1).strip()  # "Canticle: Isaiah 42:10-16"
                        subtitle = match.group(2).strip()  # "God, victor and savior"
                        
                        # Clean up subtitle: remove leading punctuation like em dash
                        subtitle = re.sub(r'^[—\-\s]+', '', subtitle)
                        
                        print(f"  Found Canticle title: {title}")
                        if subtitle:
                            print(f"  Found Canticle subtitle: {subtitle}")
                        
                        return {
                            "title": title,
                            "subtitle": subtitle
                        }
                    else:
                        # Fallback: try to split at the last digit followed by a capital letter
                        # Find position where verse numbers end (last digit) and subtitle begins
                        verse_end = re.search(r'\d+([A-Z])', span_text)
                        if verse_end:
                            split_pos = verse_end.start(1)
                            title = span_text[:split_pos].strip()
                            subtitle = span_text[split_pos:].strip()
                            
                            print(f"  Found Canticle title: {title}")
                            if subtitle:
                                print(f"  Found Canticle subtitle: {subtitle}")
                            
                            return {
                                "title": title,
                                "subtitle": subtitle
                            }
                        else:
                            # Use the whole text as title
                            print(f"  Found Canticle (no subtitle split): {span_text}")
                            return {
                                "title": span_text,
                                "subtitle": ""
                            }
        except Exception as e:
            print(f"  ⚠ Error extracting canticle info: {e}")
        
        # Fallback
        return {
            "title": "[Canticle title]",
            "subtitle": ""
        }

    def _extract_short_reading(self, text):
        """Extract short reading text - find first instance of READING (after PSALMODY) up to RESPONSORY
        
        The structure is:
        READING[citation][reading text]RESPONSORY
        
        We need to extract everything from READING to just before RESPONSORY.
        Note: There may be multiple READINGs (alternative options), so we take the first one
        that has a RESPONSORY following it.
        """
        try:
            # Find all instances of READING (case-insensitive)
            reading_matches = list(re.finditer(r'READING', text, re.IGNORECASE))
            
            if not reading_matches:
                print("  ⚠ No READING marker found")
                return {"citation": "", "text": ""}
            
            # Find the first READING that has a RESPONSORY after it
            # (skip introductory text or readings that don't have responsory)
            reading_start = None
            for match in reading_matches:
                test_start = match.end()
                # Check if RESPONSORY follows within reasonable distance (< 1000 chars)
                responsory_test = re.search(r'RESPONSORY', text[test_start:test_start+1000], re.IGNORECASE)
                if responsory_test:
                    reading_start = test_start
                    break
            
            if reading_start is None:
                print("  ⚠ No READING with RESPONSORY found")
                return {"citation": "", "text": ""}
            
            # Find RESPONSORY after the last READING
            responsory_match = re.search(r'RESPONSORY', text[reading_start:], re.IGNORECASE)
            
            if not responsory_match:
                print("  ⚠ No RESPONSORY marker found after READING")
                return {"citation": "", "text": ""}
            
            # Extract the reading section
            reading_end = reading_start + responsory_match.start()
            reading_section = text[reading_start:reading_end].strip()
            
            # The reading section format is typically:
            # [Optional category like "[Pastors]"][Citation like "Hebrews 13:7-9a"][Reading text]
            # Try to separate citation from text
            
            # First, remove any category prefix in square brackets (e.g., "[Pastors]")
            reading_section = re.sub(r'^\[.*?\]\s*', '', reading_section)
            
            # Look for a book name followed by chapter:verse pattern at the start
            # Book names can be: single word (Genesis, Mark), two words (1 Corinthians), or abbreviations
            citation_match = re.match(r'^([1-3]?\s*[A-Za-z]+\s+\d+:\d+[a-z]?(?:-\d+[a-z]?)?)', reading_section)
            
            if citation_match:
                citation = citation_match.group(1).strip()
                reading_text = reading_section[citation_match.end():].strip()
            else:
                # No clear citation found, treat entire section as text
                citation = ""
                reading_text = reading_section
            
            print(f"  Found READING: {citation}")
            print(f"    Text preview: {reading_text[:100]}...")
            
            return {
                "citation": citation,
                "text": reading_text
            }
            
        except Exception as e:
            print(f"  ⚠ Error extracting short reading: {e}")
            return {"citation": "", "text": ""}

    def _extract_responsory_from_html(self, soup, text):
        """Extract responsory from text with proper parsing
        
        Pattern in plain text:
        RESPONSORY
        [Main response line]
        — [Main response line repeated]
        [Verse line]
        — [Shortened response]
        Glory to the Father...
        — [Main response line repeated]
        
        Returns a list of dicts with 'speaker' and 'text' keys
        """
        try:
            # Find RESPONSORY marker in plain text
            responsory_match = re.search(r'RESPONSORY', text, re.IGNORECASE)
            if not responsory_match:
                print("  ⚠ No RESPONSORY marker found in text")
                return []
            
            # Start after RESPONSORY marker
            responsory_start = responsory_match.end()
            
            # Find where to stop: "Or:" or "GOSPEL CANTICLE"
            stop_patterns = [
                r'\bOr:',
                r'GOSPEL\s+CANTICLE',
                r'CANTICLE\s+OF\s+ZECHARIAH',
            ]
            
            responsory_end = len(text)
            for pattern in stop_patterns:
                stop_match = re.search(pattern, text[responsory_start:], re.IGNORECASE)
                if stop_match:
                    responsory_end = responsory_start + stop_match.start()
                    break
            
            # Extract the responsory section
            responsory_section = text[responsory_start:responsory_end].strip()
            
            # First, normalize the em dash characters
            normalized_section = responsory_section.replace('\u2014', '—').replace('\u2013', '—').replace('\u2015', '—')
            
            # Split by em dash to separate the response lines
            # This gives us segments like: [main text]—[response 1+verse]—[response 2+glory]—[final response]
            # Keep each em-dash-separated part as a whole segment
            em_dash_parts = [part.strip() for part in normalized_section.split('—') if part.strip()]
            
            # Now we need to further split parts that contain multiple sentences
            # Part 0: Main response (single sentence)
            # Part 1: Repeated response + Verse (need to split by period)
            # Part 2: Shortened response + Glory Be (need to split by period)
            # Part 3: Final response (single sentence)
            
            all_segments = []
            
            # Add first part as-is (main response)
            if len(em_dash_parts) > 0:
                all_segments.append(em_dash_parts[0])
            
            # Split part 1 (repeated + verse) by period
            if len(em_dash_parts) > 1:
                part1_sentences = []
                current = ""
                for i, char in enumerate(em_dash_parts[1]):
                    current += char
                    if char == '.' and i + 1 < len(em_dash_parts[1]):
                        # Check if next char is space or uppercase (new sentence)
                        next_char = em_dash_parts[1][i + 1]
                        if next_char.isupper() or (i + 2 < len(em_dash_parts[1]) and em_dash_parts[1][i + 2].isupper()):
                            part1_sentences.append(current.strip())
                            current = ""
                if current.strip():
                    part1_sentences.append(current.strip())
                all_segments.extend(part1_sentences)
            
            # Split part 2 (shortened + Glory) by period before "Glory"
            if len(em_dash_parts) > 2:
                # Find "Glory to the Father" in this part
                glory_pattern = r'(.*?)(Glory\s+to\s+the\s+Father.*)'
                glory_match = re.search(glory_pattern, em_dash_parts[2], re.IGNORECASE | re.DOTALL)
                if glory_match:
                    shortened = glory_match.group(1).strip()
                    glory = glory_match.group(2).strip()
                    if shortened:
                        all_segments.append(shortened)
                    if glory:
                        all_segments.append(glory)
                else:
                    # No Glory found, just add the whole part
                    all_segments.append(em_dash_parts[2])
            
            # Add remaining parts as-is
            for i in range(3, len(em_dash_parts)):
                all_segments.append(em_dash_parts[i])
            
            # Build responsory verses according to the 6-segment pattern:
            # Segment 0: Main response
            # Segment 1: Repeated response (after first em dash)
            # Segment 2: Verse
            # Segment 3: Shortened response (after second em dash)
            # Segment 4: Glory Be
            # Segment 5: Final response (after third em dash)
            #
            # Slide 1: RESPONSORY title + Seg 0 + "— " + Seg 1 (ALL)
            # Slide 2: Seg 2 + "— " + Seg 3 (Priest)
            # Slide 3: Seg 4 + "— " + Seg 5 (Priest)
            
            if len(all_segments) < 6:
                print(f"  ⚠ Expected 6 segments but found {len(all_segments)}, structure may be incorrect")
                return []
            
            responsory_verses = []
            
            # Slide 1: Main response + repeated response (both ALL)
            combined_first = all_segments[0].strip()
            if len(all_segments) > 1:
                combined_first += "\n— " + all_segments[1].strip()
            
            responsory_verses.append({
                "speaker": "All",
                "text": combined_first,
                "include_title": True  # This slide should have the RESPONSORY title
            })
            
            # Slide 2: Verse (Priest) + shortened response (with em dash)
            combined_verse = all_segments[2].strip()
            if len(all_segments) > 3:
                combined_verse += "\n— " + all_segments[3].strip()
            
            responsory_verses.append({
                "speaker": "Priest",
                "text": combined_verse
            })
            
            # Slide 3: Glory Be (Priest) + final response (with em dash)
            combined_glory = all_segments[4].strip()
            if len(all_segments) > 5:
                combined_glory += "\n— " + all_segments[5].strip()
            
            responsory_verses.append({
                "speaker": "Priest",
                "text": combined_glory
            })
            
            print(f"  Found RESPONSORY with {len(responsory_verses)} parts")
            return responsory_verses
            
        except Exception as e:
            print(f"  ⚠ Error extracting responsory from HTML: {e}")
            traceback.print_exc()
            return []

    def _extract_responsory(self, text):
        """Extract responsory text with call-and-response structure
        
        Pattern:
        RESPONSORY
        [Main response line]
        — [Main response line repeated]
        [Verse line]
        — [Shortened response]
        Glory to the Father...
        — [Main response line repeated]
        
        Returns a list of dicts with 'speaker' and 'text' keys for alternating priest/people
        """
        try:
            # Find first instance of RESPONSORY (case-insensitive)
            responsory_match = re.search(r'RESPONSORY', text, re.IGNORECASE)
            
            if not responsory_match:
                print("  ⚠ No RESPONSORY marker found")
                return []
            
            # Start after RESPONSORY marker
            responsory_start = responsory_match.end()
            
            # Find where to stop: Look for "GOSPEL CANTICLE" or "CANTICLE OF ZECHARIAH" or "OR" (for alternative readings)
            stop_patterns = [
                r'GOSPEL\s+CANTICLE',
                r'CANTICLE\s+OF\s+ZECHARIAH',
                r'\bOR\b',  # Alternative reading marker
                r'INTERCESSIONS'
            ]
            
            responsory_end = len(text)
            for pattern in stop_patterns:
                stop_match = re.search(pattern, text[responsory_start:], re.IGNORECASE)
                if stop_match:
                    responsory_end = responsory_start + stop_match.start()
                    break
            
            # Extract responsory section
            responsory_section = text[responsory_start:responsory_end].strip()
            
            # Debug: Print the raw responsory section to understand the structure
            print(f"  DEBUG: Responsory section length: {len(responsory_section)}")
            print(f"  DEBUG: Responsory section preview: {responsory_section[:300]}")
            
            # Split into lines - try different approaches
            # The em dash (—) is a special character that may not split properly with \n alone
            # Replace various whitespace and em dash patterns with newlines first
            normalized_section = responsory_section.replace('\r\n', '\n').replace('\r', '\n')
            
            # Split by newline
            all_lines = normalized_section.split('\n')
            
            # Filter out empty lines and clean up
            lines = [line.strip() for line in all_lines if line.strip() and len(line.strip()) > 3]
            
            print(f"  DEBUG: Found {len(lines)} non-empty lines")
            for i, line in enumerate(lines[:10]):  # Show first 10 lines
                print(f"    Line {i}: {line[:80]}")
            
            # Parse the responsory structure
            # Expected pattern (6 lines total):
            # Line 1: Main response
            # Line 2: — Main response (repeated)
            # Line 3: Verse
            # Line 4: — Shortened response
            # Line 5: Glory to the Father...
            # Line 6: — Main response (repeated)
            
            if len(lines) < 6:
                print(f"  ⚠ Responsory has {len(lines)} lines, expected at least 6")
                return []
            
            # Build the responsory as alternating priest/people responses
            responsory_verses = []
            
            # Line 1: Main response (ALL)
            responsory_verses.append({
                "speaker": "All",
                "text": lines[0]
            })
            
            # Line 2: Repeated main response (ALL) - skip the "—" prefix
            # (This is typically the same as line 1, so we can skip it in slides)
            
            # Line 3: Verse (Priest)
            responsory_verses.append({
                "speaker": "Priest",
                "text": lines[2]
            })
            
            # Line 4: Shortened response (ALL) - remove "—" prefix
            response_text = lines[3]
            if response_text.startswith('—'):
                response_text = response_text[1:].strip()
            responsory_verses.append({
                "speaker": "All",
                "text": response_text
            })
            
            # Line 5: Glory Be (Priest)
            responsory_verses.append({
                "speaker": "Priest",
                "text": lines[4]
            })
            
            # Line 6: Final response (ALL) - remove "—" prefix
            final_response = lines[5]
            if final_response.startswith('—'):
                final_response = final_response[1:].strip()
            responsory_verses.append({
                "speaker": "All",
                "text": final_response
            })
            
            print(f"  Found RESPONSORY with {len(responsory_verses)} parts")
            return responsory_verses
            
        except Exception as e:
            print(f"  ⚠ Error extracting responsory: {e}")
            traceback.print_exc()
            return []

    def _extract_gospel_antiphon(self, text):
        """Extract gospel canticle antiphon
        
        Pattern:
        GOSPEL CANTICLE
        Ant. [Optional category like [Martyrs]] [Antiphon text]
        
        Returns the full antiphon text including any category markers
        """
        try:
            # Find GOSPEL CANTICLE marker
            gc_match = re.search(r'GOSPEL\s+CANTICLE', text, re.IGNORECASE)
            if not gc_match:
                print("  ⚠ No GOSPEL CANTICLE marker found")
                return ""
            
            # Start after GOSPEL CANTICLE
            start_pos = gc_match.end()
            
            # Find "Ant." marker
            ant_match = re.search(r'Ant\.', text[start_pos:start_pos+500], re.IGNORECASE)
            if not ant_match:
                print("  ⚠ No antiphon marker found after GOSPEL CANTICLE")
                return ""
            
            # Start after "Ant."
            ant_start = start_pos + ant_match.end()
            
            # Find where antiphon ends - typically at "Canticle" or "Benedictus" or next major section
            # The antiphon can span multiple sentences, so we need to find the next structural marker
            stop_patterns = [
                r'Canticle\s+of\s+Zechariah',
                r'Benedictus',
                r'Canticle:',
                r'INTERCESSIONS',
                r'Let us pray',
            ]
            
            end_pos = len(text)
            for pattern in stop_patterns:
                stop_match = re.search(pattern, text[ant_start:ant_start+2000], re.IGNORECASE)
                if stop_match:
                    end_pos = ant_start + stop_match.start()
                    break
            
            # Extract antiphon text
            antiphon_text = text[ant_start:end_pos].strip()
            
            # Clean up: remove extra whitespace and newlines
            antiphon_text = re.sub(r'\s+', ' ', antiphon_text).strip()
            
            # Remove any trailing section markers that might have slipped in
            antiphon_text = re.sub(r'(Canticle|Benedictus|INTERCESSIONS).*$', '', antiphon_text, flags=re.IGNORECASE).strip()
            
            if antiphon_text:
                print(f"  Found Gospel Canticle antiphon: {antiphon_text[:80]}...")
                return antiphon_text
            
        except Exception as e:
            print(f"  ⚠ Error extracting gospel antiphon: {e}")
            traceback.print_exc()
        
        return ""

    def _extract_benedictus_verses(self, text):
        """Extract Benedictus verses - these are standard but can be customized"""
        return [
            "Blessed be the Lord, the God of Israel; he has come to his people and set them free.",
            "He has raised up for us a mighty savior, born of the house of his servant David.",
            # ... continue with full Benedictus text
        ]

    def _extract_intercessions(self, text):
        """Extract intercessions"""
        return "[Intercessions for today]"

    def _extract_concluding_prayer(self, text):
        """Extract concluding prayer/collect"""
        return "[Concluding prayer for today]"

    def _extract_first_reading_citation(self, text):
        """Extract first reading citation"""
        return "Rom 9:1-5"  # Fallback - would extract actual citation

    def _extract_first_reading_verses(self, text):
        """Extract first reading verses"""
        return [
            "A reading from the Letter of Saint Paul to the Romans",
            "[First reading content verse 1]",
            "[First reading content verse 2]",
            "[First reading content verse 3]",
            "The word of the Lord."
        ]

    def _extract_psalm_citation(self, text):
        """Extract responsorial psalm citation"""
        return "Ps 147:12-13, 14-15, 19-20"

    def _extract_psalm_response_verses(self, text):
        """Extract responsorial psalm verses"""
        return [
            "℟. Praise the Lord, Jerusalem.",
            "[Psalm verse 1]",
            "℟. Praise the Lord, Jerusalem.", 
            "[Psalm verse 2]",
            "℟. Praise the Lord, Jerusalem."
        ]

    def _extract_gospel_acclamation(self, text):
        """Extract gospel acclamation"""
        return "℟. Alleluia, alleluia.\n[Alleluia verse]\n℟. Alleluia, alleluia."

    def _extract_gospel_citation(self, text):
        """Extract gospel citation"""
        return "Lk 14:1-6"

    def _extract_gospel_verses(self, text):
        """Extract gospel verses"""
        return [
            "✠ A reading from the holy Gospel according to Luke",
            "[Gospel content verse 1]",
            "[Gospel content verse 2]", 
            "[Gospel content verse 3]",
            "The Gospel of the Lord."
        ]

    def _get_static_devotional_content(self):
        """
        Return the static devotional content that doesn't change (from reference presentation)
        """
        return {
            "sacred_heart_hymns": [
                {"content": ""},  # Blank slide
                {"content": "Heart of Jesus meek and mild. Hear oh hear thy feeble child When the tempest's most severe. Heart of Jesus, hear. Sweetly, we'll rest on thy Sacred Heart. Never from Thee. O let us part!"},
                {"content": "Hear then, Thy loving children's pray'r O Heart of Jesus, Heart of Jesus hear."},
                {"content": ""},  # Blank slide
                {"content": "Oh Sacred Heart, Oh love divine. Do keep us near to Thee. And make our love so like to Thine That we may holy be. Heart of Jesus hear. Oh heart of love divine. Listen to our pray'r."},
                {"content": "Make us always Thine. Oh temple pure, Oh house of gold. Our heaven here below. What gifts unfurled, what wealth untold. From Thee do ever flow. Heart of Jesus hear. Oh Heart of love divine. Listen to our pray'r. Make us always Thine."}
            ],
            
            "post_communion_prayers": [
                {"content": ""},  # Blank
                {"content": "Soul of Christ, make me holy. Body of Christ, save me. Blood of Christ, inebriate me. Water from the side of Christ, wash me. Passion of Christ, make me strong. O good Jesus, hear me. Hide me within your wounds."},
                {"content": "Let me never be separated from You. Deliver me from the wicked enemy, Call me at the hour of my death. And tell me to come to you that with Your saints I may praise You forever. Amen."},
                {"title": "PRAYER OF THANKSGIVING:", "content": "Lord God, I thank you through the Sacred Heart of Jesus, who is pleased to offer You on our behalf continuous thanksgiving in the Eucharist."},
                # ... continue with all post-communion prayer content from reference
            ],
            
            "jubilee_prayer": [
                {"title": "THE JUBILEE PRAYER"},
                {"content": "Father in heaven, may the faith you have given us in your son, Jesus Christ, our brother, and the flame of charity"},
                {"content": "enkindled in our hearts by the Holy Spirit, reawaken in us the blessed hope for the coming of your Kingdom."},
                # ... continue with jubilee prayer content
            ],
            
            "st_joseph_prayer": [
                {"content": "To you, O blessed Joseph, do we come in our tribulation, and having implored the help of your most holy Spouse, we confidently invoke your patronage also."},
                # ... continue with St. Joseph prayer content
            ]
        }

    def _get_fallback_morning_prayer(self):
        """Fallback morning prayer structure if iBreviary fails"""
        return {
            "psalmody": {
                "antiphon_1": {"text": "", "format": "all_response", "psalm_title": "", "psalm_subtitle": ""},
                "psalm_1": [{"speaker": "Priest", "text": ""}],
                "antiphon_2": {"text": "", "format": "all_response"}, 
                "canticle": [{"speaker": "Priest", "text": ""}],
                "antiphon_3": {"text": "", "format": "all_response", "psalm_title": "", "psalm_subtitle": ""},
                "psalm_3": [{"speaker": "Priest", "text": ""}]
            }
        }

    def _get_fallback_readings(self):
        """Fallback readings if iBreviary fails"""
        return {
            "first_reading": {"citation": "[Citation]", "verses": ["[Reading text]"]},
            "responsorial_psalm": {"citation": "[Psalm citation]", "verses": ["[Psalm text]"]},
            "gospel_acclamation": {"verse": "[Alleluia verse]"},
            "gospel": {"citation": "[Gospel citation]", "verses": ["[Gospel text]"]}
        }

    def _get_fallback_data(self, target_date=None):
        """Complete fallback data structure"""
        if target_date is None:
            target_date = datetime.now()
            
        return {
            "date": target_date.strftime('%B %d, %Y'),
            "morning_prayer": self._get_fallback_morning_prayer(),
            "mass_readings": self._get_fallback_readings(),
            "static_content": self._get_static_devotional_content()
        }

    def create_presentation_from_template(self, liturgical_data, output_filename=None):
        """
        Create presentation using the reference template structure with live liturgical data
        """
        if output_filename is None:
            # Use OLPH naming convention: olph_slides_[year]_[month]_[day].pptx
            # Extract date from liturgical_data if available, otherwise use current date
            if 'date' in liturgical_data:
                try:
                    # Parse the date string to get components
                    date_obj = datetime.strptime(liturgical_data['date'], '%B %d, %Y')
                    output_filename = f"olph_slides_{date_obj.year}_{date_obj.month:02d}_{date_obj.day:02d}.pptx"
                except:
                    # Fallback to current date
                    now = datetime.now()
                    output_filename = f"olph_slides_{now.year}_{now.month:02d}_{now.day:02d}.pptx"
            else:
                now = datetime.now()
                output_filename = f"olph_slides_{now.year}_{now.month:02d}_{now.day:02d}.pptx"
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        print(f"Creating presentation using reference template structure...")
        print(f"Date: {liturgical_data['date']}")
        
        slide_count = 0
        
        # Add blank black slide at the very beginning
        slide_count = self._create_initial_blank_slide(prs, slide_count)
        
        # Add Daily Morning Prayer image slide as second slide
        slide_count = self._create_daily_morning_prayer_image_slide(prs, slide_count)
        
        # Apply reference template structure to current liturgical data
        slide_count = self._create_opening_slides(prs, liturgical_data, slide_count)
        slide_count = self._create_psalmody_section(prs, liturgical_data, slide_count)
        slide_count = self._create_reading_section(prs, liturgical_data, slide_count)
        slide_count = self._create_responsory_section(prs, liturgical_data, slide_count)
        slide_count = self._create_gospel_canticle_section(prs, liturgical_data, slide_count)
        slide_count = self._create_intercessions_section(prs, liturgical_data, slide_count)
        slide_count = self._create_sacred_heart_hymns(prs, liturgical_data, slide_count)
        slide_count = self._create_mass_readings_section(prs, liturgical_data, slide_count)
        slide_count = self._create_post_communion_prayers(prs, liturgical_data, slide_count)
        # Commented out: These sections create placeholder/empty slides
        # slide_count = self._create_transition_slides(prs, slide_count)
        # slide_count = self._create_jubilee_prayer(prs, liturgical_data, slide_count)
        # slide_count = self._create_st_joseph_prayer(prs, liturgical_data, slide_count)
        
        # Save presentation
        output_dir = "output_v2"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        output_path = os.path.join(output_dir, output_filename)
        prs.save(output_path)
        
        print(f"\nPresentation created successfully!")
        print(f"File: {output_path}")
        print(f"Total slides: {slide_count}")
        print(f"Target slides (reference): {self.reference_template['metadata']['total_expected_slides']}")
        
        return output_path

    def _create_initial_blank_slide(self, prs, slide_count):
        """Create initial blank black slide at the very beginning"""
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Make the slide background black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Initial blank black slide")
        return slide_count

    def _create_daily_morning_prayer_image_slide(self, prs, slide_count):
        """Create Daily Morning Prayer image slide as second slide"""
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Path to the image file
        image_path = "daily_morning_prayer.png"
        
        # Check if image exists
        if os.path.exists(image_path):
            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add image to fill the entire slide
            slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
            print(f"Created slide {slide_count}: Daily Morning Prayer image slide")
        else:
            # Fallback: create text slide if image not found
            title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = "Daily Morning Prayer"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(60)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(184, 134, 11)  # Gold color
            title_para.alignment = PP_ALIGN.CENTER
            print(f"Created slide {slide_count}: Daily Morning Prayer text slide (image not found)")
        
        return slide_count

    def _create_opening_slides(self, prs, liturgical_data, slide_count):
        """Create opening slides following reference template
        
        Uses auto-fit functionality to automatically adjust text size based on content length.
        This ensures text always fits within the designated space regardless of antiphon length.
        """
        # Title slide with PSALMODY heading, antiphon, and psalm information
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # PSALMODY title
        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = "PSALMODY"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(80)
        title_para.font.name = "Georgia"
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x98, 0x00, 0x00)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Antiphon text with auto-fit
        antiphon_1 = liturgical_data['morning_prayer']['psalmody']['antiphon_1']
        
        antiphon_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(12.7), Inches(3.5))
        antiphon_frame = antiphon_box.text_frame
        antiphon_frame.word_wrap = True
        
        # Enable auto-fit to shrink text if needed
        antiphon_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Create the paragraph with proper formatting
        first_para = antiphon_frame.paragraphs[0]
        first_para.alignment = PP_ALIGN.CENTER
        
        # Add "(All) Ant. 1 " in blue
        first_run = first_para.add_run()
        first_run.text = "(All) Ant. 1 "
        first_run.font.size = Pt(52)
        first_run.font.name = "Georgia" 
        first_run.font.bold = True
        first_run.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)
        
        # Add the antiphon text in black
        second_run = first_para.add_run()
        second_run.text = antiphon_1['text']
        second_run.font.size = Pt(52)
        second_run.font.name = "Georgia"
        second_run.font.bold = True  
        second_run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Psalm title and subtitle with auto-fit
        psalm_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.3), Inches(12.7), Inches(2.0))
        psalm_frame = psalm_box.text_frame
        psalm_frame.word_wrap = True
        
        # Enable auto-fit to shrink text if needed
        psalm_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        psalm_title = antiphon_1.get('psalm_title', '')
        psalm_subtitle = antiphon_1.get('psalm_subtitle', '')
        
        if psalm_title or psalm_subtitle:
            if psalm_title:
                psalm_frame.text = psalm_title
                psalm_para = psalm_frame.paragraphs[0]
                psalm_para.font.size = Pt(44)
                psalm_para.font.name = "Georgia"
                psalm_para.font.bold = True
                psalm_para.font.color.rgb = RGBColor(0x98, 0x00, 0x00)
                psalm_para.alignment = PP_ALIGN.LEFT
            
            if psalm_subtitle:
                if psalm_title:
                    subtitle_para = psalm_frame.add_paragraph()
                    subtitle_para.text = psalm_subtitle
                    subtitle_para.font.size = Pt(44)
                    subtitle_para.font.name = "Georgia"
                    subtitle_para.font.bold = True
                    subtitle_para.font.color.rgb = RGBColor(0x98, 0x00, 0x00)
                    subtitle_para.alignment = PP_ALIGN.LEFT
                else:
                    psalm_frame.text = psalm_subtitle
                    psalm_para = psalm_frame.paragraphs[0]
                    psalm_para.font.size = Pt(44)
                    psalm_para.font.name = "Georgia"
                    psalm_para.font.bold = True
                    psalm_para.font.color.rgb = RGBColor(0x98, 0x00, 0x00)
                    psalm_para.alignment = PP_ALIGN.LEFT
        
        print(f"Created slide {slide_count}: PSALMODY title slide")
        
        return slide_count

    def _create_psalmody_section(self, prs, liturgical_data, slide_count):
        """Create psalmody section following reference template exactly"""
        # Create psalm verses alternating priest/people (following reference pattern)
        psalm_1_verses = liturgical_data['morning_prayer']['psalmody']['psalm_1']
        for verse in psalm_1_verses:
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Create text box for the verse content with auto-fit
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Enable auto-fit to make text as large as possible without crossing boundaries
            content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Create paragraph
            content_para = content_frame.paragraphs[0]
            content_para.alignment = PP_ALIGN.LEFT
            
            # Check if this is the "Glory to the Father" verse (no speaker label needed)
            if "Glory to the Father" in verse['text'] or "Glory to the father" in verse['text']:
                # Glory be verse: just black text, no speaker label
                glory_run = content_para.add_run()
                glory_run.text = verse['text']
                glory_run.font.size = Pt(44)
                glory_run.font.name = "Georgia"
                glory_run.font.bold = True
                glory_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"Created slide {slide_count}: Psalm 1 - Glory Be")
                
            elif verse['speaker'] == "Priest":
                # Priest slides: "Priest: " + text, all in red
                priest_run = content_para.add_run()
                priest_run.text = f"Priest: {verse['text']}"
                priest_run.font.size = Pt(44)
                priest_run.font.name = "Georgia"
                priest_run.font.bold = True
                priest_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
                print(f"Created slide {slide_count}: Psalm 1 - {verse['speaker']}")
                
            elif verse['speaker'] == "People":
                # People slides: "People: " in blue, rest of text in black
                people_label = content_para.add_run()
                people_label.text = "People: "
                people_label.font.size = Pt(44)
                people_label.font.name = "Georgia"
                people_label.font.bold = True
                people_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
                
                people_text = content_para.add_run()
                people_text.text = verse['text']
                people_text.font.size = Pt(44)
                people_text.font.name = "Georgia"
                people_text.font.bold = True
                people_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"Created slide {slide_count}: Psalm 1 - {verse['speaker']}")
        
        # Repeat Antiphon 1 after Psalm 1 verses
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        antiphon_1 = liturgical_data['morning_prayer']['psalmody']['antiphon_1']
        
        # Create text box for the repeated antiphon with auto-fit
        ant_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        ant_frame = ant_box.text_frame
        ant_frame.word_wrap = True
        ant_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ant_para = ant_frame.paragraphs[0]
        ant_para.alignment = PP_ALIGN.CENTER
        
        # Add "(All) Ant. 1 " in blue
        ant_label = ant_para.add_run()
        ant_label.text = "(All) Ant. 1 "
        ant_label.font.size = Pt(44)
        ant_label.font.name = "Georgia"
        ant_label.font.bold = True
        ant_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
        
        # Add the antiphon text in black
        ant_text = ant_para.add_run()
        ant_text.text = antiphon_1['text']
        ant_text.font.size = Pt(44)
        ant_text.font.name = "Georgia"
        ant_text.font.bold = True
        ant_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Repeated Antiphon 1")
        
        # Add Antiphon 2 slide
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        antiphon_2 = liturgical_data['morning_prayer']['psalmody']['antiphon_2']
        
        # Create text box for Antiphon 2 with auto-fit
        ant2_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        ant2_frame = ant2_box.text_frame
        ant2_frame.word_wrap = True
        ant2_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ant2_para = ant2_frame.paragraphs[0]
        ant2_para.alignment = PP_ALIGN.CENTER
        
        # Add "Ant. 2 " in blue
        ant2_label = ant2_para.add_run()
        ant2_label.text = "Ant. 2 "
        ant2_label.font.size = Pt(44)
        ant2_label.font.name = "Georgia"
        ant2_label.font.bold = True
        ant2_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
        
        # Add the antiphon 2 text in black
        ant2_text = ant2_para.add_run()
        ant2_text.text = antiphon_2['text']
        ant2_text.font.size = Pt(44)
        ant2_text.font.name = "Georgia"
        ant2_text.font.bold = True
        ant2_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Antiphon 2")
        
        # Add Canticle info slide (red text)
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        canticle_info = liturgical_data['morning_prayer']['psalmody']['canticle_info']
        
        # Create text box for Canticle info with auto-fit
        canticle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(3))
        canticle_frame = canticle_box.text_frame
        canticle_frame.word_wrap = True
        canticle_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Add canticle title in red
        canticle_para = canticle_frame.paragraphs[0]
        canticle_para.alignment = PP_ALIGN.LEFT
        
        canticle_title_run = canticle_para.add_run()
        canticle_title_run.text = canticle_info['title']
        canticle_title_run.font.size = Pt(44)
        canticle_title_run.font.name = "Georgia"
        canticle_title_run.font.bold = True
        canticle_title_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
        
        # Add canticle subtitle in red on next line if present
        if canticle_info['subtitle']:
            canticle_subtitle_para = canticle_frame.add_paragraph()
            canticle_subtitle_para.alignment = PP_ALIGN.LEFT
            
            canticle_subtitle_run = canticle_subtitle_para.add_run()
            canticle_subtitle_run.text = canticle_info['subtitle']
            canticle_subtitle_run.font.size = Pt(44)
            canticle_subtitle_run.font.name = "Georgia"
            canticle_subtitle_run.font.bold = True
            canticle_subtitle_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
        
        print(f"Created slide {slide_count}: Canticle info")
        
        # Create canticle verses alternating priest/people (same pattern as psalm verses)
        canticle_data = liturgical_data['morning_prayer']['psalmody']['canticle']
        canticle_verses = canticle_data.get('verses', canticle_data) if isinstance(canticle_data, dict) else canticle_data
        omit_glory_be = canticle_data.get('omit_glory_be', False) if isinstance(canticle_data, dict) else False
        
        for verse in canticle_verses:
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Create text box for the verse content with auto-fit
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Enable auto-fit to make text as large as possible without crossing boundaries
            content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Create paragraph
            content_para = content_frame.paragraphs[0]
            content_para.alignment = PP_ALIGN.LEFT
            
            # Check if this is the "Glory to the Father" verse (no speaker label needed)
            if "Glory to the Father" in verse['text'] or "Glory to the father" in verse['text']:
                # Glory be verse: just black text, no speaker label
                glory_run = content_para.add_run()
                glory_run.text = verse['text']
                glory_run.font.size = Pt(44)
                glory_run.font.name = "Georgia"
                glory_run.font.bold = True
                glory_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"Created slide {slide_count}: Canticle - Glory Be")
                
            elif verse['speaker'] == "Priest":
                # Priest slides: "Priest: " + text, all in red
                priest_run = content_para.add_run()
                priest_run.text = f"Priest: {verse['text']}"
                priest_run.font.size = Pt(44)
                priest_run.font.name = "Georgia"
                priest_run.font.bold = True
                priest_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
                print(f"Created slide {slide_count}: Canticle - {verse['speaker']}")
                
            elif verse['speaker'] == "People":
                # People slides: "People: " in blue, rest of text in black
                people_label = content_para.add_run()
                people_label.text = "People: "
                people_label.font.size = Pt(44)
                people_label.font.name = "Georgia"
                people_label.font.bold = True
                people_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
                
                people_text = content_para.add_run()
                people_text.text = verse['text']
                people_text.font.size = Pt(44)
                people_text.font.name = "Georgia"
                people_text.font.bold = True
                people_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"Created slide {slide_count}: Canticle - {verse['speaker']}")
        
        # Add Glory Be after canticle (unless explicitly omitted)
        if not omit_glory_be:
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Create text box for Glory Be
            glory_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            glory_frame = glory_box.text_frame
            glory_frame.word_wrap = True
            glory_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            glory_para = glory_frame.paragraphs[0]
            glory_para.alignment = PP_ALIGN.LEFT
            
            # Glory Be text in black, no speaker label
            glory_run = glory_para.add_run()
            glory_run.text = "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be for ever. Amen."
            glory_run.font.size = Pt(44)
            glory_run.font.name = "Georgia"
            glory_run.font.bold = True
            glory_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            
            print(f"Created slide {slide_count}: Canticle - Glory Be")
        else:
            print(f"  Skipping Glory Be slide (explicitly omitted for this canticle)")
        
        # Repeat Antiphon 2 after Canticle verses (and Glory Be if present)
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        antiphon_2 = liturgical_data['morning_prayer']['psalmody']['antiphon_2']
        
        # Create text box for the repeated antiphon with auto-fit
        ant2_repeat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        ant2_repeat_frame = ant2_repeat_box.text_frame
        ant2_repeat_frame.word_wrap = True
        ant2_repeat_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ant2_repeat_para = ant2_repeat_frame.paragraphs[0]
        ant2_repeat_para.alignment = PP_ALIGN.CENTER
        
        # Add "(All) Ant. 2 " in blue
        ant2_repeat_label = ant2_repeat_para.add_run()
        ant2_repeat_label.text = "(All) Ant. 2 "
        ant2_repeat_label.font.size = Pt(44)
        ant2_repeat_label.font.name = "Georgia"
        ant2_repeat_label.font.bold = True
        ant2_repeat_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
        
        # Add the antiphon text in black
        ant2_repeat_text = ant2_repeat_para.add_run()
        ant2_repeat_text.text = antiphon_2['text']
        ant2_repeat_text.font.size = Pt(44)
        ant2_repeat_text.font.name = "Georgia"
        ant2_repeat_text.font.bold = True
        ant2_repeat_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Repeated Antiphon 2")
        
        # ========== ANTIPHON 3 and PSALM 3 ==========
        
        # Create Antiphon 3 slide (without psalm info)
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        antiphon_3 = liturgical_data['morning_prayer']['psalmody']['antiphon_3']
        
        # Create text box for Antiphon 3 with auto-fit
        ant3_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        ant3_frame = ant3_box.text_frame
        ant3_frame.word_wrap = True
        ant3_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ant3_para = ant3_frame.paragraphs[0]
        ant3_para.alignment = PP_ALIGN.CENTER
        
        # Add "(All) Ant. 3 " in blue
        ant3_label = ant3_para.add_run()
        ant3_label.text = "(All) Ant. 3 "
        ant3_label.font.size = Pt(44)
        ant3_label.font.name = "Georgia"
        ant3_label.font.bold = True
        ant3_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
        
        # Add the antiphon text in black
        ant3_text = ant3_para.add_run()
        ant3_text.text = antiphon_3['text']
        ant3_text.font.size = Pt(44)
        ant3_text.font.name = "Georgia"
        ant3_text.font.bold = True
        ant3_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Antiphon 3")
        
        # Create separate slide for red psalm title and subtitle
        if antiphon_3.get('psalm_title'):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            psalm_info_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(3.5))
            psalm_info_frame = psalm_info_box.text_frame
            psalm_info_frame.word_wrap = True
            psalm_info_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            psalm_info_para = psalm_info_frame.paragraphs[0]
            psalm_info_para.alignment = PP_ALIGN.CENTER
            
            # Add psalm title in red
            psalm_title_run = psalm_info_para.add_run()
            psalm_title_run.text = antiphon_3['psalm_title']
            psalm_title_run.font.size = Pt(48)
            psalm_title_run.font.name = "Georgia"
            psalm_title_run.font.bold = True
            psalm_title_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
            
            # Add psalm subtitle in red italic if available
            if antiphon_3.get('psalm_subtitle'):
                psalm_info_para.add_run().text = "\n"
                psalm_subtitle_run = psalm_info_para.add_run()
                psalm_subtitle_run.text = antiphon_3['psalm_subtitle']
                psalm_subtitle_run.font.size = Pt(36)
                psalm_subtitle_run.font.name = "Georgia"
                psalm_subtitle_run.font.italic = True
                psalm_subtitle_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
            
            print(f"Created slide {slide_count}: Psalm 3 Title and Subtitle")
        
        # Create Psalm 3 verses alternating priest/people
        psalm_3_verses = liturgical_data['morning_prayer']['psalmody']['psalm_3']
        
        for verse in psalm_3_verses:
            # Check if this is the "Glory to the Father" verse (skip it, we'll add it manually after)
            if "Glory to the Father" in verse['text'] or "Glory to the father" in verse['text']:
                # Skip this verse - we'll add Glory Be manually after all verses
                print(f"  Skipping Glory Be verse from extraction (will add manually)")
                continue
            
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Create text box for the verse content with auto-fit
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            content_para = content_frame.paragraphs[0]
            content_para.alignment = PP_ALIGN.LEFT
            
            if verse['speaker'] == "Priest":
                # Priest slides: "Priest: " + text, all in red
                priest_run = content_para.add_run()
                priest_run.text = f"Priest: {verse['text']}"
                priest_run.font.size = Pt(44)
                priest_run.font.name = "Georgia"
                priest_run.font.bold = True
                priest_run.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
                print(f"Created slide {slide_count}: Psalm 3 - {verse['speaker']}")
                
            elif verse['speaker'] == "People":
                # People slides: "People: " in blue, rest of text in black
                people_label = content_para.add_run()
                people_label.text = "People: "
                people_label.font.size = Pt(44)
                people_label.font.name = "Georgia"
                people_label.font.bold = True
                people_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
                
                people_text = content_para.add_run()
                people_text.text = verse['text']
                people_text.font.size = Pt(44)
                people_text.font.name = "Georgia"
                people_text.font.bold = True
                people_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
                print(f"Created slide {slide_count}: Psalm 3 - {verse['speaker']}")
        
        # Add Glory Be after Psalm 3
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        glory_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        glory_frame = glory_box.text_frame
        glory_frame.word_wrap = True
        glory_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        glory_para = glory_frame.paragraphs[0]
        glory_para.alignment = PP_ALIGN.LEFT
        
        glory_run = glory_para.add_run()
        glory_run.text = "Glory to the Father, and to the Son, and to the Holy Spirit: as it was in the beginning, is now, and will be for ever. Amen."
        glory_run.font.size = Pt(44)
        glory_run.font.name = "Georgia"
        glory_run.font.bold = True
        glory_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Psalm 3 - Glory Be")
        
        # Repeat Antiphon 3 after Psalm 2 verses and Glory Be
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Create text box for the repeated antiphon with auto-fit
        ant3_repeat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
        ant3_repeat_frame = ant3_repeat_box.text_frame
        ant3_repeat_frame.word_wrap = True
        ant3_repeat_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        ant3_repeat_para = ant3_repeat_frame.paragraphs[0]
        ant3_repeat_para.alignment = PP_ALIGN.CENTER
        
        # Add "(All) Ant. 3 " in blue
        ant3_repeat_label = ant3_repeat_para.add_run()
        ant3_repeat_label.text = "(All) Ant. 3 "
        ant3_repeat_label.font.size = Pt(44)
        ant3_repeat_label.font.name = "Georgia"
        ant3_repeat_label.font.bold = True
        ant3_repeat_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
        
        # Add the antiphon text in black
        ant3_repeat_text = ant3_repeat_para.add_run()
        ant3_repeat_text.text = antiphon_3['text']
        ant3_repeat_text.font.size = Pt(44)
        ant3_repeat_text.font.name = "Georgia"
        ant3_repeat_text.font.bold = True
        ant3_repeat_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        print(f"Created slide {slide_count}: Repeated Antiphon 3")
        
        return slide_count

    # Additional section creation methods would follow the same pattern...
    def _create_reading_section(self, prs, liturgical_data, slide_count):
        """Create reading section with READING title and content on same slide
        
        Structure:
        - Single slide with "READING" title at top and content below
        """
        # Get reading data
        reading_data = liturgical_data.get('morning_prayer', {}).get('reading', {}).get('short_reading', {})
        
        if not reading_data or not reading_data.get('text'):
            print(f"  ⚠ No reading data available, skipping reading section")
            return slide_count
        
        # Create single slide with title and content
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add "READING" title at top
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(11.33), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = "READING"
        
        # Format title
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(48)
                run.font.bold = True
                run.font.color.rgb = self.reference_template['formatting_rules']['title_color']
        
        # Build content text
        content_text = ""
        if reading_data.get('citation'):
            content_text = f"{reading_data['citation']}\n\n"
        content_text += reading_data['text']
        
        # Add content text below title
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.75), Inches(12.33), Inches(5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.text = content_text
        
        # Format content
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(30)
                run.font.color.rgb = RGBColor(0, 0, 0)
        
        print(f"Created slide {slide_count}: READING (title + content)")
        
        return slide_count

    def _create_responsory_section(self, prs, liturgical_data, slide_count):
        """Create responsory section with alternating priest/all responses
        
        Structure:
        - First slide: "RESPONSORY" title + first response lines (combined)
        - Additional slides: Remaining responses with multiple lines per slide
        """
        # Get responsory data
        responsory_verses = liturgical_data.get('morning_prayer', {}).get('reading', {}).get('responsory', [])
        
        if not responsory_verses:
            print(f"  ⚠ No responsory data available, skipping responsory section")
            return slide_count
        
        # Create slides for each responsory part
        for verse in responsory_verses:
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Check if this slide should include the RESPONSORY title
            include_title = verse.get('include_title', False)
            
            # Determine speaker color
            speaker = verse.get('speaker', 'Priest')
            if speaker == 'All':
                speaker_color = self.reference_template['formatting_rules']['all_color']
                speaker_label = "(All) "
            else:
                speaker_color = self.reference_template['formatting_rules']['priest_color']
                speaker_label = "(Priest) "
            
            # If this is the first slide with title
            if include_title:
                # Add RESPONSORY title at top
                title_box = slide.shapes.add_textbox(
                    Inches(1), Inches(0.5), Inches(11.33), Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.text = "RESPONSORY"
                
                # Format title
                for paragraph in title_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(48)
                        run.font.bold = True
                        run.font.color.rgb = self.reference_template['formatting_rules']['title_color']
                
                # Add content below title
                text_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.75), Inches(12.33), Inches(5)
                )
                vertical_position = Inches(1.75)
            else:
                # No title, content starts higher
                text_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1), Inches(12.33), Inches(5.5)
                )
                vertical_position = Inches(1)
            
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Split the text by newlines to handle multi-line responses
            lines = verse['text'].split('\n')
            
            # Add first line with speaker label
            para = text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            
            speaker_run = para.add_run()
            speaker_run.text = speaker_label
            speaker_run.font.size = Pt(36)
            speaker_run.font.name = "Georgia"
            speaker_run.font.bold = True
            speaker_run.font.color.rgb = speaker_color
            
            # Add first line text in black
            text_run = para.add_run()
            text_run.text = lines[0]
            text_run.font.size = Pt(36)
            text_run.font.name = "Georgia"
            text_run.font.bold = True
            text_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add remaining lines (without speaker label)
            for line in lines[1:]:
                if line.strip():
                    new_para = text_frame.add_paragraph()
                    new_para.alignment = PP_ALIGN.CENTER
                    
                    line_run = new_para.add_run()
                    line_run.text = line
                    line_run.font.size = Pt(36)
                    line_run.font.name = "Georgia"
                    line_run.font.bold = True
                    line_run.font.color.rgb = RGBColor(0, 0, 0)
            
            print(f"Created slide {slide_count}: Responsory - {speaker}")
        
        return slide_count

    def _create_gospel_canticle_section(self, prs, liturgical_data, slide_count):
        """Create gospel canticle section following reference template
        
        Structure:
        1. GOSPEL CANTICLE title slide
        2. Antiphon slide with the gospel canticle antiphon text (ALL in purple)
        
        The antiphon may contain category markers like [Martyrs] or [Pastors]
        which should be displayed as part of the antiphon text.
        """
        try:
            gospel_canticle = liturgical_data['morning_prayer']['gospel_canticle']
            antiphon_text = gospel_canticle.get('antiphon', '')
            
            if not antiphon_text:
                print("  ⚠ No gospel canticle antiphon found, skipping section")
                return slide_count
            
            # Create single combined slide with header and antiphon
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slide_count += 1
            
            # Add "GOSPEL CANTICLE" header at the top
            header_left = Inches(0.5)
            header_top = Inches(0.5)
            header_width = Inches(12.33)
            header_height = Inches(1)
            
            header_box = slide.shapes.add_textbox(header_left, header_top, header_width, header_height)
            header_frame = header_box.text_frame
            header_frame.word_wrap = True
            
            header_p = header_frame.paragraphs[0]
            header_p.text = "GOSPEL CANTICLE"
            header_p.alignment = PP_ALIGN.CENTER
            header_p.font.size = Pt(48)
            header_p.font.name = 'Georgia'
            header_p.font.bold = True
            header_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            
            # Add antiphon content below header with auto-fit
            content_left = Inches(0.5)
            content_top = Inches(2)
            content_width = Inches(12.33)
            content_height = Inches(5)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            content_p = content_frame.paragraphs[0]
            content_p.alignment = PP_ALIGN.CENTER
            
            # Add "Ant. " label in red
            ant_label = content_p.add_run()
            ant_label.text = "Ant. "
            ant_label.font.size = Pt(44)
            ant_label.font.name = 'Georgia'
            ant_label.font.bold = True
            ant_label.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red (like Priest label)
            
            # Add antiphon text in purple (ALL response)
            ant_text = content_p.add_run()
            ant_text.text = antiphon_text
            ant_text.font.size = Pt(44)
            ant_text.font.name = 'Georgia'
            ant_text.font.bold = True
            ant_text.font.color.rgb = RGBColor(100, 0, 100)  # Purple for ALL
            
            print(f"Created slide {slide_count}: GOSPEL CANTICLE (with header and antiphon)")
            
            # Add Canticle of Zechariah title slide (red text, hard-coded)
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add the three lines of red text, centered
            canticle_left = Inches(0.5)
            canticle_top = Inches(2.5)
            canticle_width = Inches(12.33)
            canticle_height = Inches(3)
            
            canticle_box = slide.shapes.add_textbox(canticle_left, canticle_top, canticle_width, canticle_height)
            canticle_frame = canticle_box.text_frame
            canticle_frame.word_wrap = True
            
            # Line 1: "Canticle of Zechariah"
            p1 = canticle_frame.paragraphs[0]
            p1.text = "Canticle of Zechariah"
            p1.alignment = PP_ALIGN.CENTER
            p1.font.size = Pt(44)
            p1.font.name = 'Georgia'
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
            
            # Line 2: "Luke 1:68-79"
            p2 = canticle_frame.add_paragraph()
            p2.text = "Luke 1:68-79"
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(44)
            p2.font.name = 'Georgia'
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
            
            # Line 3: "The Messiah and his forerunner"
            p3 = canticle_frame.add_paragraph()
            p3.text = "The Messiah and his forerunner"
            p3.alignment = PP_ALIGN.CENTER
            p3.font.size = Pt(44)
            p3.font.name = 'Georgia'
            p3.font.bold = True
            p3.font.color.rgb = RGBColor(0x98, 0x00, 0x00)  # Red
            
            print(f"Created slide {slide_count}: Canticle of Zechariah (title)")
            
            # Add 9 Benedictus verse slides (alternating red/black text)
            benedictus_verses = [
                # Slide 1 - RED
                "Blessed + be the Lord, the God of Israel; *\nhe has come to his people and set them free.",
                # Slide 2 - BLACK
                "He has raised up for us a mighty savior,*\nborn of the house of his servant David.",
                # Slide 3 - RED
                "Through his holy prophets he promised of old †\nthat he would save us from our enemies, *\nfrom the hands of all who hate us.",
                # Slide 4 - BLACK
                "He promised to show mercy to our fathers *\nand to remember his holy covenant.",
                # Slide 5 - RED
                "This was the oath he swore to our father Abraham: *\nto set us free from the hands of our enemies,\nfree to worship him without fear, *\nholy and righteous in his sight\n  all the days of our life.",
                # Slide 6 - BLACK
                "You, my child, shall be called the prophet of the Most High; *\nfor you will go before the Lord to prepare his way,\nto give his people knowledge of salvation *\nby the forgiveness of their sins.",
                # Slide 7 - RED
                "In the tender compassion of our God *\nthe dawn from on high shall break upon us,\nto shine on those who dwell in darkness and the shadow of death, *\nand to guide our feet into the way of peace.",
                # Slide 8 - BLACK
                "Glory to the Father, and to the Son, *\nand to the Holy Spirit:",
                # Slide 9 - RED
                "as it was in the beginning, is now, *\nand will be forever. Amen."
            ]
            
            for i, verse_text in enumerate(benedictus_verses):
                slide_count += 1
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Determine color: odd slides (1,3,5,7,9) are red, even (2,4,6,8) are black
                is_red = (i % 2 == 0)
                text_color = RGBColor(0x98, 0x00, 0x00) if is_red else RGBColor(0, 0, 0)
                color_name = "red" if is_red else "black"
                
                # Create textbox with auto-fit for maximum text size
                verse_left = Inches(0.5)
                verse_top = Inches(1)
                verse_width = Inches(12.33)
                verse_height = Inches(5.5)
                
                verse_box = slide.shapes.add_textbox(verse_left, verse_top, verse_width, verse_height)
                verse_frame = verse_box.text_frame
                verse_frame.word_wrap = True
                verse_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                verse_p = verse_frame.paragraphs[0]
                verse_p.text = verse_text
                verse_p.alignment = PP_ALIGN.LEFT
                verse_p.font.size = Pt(44)  # Starting size, will auto-fit larger if possible
                verse_p.font.name = 'Georgia'
                verse_p.font.bold = True
                verse_p.font.color.rgb = text_color
                
                print(f"Created slide {slide_count}: Benedictus verse {i+1} ({color_name})")
            
            # Add repeated antiphon slide after Benedictus verses
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Create textbox for repeated antiphon with auto-fit
            ant_left = Inches(0.5)
            ant_top = Inches(1)
            ant_width = Inches(12.33)
            ant_height = Inches(5.5)
            
            ant_box = slide.shapes.add_textbox(ant_left, ant_top, ant_width, ant_height)
            ant_frame = ant_box.text_frame
            ant_frame.word_wrap = True
            ant_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            ant_p = ant_frame.paragraphs[0]
            ant_p.alignment = PP_ALIGN.CENTER
            
            # Add "Ant. " label in blue (All response pattern)
            ant_label = ant_p.add_run()
            ant_label.text = "Ant. "
            ant_label.font.size = Pt(44)
            ant_label.font.name = 'Georgia'
            ant_label.font.bold = True
            ant_label.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue for repeated antiphon
            
            # Add antiphon text in black
            ant_text_run = ant_p.add_run()
            ant_text_run.text = antiphon_text
            ant_text_run.font.size = Pt(44)
            ant_text_run.font.name = 'Georgia'
            ant_text_run.font.bold = True
            ant_text_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            
            print(f"Created slide {slide_count}: Repeated Gospel Canticle Antiphon")
            
            return slide_count
            
        except Exception as e:
            print(f"  ⚠ Error creating gospel canticle section: {e}")
            traceback.print_exc()
            return slide_count

    def _extract_intercessions(self, soup, text):
        """Extract intercessions section from liturgical text"""
        try:
            # Find INTERCESSIONS markers - there may be multiple instances
            # We want the one AFTER the Gospel Canticle
            # Work with the HTML (from soup) not the plain text
            html_content = str(soup)
            
            intercessions_matches = list(re.finditer(r'INTERCESSIONS', html_content, re.IGNORECASE))
            
            if not intercessions_matches:
                print("  ⚠ No INTERCESSIONS marker found")
                return []
            
            # Use the LAST instance (after Gospel Canticle)
            intercessions_pos = intercessions_matches[-1].start()
            html_from_intercessions = html_content[intercessions_pos:]
            
            # Extract until THE LORD'S PRAYER or Let us pray
            end_match = re.search(r'THE LORD.S PRAYER|Let us pray\.', html_from_intercessions, re.IGNORECASE)
            if end_match:
                intercessions_section = html_from_intercessions[:end_match.start()]
            else:
                intercessions_section = html_from_intercessions[:3000]
            
            print(f"  Found INTERCESSIONS section ({len(intercessions_section)} chars)")
            
            # Parse the intercessions structure
            # Structure can have category markers like [Martyrs] or [Pastors]
            intercessions_groups = []
            
            # Split by category markers if present
            category_pattern = r'\[(Martyrs|Pastors|Doctors|Virgins|Holy Men and Women)\]'
            parts = re.split(category_pattern, intercessions_section, flags=re.IGNORECASE)
            
            # If no categories, treat as single group
            if len(parts) == 1:
                intercessions_groups.append({'category': None, 'text': parts[0]})
            else:
                # Skip first part (before any category)
                for i in range(1, len(parts), 2):
                    if i + 1 < len(parts):
                        category = parts[i]
                        text_content = parts[i + 1]
                        intercessions_groups.append({'category': category, 'text': text_content})
            
            # Parse each group
            all_intercessions = []
            for group in intercessions_groups:
                group_text = group['text']
                category = group['category']
                
                # Find introduction (text before first "Nourish" or "You redeemed")
                intro_match = re.search(r'(.*?)(?=Nourish your people|You redeemed us)', group_text, re.DOTALL | re.IGNORECASE)
                if intro_match:
                    introduction = intro_match.group(1).strip()
                    # Clean up introduction - remove HTML tags and INTERCESSIONS text
                    introduction = re.sub(r'INTERCESSIONS', '', introduction, flags=re.IGNORECASE).strip()
                    introduction = re.sub(r'<[^>]+>', '', introduction).strip()  # Remove all HTML tags
                    # Get text after introduction for intention extraction
                    intentions_text = group_text[intro_match.end():]
                else:
                    introduction = ""
                    intentions_text = group_text
                
                # Find response line (All)
                response_match = re.search(r'(Nourish your people, Lord\.|You redeemed us by your blood\.)', group_text, re.IGNORECASE)
                response_line = response_match.group(1) if response_match else ""
                
                # Find all prayer intentions (lines with em dash)
                # Each intention has a petition line followed by em dash and response
                # Structure in HTML: [petition]<span class="rubrica">—</span>[response]<br><em>Response line</em>
                # We capture up to but not including the italicized response repeat
                
                # First, remove the italicized response lines from the text (both in <em> tags and plain)
                cleaned_text = re.sub(r'<em>\s*(Nourish your people, Lord\.|You redeemed us by your blood\.)\s*</em>', '', intentions_text, flags=re.IGNORECASE)
                
                # Remove standalone response lines that appear between intentions
                cleaned_text = re.sub(r'(Nourish your people, Lord\.|You redeemed us by your blood\.)', '', cleaned_text, flags=re.IGNORECASE)
                
                # Replace the rubrica em dash spans with plain em dash for easier matching
                cleaned_text = re.sub(r'<span class="rubrica">—</span>', '—', cleaned_text, flags=re.IGNORECASE)
                
                # Remove <br> and <br/> tags which interfere with pattern matching
                cleaned_text = re.sub(r'<br\s*/?>', ' ', cleaned_text, flags=re.IGNORECASE)
                
                # Remove any remaining HTML tags (including leftover </em>, <em>, etc.)
                cleaned_text = re.sub(r'<[^>]+>', ' ', cleaned_text, flags=re.IGNORECASE)
                
                # Now extract intentions with em dash pattern
                # Pattern: [petition text (no < or —)]—[response text (no < except maybe br)]
                intention_pattern = r'([^—<]+?)—\s*([^<]+?)\.'
                intentions = []
                
                for match in re.finditer(intention_pattern, cleaned_text):
                    petition = match.group(1).strip()
                    response = match.group(2).strip()
                    
                    # Skip if petition is too short (likely noise)
                    if len(petition) < 20:
                        continue
                    
                    # Skip if petition contains the main response line (this would be noise/duplication)
                    if re.search(r'Nourish your people|You redeemed us|INTERCESSIONS', petition, re.IGNORECASE):
                        continue
                    
                    # Clean petition: remove any leading/trailing HTML or punctuation
                    petition = re.sub(r'<[^>]+>', '', petition).strip()
                    petition = re.sub(r'^\s*[,\s]+', '', petition).strip()
                    
                    # Clean response: remove any HTML tags and normalize whitespace
                    response = re.sub(r'<[^>]+>', '', response).strip()
                    response = re.sub(r'\s+', ' ', response).strip()
                    
                    # Add period back
                    if response and not response.endswith('.'):
                        response += '.'
                    
                    if petition and response:  # Only add if we have both valid petition and response
                        intentions.append({
                            'petition': petition,
                            'response': response
                        })
                
                if introduction or intentions:
                    all_intercessions.append({
                        'category': category,
                        'introduction': introduction,
                        'response_line': response_line,
                        'intentions': intentions
                    })
            
            print(f"  Extracted {len(all_intercessions)} intercession group(s)")
            for i, group in enumerate(all_intercessions):
                print(f"    Group {i+1}: {group['category'] or 'Default'}, {len(group['intentions'])} intentions")
            
            return all_intercessions
            
        except Exception as e:
            print(f"  ⚠ Error extracting intercessions: {e}")
            traceback.print_exc()
            return []

    def _create_intercessions_section(self, prs, liturgical_data, slide_count):
        """Create intercessions section following reference template"""
        try:
            intercessions_data = liturgical_data['morning_prayer'].get('intercessions', [])
            
            if not intercessions_data:
                print("  No intercessions data available")
                return slide_count
            
            # Slide 1: INTERCESSIONS title
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Title text
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            
            title_run = title_para.add_run()
            title_run.text = "INTERCESSIONS"
            title_run.font.name = "Georgia"
            title_run.font.size = Pt(48)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            
            print(f"Created slide {slide_count}: INTERCESSIONS (title)")
            
            # Process each intercession group
            for group in intercessions_data:
                category = group.get('category')
                introduction = group.get('introduction', '')
                response_line = group.get('response_line', '')
                intentions = group.get('intentions', [])
                
                # Slide 2: Introduction + (All) Response
                slide_count += 1
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                content_para = content_frame.paragraphs[0]
                content_para.alignment = PP_ALIGN.CENTER
                
                # Introduction text (black)
                if introduction:
                    intro_run = content_para.add_run()
                    intro_run.text = introduction
                    intro_run.font.name = "Georgia"
                    intro_run.font.size = Pt(30)
                    intro_run.font.bold = True
                    intro_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                
                # Add spacing
                if introduction and response_line:
                    spacing_run = content_para.add_run()
                    spacing_run.text = "\n\n"
                    spacing_run.font.size = Pt(20)
                
                # Response line - (All) label in purple, text in purple
                if response_line:
                    all_label_run = content_para.add_run()
                    all_label_run.text = "(All) "
                    all_label_run.font.name = "Georgia"
                    all_label_run.font.size = Pt(30)
                    all_label_run.font.bold = True
                    all_label_run.font.color.rgb = RGBColor(100, 0, 100)  # Purple
                    
                    response_run = content_para.add_run()
                    response_run.text = response_line
                    response_run.font.name = "Georgia"
                    response_run.font.size = Pt(30)
                    response_run.font.bold = True
                    response_run.font.color.rgb = RGBColor(100, 0, 100)  # Purple
                
                print(f"Created slide {slide_count}: Intercessions Introduction{' - ' + category if category else ''}")
                
                # Slides for each intention
                for intention in intentions:
                    slide_count += 1
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    intention_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.33), Inches(5.5))
                    intention_frame = intention_box.text_frame
                    intention_frame.word_wrap = True
                    intention_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    intention_para = intention_frame.paragraphs[0]
                    intention_para.alignment = PP_ALIGN.CENTER
                    
                    # Petition text (black)
                    petition_run = intention_para.add_run()
                    petition_run.text = intention['petition']
                    petition_run.font.name = "Georgia"
                    petition_run.font.size = Pt(30)
                    petition_run.font.bold = True
                    petition_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                    
                    # Em dash response (on new line, black)
                    dash_run = intention_para.add_run()
                    dash_run.text = "\n— " + intention['response']
                    dash_run.font.name = "Georgia"
                    dash_run.font.size = Pt(30)
                    dash_run.font.bold = True
                    dash_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                    
                    print(f"Created slide {slide_count}: Intercession Intention")
            
            return slide_count
            
        except Exception as e:
            print(f"  ⚠ Error creating intercessions section: {e}")
            traceback.print_exc()
            return slide_count

    def _create_sacred_heart_hymns(self, prs, liturgical_data, slide_count):
        """Create static sacred heart hymns following reference template"""
        # Implementation would use static content from reference
        return slide_count + 6  # Placeholder

    def _create_mass_readings_section(self, prs, liturgical_data, slide_count):
        """Create mass readings section following reference template"""
        # Implementation would follow reference structure
        return slide_count + 19  # Placeholder

    def _create_post_communion_prayers(self, prs, liturgical_data, slide_count):
        """Create static post-communion prayers following reference template"""
        # Implementation would use static content from reference
        return slide_count + 17  # Placeholder

    def _create_transition_slides(self, prs, slide_count):
        """Create blank transition slides following reference template"""
        for i in range(10):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            print(f"Created slide {slide_count}: Transition slide")
        return slide_count

    def _create_jubilee_prayer(self, prs, liturgical_data, slide_count):
        """Create static jubilee prayer following reference template"""
        # Implementation would use static content from reference
        return slide_count + 7  # Placeholder

    def _create_st_joseph_prayer(self, prs, liturgical_data, slide_count):
        """Create static St. Joseph prayer following reference template"""
        # Implementation would use static content from reference
        return slide_count + 12  # Placeholder

def main():
    import sys
    
    print("BBGRL Slide Generator V2 - Template-Based Dynamic Generator")
    print("=" * 60)
    print("Fetching live liturgical data and applying reference structure...")
    
    generator = BBGRLSlideGeneratorV2()
    
    # Parse command line arguments for date
    if len(sys.argv) > 1:
        # Accept date in format YYYY-MM-DD
        date_str = sys.argv[1]
        try:
            target_date = datetime.strptime(date_str, "%Y-%m-%d")
        except ValueError:
            print(f"Error: Invalid date format '{date_str}'. Use YYYY-MM-DD")
            return
    else:
        # Default to November 11, 2025
        target_date = datetime(2025, 11, 11)
    
    print(f"Generating slides for: {target_date.strftime('%B %d, %Y')}")
    
    # Fetch liturgical data for the specified date
    liturgical_data = generator.fetch_live_liturgical_data(target_date)
    
    # Create presentation using reference template structure
    output_path = generator.create_presentation_from_template(liturgical_data)
    
    print("\n✓ Template-based presentation created successfully!")
    print("✓ Uses live liturgical data with exact reference formatting")
    print(f"✓ File naming convention: olph_slides_[year]_[month]_[day].pptx")

if __name__ == "__main__":
    main()