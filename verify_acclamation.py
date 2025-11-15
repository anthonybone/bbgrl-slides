"""Verify Gospel Acclamation slides content"""

from pptx import Presentation

def verify_acclamation_slides(filename):
    """Check the Gospel Acclamation slides in a presentation"""
    print(f"\n{'='*70}")
    print(f"Verifying: {filename}")
    print(f"{'='*70}\n")
    
    prs = Presentation(filename)
    
    # Look for Acclamation slides
    found_header = False
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Check for header slide
                if "Acclamation before the Gospel" in shape.text:
                    print(f"Slide {idx} (HEADER):")
                    print(f"{shape.text}")
                    print("-" * 70)
                    found_header = True
                # Check for verse slide (next slide after header)
                elif found_header and ("Come after me" in shape.text or "I give you" in shape.text or "Remain in my love" in shape.text or "Alleluia" in shape.text):
                    print(f"Slide {idx} (VERSE):")
                    print(f"{shape.text}")
                    print("-" * 70)
                    found_header = False

# Check all three dates
dates = [
    "output_v2/olph_slides_2025_11_10.pptx",
    "output_v2/olph_slides_2025_11_11.pptx",
    "output_v2/olph_slides_2025_11_12.pptx"
]

for date_file in dates:
    try:
        verify_acclamation_slides(date_file)
    except Exception as e:
        print(f"Error checking {date_file}: {e}")
