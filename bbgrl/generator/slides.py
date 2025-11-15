from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR


def create_initial_blank_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    print(f"Created slide {slide_count}: Initial blank black slide")
    return slide_count


def create_daily_morning_prayer_image_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = "png/daily_morning_prayer.png"
    if hasattr(prs, "slide_width"):
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    else:
        slide_width = None
        slide_height = None
    import os
    if os.path.exists(image_path) and slide_width is not None:
        slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
        print(f"Created slide {slide_count}: Daily Morning Prayer image slide")
    else:
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "Daily Morning Prayer"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(184, 134, 11)
        title_para.alignment = PP_ALIGN.CENTER
        print(f"Created slide {slide_count}: Daily Morning Prayer text slide (image not found)")
    return slide_count


def create_heart_of_jesus_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = "png/heart_of_jesus_slide.png"
    import os
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        print(f"Created slide {slide_count}: Heart of Jesus image slide")
    else:
        print(f"  WARNING: Heart of Jesus image not found at {image_path}")
    return slide_count


def create_heart_of_jesus_prayer_slides(prs, slide_count):
    try:
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Heart of Jesus meek and mild.\n"
            "Hear oh hear thy feeble child\n"
            "When the tempest's most severe. Heart of Jesus, hear.\n"
            "Sweetly, we'll rest on thy Sacred Heart. Never from Thee. O let us part!"
        )
        run = para.add_run()
        run.text = prayer_text
        run.font.name = "Georgia"
        run.font.size = Pt(44)
        run.font.color.rgb = RGBColor(0, 0, 0)
        print(f"Created slide {slide_count}: Heart of Jesus Prayer (1/2)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(3.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Hear then, Thy loving children's pray'r\n"
            "O Heart of Jesus,\n"
            "Heart of Jesus hear."
        )
        run = para.add_run()
        run.text = prayer_text
        run.font.name = "Georgia"
        run.font.size = Pt(44)
        run.font.color.rgb = RGBColor(0, 0, 0)
        print(f"Created slide {slide_count}: Heart of Jesus Prayer (2/2)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Heart of Jesus prayer slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_oh_sacred_heart_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = "png/oh_sacred_heart.png"
    import os
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        print(f"Created slide {slide_count}: Oh Sacred Heart image slide")
    else:
        print(f"  WARNING: Oh Sacred Heart image not found at {image_path}")
    return slide_count


def create_oh_sacred_heart_prayer_slides(prs, slide_count):
    try:
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Oh Sacred Heart, Oh love divine.\n"
            "Do keep us near to Thee.\n"
            "And make our love so like to Thine\n"
            "That we may holy be.\n"
            "Heart of Jesus hear.\n"
            "Oh heart of love divine.\n"
            "Listen to our pray'r."
        )
        run = para.add_run(); run.text = prayer_text; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Oh Sacred Heart Prayer (1/2)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Make us always Thine.\n"
            "Oh temple pure, Oh house of gold. Our heaven here below.\n"
            "What gifts unfurled, what wealth untold. From Thee do ever flow.\n"
            "Heart of Jesus hear. Oh Heart of love divine. Listen to our pray'r. Make us always Thine."
        )
        run = para.add_run(); run.text = prayer_text; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Oh Sacred Heart Prayer (2/2)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Oh Sacred Heart prayer slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_novena_sacred_heart_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = "png/novena_sacred_heart.png"
    import os
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        print(f"Created slide {slide_count}: Novena to the Sacred Heart image slide")
    else:
        print(f"  WARNING: Novena Sacred Heart image not found at {image_path}")
    return slide_count


def create_soul_of_christ_slides(prs, slide_count):
    try:
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Soul of Christ, make me holy.\n"
            "Body of Christ, save me.\n"
            "Blood of Christ, inebriate me.\n"
            "Water from the side of Christ, wash me.\n"
            "Passion of Christ, make me strong.\n"
            "O good Jesus, hear me.\n"
            "Hide me within your wounds."
        )
        run = para.add_run(); run.text = prayer_text; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Soul of Christ Prayer (1/2)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        prayer_text = (
            "Let me never be separated from You.\n"
            "Deliver me from the wicked enemy,\n"
            "Call me at the hour of my death.\n"
            "And tell me to come to you\n"
            "that with Your saints I may praise You forever.\n"
            "Amen."
        )
        run = para.add_run(); run.text = prayer_text; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Soul of Christ Prayer (2/2)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Soul of Christ prayer slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_prayer_of_thanksgiving_slides(prs, slide_count):
    try:
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(0)
        run = para.add_run(); run.text = "We carry out the Lord's command, instructed by his divine teaching, we dare to say:"; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Prayer of Thanksgiving (1/4)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER; para.space_before = Pt(0)
        run = para.add_run(); run.text = "WITH PROFOUND ADORATION OF YOUR DIVINE MAJESTY, WE HUMBLE OURSELVES BEFORE YOU AND THANK YOU MOST HEARTILY FOR THE GOODNESS YOU HAVE SHOWN US.  FATHER, WE GIVE YOU THANKS FOR THE GREAT LOVE YOU HAVE GIVEN US IN THE SACRED HEART OF JESUS, YOUR BELOVED SON. THROUGH HIM, WE OFFER YOU OUR  AND FRUSTRATIONS, OUR JOYS AND OUR SORROWS"; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Prayer of Thanksgiving (2/4)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame; text_frame.word_wrap = True; text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER; para.space_before = Pt(0)
        run = para.add_run(); run.text = "I wish to thank You, Lord Jesus, on behalf of myself and of all creatures, and to make amends to You, as far as I am able, for the ingratitude of many, of which You complain so vehemently."; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Prayer of Thanksgiving (3/4)")
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
        text_frame = text_box.text_frame; text_frame.word_wrap = True; text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER; para.space_before = Pt(0)
        run = para.add_run(); run.text = "I wish that I was able to direct the hearts and minds of all to You and, together with them and for them, love You perfectly in return just as You rightly expect."; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.color.rgb = RGBColor(0,0,0)
        print(f"Created slide {slide_count}: Prayer of Thanksgiving (4/4)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Prayer of Thanksgiving slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_novena_of_confidence_slides(prs, slide_count):
    try:
        slides = [
            {"header": "NOVENA OF CONFIDENCE:", "lines": [
                "O Lord, Jesus Christ, to Your Most Sacred Heart I confide this intention.....",
                "(Mention your request)"
            ]},
            {"lines": [
                "Only look upon me, then do what Your heart inspires. Let Your Sacred Heart decide. I count on You. I trust in You. I throw myself on Your mercy."
            ]},
            {"lines": [
                "Lord Jesus! You will not fail me. Sacred Heart of Jesus, I believe in Your love for me. O Sacred Heart Of Jesus, Your kingdom come."
            ]},
            {"lines": [
                "O Sacred Heart of Jesus, I have asked for many favors, but I earnestly implore this one. Take it; place it in Your Sacred Heart. When the Eternal Father sees it covered with Your Precious Blood,"
            ]},
            {"lines": [
                "He will not refuse it. It will be no longer my prayer but Yours, O Jesus. O Sacred Heart of Jesus, I place my trust in you. Let me never be confounded. Amen."
            ]}
        ]
        for idx, data in enumerate(slides, start=1):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
            frame = box.text_frame
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.TOP
            if data.get("header"):
                header_para = frame.paragraphs[0]
                header_para.alignment = PP_ALIGN.CENTER
                run = header_para.add_run(); run.text = data["header"]; run.font.name = "Georgia"; run.font.size = Pt(54); run.font.bold = True; run.font.color.rgb = RGBColor(0x98,0x00,0x00)
                spacer = frame.add_paragraph(); spacer.text = ""; spacer.alignment = PP_ALIGN.CENTER
            else:
                first_para = frame.paragraphs[0]; first_para.alignment = PP_ALIGN.CENTER
                run = first_para.add_run(); run.text = data["lines"][0]; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = RGBColor(0,0,0)
                for line in data["lines"][1:]:
                    p = frame.add_paragraph(); p.alignment = PP_ALIGN.CENTER
                    run = p.add_run(); run.text = line; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = RGBColor(0,0,0)
            if data.get("header"):
                for line in data["lines"]:
                    p = frame.add_paragraph(); p.alignment = PP_ALIGN.CENTER
                    run = p.add_run(); run.text = line; run.font.name = "Georgia"; run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = RGBColor(0,0,0)
            print(f"Created slide {slide_count}: Novena of Confidence ({idx}/5)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Novena of Confidence slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_novena_prayer_slides(prs, slide_count):
    try:
        slides = [
            {"header": "NOVENA PRAYER:", "lines": [
                "O most holy Heart of Jesus, fountain of every blessing, I adore You, I love You, and with a lively sorrow for my sins,"
            ]},
            {"lines": [
                "I offer You this poor heart of mine. Make me humble, patient, pure and wholly obedient to Your will. Grant, good Jesus that I may live in You and for You."
            ]},
            {"lines": [
                "Protect me in the midst of danger; comfort me in my afflictions; give me health of body, assistance in my temporal needs, Your blessing on all that I do, and the grace of a holy death.......,",
                "Our Father…. Hail Mary.... Glory be to the Father....."
            ]},
            {"lines": [
                "(Priest) Let us Pray,",
                "Heavenly Father, we rejoice in the gifts of love we have received from the Heart of Jesus, your Son."
            ]},
            {"lines": [
                "Open our hearts to share His life and continue to bless us with His love. We ask this in the name of Jesus the Lord.",
                "AMEN."
            ]}
        ]
        def style_line(paragraph, text):
            paragraph.alignment = PP_ALIGN.CENTER
            yellow = RGBColor(255, 215, 0)
            if text.startswith("Priest:"):
                label = "Priest:"; rest = text[len(label):].lstrip()
                r1 = paragraph.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(44); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = RGBColor(0x98,0x00,0x00)
                r2 = paragraph.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(44); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
            elif text.startswith("All:"):
                label = "All:"; rest = text[len(label):].lstrip()
                r1 = paragraph.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(44); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = yellow
                r2 = paragraph.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(44); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
            elif text.startswith("(Priest)"):
                after = text[len("(Priest)"):]
                r_open = paragraph.add_run(); r_open.text = "("; r_open.font.name = "Georgia"; r_open.font.size = Pt(44); r_open.font.bold = True; r_open.font.color.rgb = RGBColor(0,0,0)
                r_mid = paragraph.add_run(); r_mid.text = "Priest"; r_mid.font.name = "Georgia"; r_mid.font.size = Pt(44); r_mid.font.bold = True; r_mid.font.underline = True; r_mid.font.color.rgb = RGBColor(0x98,0x00,0x00)
                r_close = paragraph.add_run(); r_close.text = ") "; r_close.font.name = "Georgia"; r_close.font.size = Pt(44); r_close.font.bold = True; r_close.font.color.rgb = RGBColor(0,0,0)
                r_rest = paragraph.add_run(); r_rest.text = after.lstrip(); r_rest.font.name = "Georgia"; r_rest.font.size = Pt(44); r_rest.font.bold = True; r_rest.font.color.rgb = RGBColor(0,0,0)
            else:
                r = paragraph.add_run(); r.text = text; r.font.name = "Georgia"; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = RGBColor(0,0,0)
        for idx, data in enumerate(slides, start=1):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
            frame = box.text_frame; frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.TOP; frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if data.get("header"):
                header_para = frame.paragraphs[0]; header_para.alignment = PP_ALIGN.CENTER
                run = header_para.add_run(); run.text = data["header"]; run.font.name = "Georgia"; run.font.size = Pt(54); run.font.bold = True; run.font.color.rgb = RGBColor(0x98,0x00,0x00)
                spacer = frame.add_paragraph(); spacer.text = ""; spacer.alignment = PP_ALIGN.CENTER
            else:
                first_para = frame.paragraphs[0]
                style_line(first_para, data["lines"][0] if data["lines"] else "")
                for line in data["lines"][1:]:
                    p = frame.add_paragraph(); style_line(p, line)
            if data.get("header"):
                for line in data["lines"]:
                    p = frame.add_paragraph(); style_line(p, line)
            print(f"Created slide {slide_count}: Novena Prayer ({idx}/5)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Novena Prayer slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_salve_regina_slides(prs, slide_count):
    try:
        slides = [
            {"header_line": "Salve Regina", "lines": ["Salve, Regina, Mater misericordiae,", "vita, dulcedo, et spes nostra, salve."]},
            {"lines": ["ad te clamamus", "exsules filii Evae,", "ad te suspiramus, gementes et flentes", "in hac lacrimarum valle."]},
            {"lines": ["Eia, ergo, advocata nostra, illos tuos", "misericordes oculos ad nos converte;", "et Iesum"]},
            {"lines": ["benedictum fructum ventris tui,", "nobis post hoc exsilium ostende.", "O clemens, O pia, O dulcis Virgo Maria."]},
            {"lines": ["Priest: Pray for us O holy Mother of God, ", "", "All: that we may be made worthy of the promises of Christ.", "", "Priest: Let us pray"]},
            {"lines": ["Grant, Lord God, that we Your servants may rejoice in unfailing health of mind and body and, through the glorious intercession of Blessed Mary ever virgin, may we be set free from present sorrow and come to enjoy eternal happiness"]},
            {"lines": ["Through our Lord Jesus Christ Your Son, who lives and reigns with You in the unity of the Holy Spirit, one God, forever and ever.", "", "All: AMEN."]},
        ]
        red = RGBColor(0x98, 0x00, 0x00); yellow = RGBColor(255, 215, 0)
        for idx, data in enumerate(slides, start=1):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
            frame = box.text_frame; frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.TOP; frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if data.get("header_line"):
                p = frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = data["header_line"]; r.font.name = "Georgia"; r.font.size = Pt(60); r.font.bold = True; r.font.underline = True; r.font.color.rgb = red
                spacer = frame.add_paragraph(); spacer.text = ""; spacer.alignment = PP_ALIGN.CENTER
                for line in data["lines"]:
                    p2 = frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
                    if line.startswith("Priest:"):
                        label = "Priest:"; rest = line[len(label):].lstrip()
                        r1 = p2.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = red
                        r2 = p2.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    elif line.startswith("All:"):
                        label = "All:"; rest = line[len(label):].lstrip()
                        r1 = p2.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = yellow
                        r2 = p2.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    else:
                        r2 = p2.add_run(); r2.text = line; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
            else:
                if data["lines"]:
                    p0 = frame.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
                    first_line = data["lines"][0]
                    if first_line.startswith("Priest:"):
                        label = "Priest:"; rest = first_line[len(label):].lstrip()
                        r1 = p0.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = red
                        r2 = p0.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    elif first_line.startswith("All:"):
                        label = "All:"; rest = first_line[len(label):].lstrip()
                        r1 = p0.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = yellow
                        r2 = p0.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    else:
                        r0 = p0.add_run(); r0.text = first_line; r0.font.name = "Georgia"; r0.font.size = Pt(52); r0.font.bold = True; r0.font.color.rgb = RGBColor(0,0,0)
                for line in data["lines"][1:]:
                    p = frame.add_paragraph(); p.alignment = PP_ALIGN.CENTER
                    if line.startswith("Priest:"):
                        label = "Priest:"; rest = line[len(label):].lstrip()
                        r1 = p.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = red
                        r2 = p.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    elif line.startswith("All:"):
                        label = "All:"; rest = line[len(label):].lstrip()
                        r1 = p.add_run(); r1.text = label + " "; r1.font.name = "Georgia"; r1.font.size = Pt(52); r1.font.bold = True; r1.font.underline = True; r1.font.color.rgb = yellow
                        r2 = p.add_run(); r2.text = rest; r2.font.name = "Georgia"; r2.font.size = Pt(52); r2.font.bold = True; r2.font.color.rgb = RGBColor(0,0,0)
                    else:
                        r = p.add_run(); r.text = line; r.font.name = "Georgia"; r.font.size = Pt(52); r.font.bold = True; r.font.color.rgb = RGBColor(0,0,0)
            print(f"Created slide {slide_count}: Salve Regina ({idx}/7)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Salve Regina slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_prayer_to_st_michael_slides(prs, slide_count):
    try:
        slides = [["Saint Michael the Archangel, defend us in battle. Be our protection against the wickedness and snares of the devil."],["May God rebuke him we humbly pray; and do thou. O Prince of the heavenly host, by the power of God,"],["Cast into Hell Satan and all the evil spirits who prowl about the world seeking the ruin of souls. Amen."]]
        for idx, lines in enumerate(slides, start=1):
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(11.93), Inches(6.1))
            frame = box.text_frame; frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.TOP; frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            header_para = frame.paragraphs[0]; header_para.alignment = PP_ALIGN.CENTER
            header_run = header_para.add_run(); header_run.text = "Prayer to St. Michael"; header_run.font.name = "Georgia"; header_run.font.size = Pt(54); header_run.font.bold = True; header_run.font.color.rgb = RGBColor(0x98,0x00,0x00)
            spacer = frame.add_paragraph(); spacer.text = ""; spacer.alignment = PP_ALIGN.CENTER
            for text in lines:
                p = frame.add_paragraph(); p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = text; r.font.name = "Georgia"; r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = RGBColor(0,0,0)
            print(f"Created slide {slide_count}: Prayer to St. Michael ({idx}/3)")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Prayer to St. Michael slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_jubilee_prayer_slides(prs, slide_count):
    try:
        burgundy = RGBColor(139, 0, 0)
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.33), Inches(3.0))
        title_frame = title_box.text_frame; title_frame.word_wrap = True; title_frame.text = "THE JUBILEE PRAYER"
        title_para = title_frame.paragraphs[0]; title_para.alignment = PP_ALIGN.CENTER; title_para.font.name = "Georgia"; title_para.font.bold = True; title_para.font.size = Pt(80); title_para.font.color.rgb = burgundy
        print(f"Created slide {slide_count}: THE JUBILEE PRAYER (title)")
        def add_body_slide(lines):
            nonlocal slide_count
            slide_count += 1
            s = prs.slides.add_slide(prs.slide_layouts[6])
            box = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.8))
            tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
            for idx, ln in enumerate(lines):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = ln; p.alignment = PP_ALIGN.CENTER; p.font.name = "Georgia"; p.font.bold = True; p.font.size = Pt(48); p.font.color.rgb = RGBColor(0,0,0)
            print(f"Created slide {slide_count}: Jubilee Prayer body")
        add_body_slide(["Father in heaven,", "may the faith you have given us in your son, Jesus Christ, our brother, and the flame of charity"])
        add_body_slide(["enkindled in our hearts by the Holy Spirit, reawaken in us the blessed hope for the coming of your Kingdom."])
        add_body_slide(["May your grace transform us", "into tireless cultivators of the seeds of the Gospel. May those seeds transform from within both humanity and the whole"])
        add_body_slide(["world into your own family. ", "By the grace of the Jubilee Year may we be empowered to bring your love to the world, both near and far,"])
        add_body_slide(["so that by your Spirit we may witness in faith and hope to the great work of your love that began in your Son. Amen."])
        slide_count += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.8))
        tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.clear()
        p = tf.paragraphs[0]; p.text = "Amen."; p.alignment = PP_ALIGN.CENTER; p.font.name = "Georgia"; p.font.bold = True; p.font.size = Pt(48); p.font.color.rgb = burgundy
        print(f"Created slide {slide_count}: Jubilee Prayer closing with Amen.")
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating Jubilee Prayer slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_st_joseph_prayer_image_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = "png/st_joseph_prayer.png"
    import os
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        print(f"Created slide {slide_count}: St. Joseph Prayer image slide")
    else:
        print(f"  WARNING: St. Joseph Prayer image not found at {image_path}")
    return slide_count


def create_st_joseph_prayer_text_slides(prs, slide_count):
    try:
        burgundy = RGBColor(139, 0, 0)
        slides = [
            ["To you, O blessed Joseph, do we come in our tribulation, and having implored the help of your most holy Spouse, we confidently invoke your patronage also."],
            ["Through that charity which bound you to the Immaculate Virgin Mother of God and through the"],
            ["paternal love with which you embraced the Child Jesus, we humbly beg you graciously to regard the inheritance"],
            ["which Jesus Christ has purchased by his Blood, and with your power and strength to aid us in our necessities"],
            ["O most watchful guardian of the Holy Family, defend the chosen children of Jesus Christ;"],
            ["O most loving father, ward off from us every contagion of error and corrupting influence;"],
            ["O our most mighty protector, be kind to us and from heaven assist us in our struggle with the power of darkness"],
            ["As once you rescued the Child Jesus from deadly peril, so now protect God's Holy Church from the snares of the"],
            ["enemy and from all adversity; shield, too, each one of us by your constant protection,"],
            ["so that, supported by your example and your aid, we may be able to live piously, to die in holiness, and to obtain eternal happiness in heaven. Amen."],
            ["St. Joseph, husband of the Blessed Virgin Mary", "All praise, O God, for Joseph, The guardian of your Son, Who saved him from King Herod, When safety there was none."],
            ["He taught the trade of builder, When they to Naz'reth came, And Joseph's love made \"Father\" To be, for Christ, God's name."]
        ]
        phrase = "St. Joseph, husband of the Blessed Virgin Mary"
        import builtins
        def add_slide(lines):
            nonlocal slide_count
            slide_count += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(11.93), Inches(6.1))
            tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
            for idx, line in enumerate(lines):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if phrase in line:
                    before, mid, after = line.partition(phrase)
                    if before:
                        r = p.add_run(); r.text = before; r.font.name = "Georgia"; r.font.bold = True; r.font.size = Pt(52); r.font.color.rgb = RGBColor(0,0,0)
                    r = p.add_run(); r.text = phrase; r.font.name = "Georgia"; r.font.bold = True; r.font.size = Pt(52); r.font.color.rgb = burgundy
                    if after:
                        r = p.add_run(); r.text = after; r.font.name = "Georgia"; r.font.bold = True; r.font.size = Pt(52); r.font.color.rgb = RGBColor(0,0,0)
                else:
                    r = p.add_run(); r.text = line; r.font.name = "Georgia"; r.font.bold = True; r.font.size = Pt(52); r.font.color.rgb = RGBColor(0,0,0)
            print(f"Created slide {slide_count}: St. Joseph prayer text")
        for lines in slides:
            add_slide(lines)
        return slide_count
    except Exception as e:
        print(f"  WARNING: Error creating St. Joseph prayer text slides: {e}")
        import traceback; traceback.print_exc()
        return slide_count


def create_lords_prayer_slide(prs, slide_count):
    slide_count += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4.5))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "THE LORD'S PRAYER"; r.font.name = "Georgia"; r.font.bold = True; r.font.size = Pt(60); r.font.color.rgb = RGBColor(0x98,0x00,0x00)
    print(f"Created slide {slide_count}: THE LORD'S PRAYER (title)")
    return slide_count
